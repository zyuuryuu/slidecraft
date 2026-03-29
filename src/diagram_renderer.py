"""
DiagramSpec → PPTX Renderer

Converts a DiagramSpec (intermediate JSON) into an editable PPTX slide
with auto-layouted shapes and connectors.

Rendering pipeline:
  1. Parse DiagramSpec → validated data structures
  2. Auto-layout: assign layers (topological sort) → position nodes
  3. Draw: shapes → group zones → connectors (with PoC-proven helpers)

Theme integration:
  Uses ThemeConfig (from diagram_theme.py) to apply consistent colors
  and fonts matching the target template (default: Midnight Executive).
  Can also render onto an existing template PPTX for header bar / slide
  number inheritance.

See DIAGRAM_PIPELINE_SPEC.md §3 for algorithm documentation.
"""

from __future__ import annotations
from collections import defaultdict, deque
from dataclasses import dataclass
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from lxml import etree

from diagram_schema import (
    DiagramSpec, Node, Edge, Group, Lane, NodeStyle, EdgeStyle,
    LaneStyle, LayoutConfig, parse_diagram_json,
)
from diagram_theme import ThemeConfig, DEFAULT_THEME
from diagram_icons import get_icon_png_path


# ══════════════════════════════════════════════
# Constants
# ══════════════════════════════════════════════

_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'

SHAPE_MAP = {
    "rect":         MSO_SHAPE.RECTANGLE,
    "rounded_rect": MSO_SHAPE.ROUNDED_RECTANGLE,
    "diamond":      MSO_SHAPE.DIAMOND,
    "circle":       MSO_SHAPE.OVAL,
    "oval":         MSO_SHAPE.OVAL,
    "hexagon":      MSO_SHAPE.HEXAGON,
}

SLIDE_W = 13.333  # inches (widescreen 16:9)
SLIDE_H = 7.5


# ══════════════════════════════════════════════
# Layout Engine
# ══════════════════════════════════════════════

@dataclass
class NodePosition:
    """Computed position and size for a node (in inches)."""
    node_id: str
    x: float       # left
    y: float       # top
    w: float        # width
    h: float        # height
    layer: int      # assigned layer (0 = root)
    order: int      # order within layer
    scale: float = 1.0   # layout scale factor (1.0 = no shrink)


def _build_adjacency(spec: DiagramSpec) -> tuple[dict, dict]:
    """Build forward and reverse adjacency lists from edges."""
    fwd = defaultdict(list)   # parent → children
    rev = defaultdict(list)   # child → parents
    for e in spec.edges:
        fwd[e.from_id].append(e.to_id)
        rev[e.to_id].append(e.from_id)
    return dict(fwd), dict(rev)


def _find_back_edges(fwd: dict, node_ids: list[str]) -> set[tuple[str, str]]:
    """Detect back-edges in a directed graph using DFS.

    A back-edge (u→v) is one where v is an ancestor of u in the DFS tree,
    indicating a cycle.  These edges are excluded from layer assignment
    but still rendered as connections.
    """
    WHITE, GRAY, BLACK = 0, 1, 2
    color = {nid: WHITE for nid in node_ids}
    back_edges: set[tuple[str, str]] = set()

    def dfs(u: str) -> None:
        color[u] = GRAY
        for v in fwd.get(u, []):
            if v not in color:
                continue
            if color[v] == GRAY:
                back_edges.add((u, v))
            elif color[v] == WHITE:
                dfs(v)
        color[u] = BLACK

    for nid in node_ids:
        if color[nid] == WHITE:
            dfs(nid)

    return back_edges


def _assign_layers(spec: DiagramSpec) -> dict[str, int]:
    """Assign each node to a layer using longest-path from roots (BFS).

    Roots = nodes with no incoming edges (after removing back-edges).
    Back-edges (cycles) are detected via DFS and excluded from the DAG
    used for layer computation, so cycles don't inflate layer counts.
    """
    fwd, rev = _build_adjacency(spec)
    node_ids = [n.id for n in spec.nodes]

    # Detect and remove back-edges to get a DAG
    back_edges = _find_back_edges(fwd, node_ids)
    dag_fwd: dict[str, list[str]] = defaultdict(list)
    dag_rev: dict[str, list[str]] = defaultdict(list)
    for u, children in fwd.items():
        for v in children:
            if (u, v) not in back_edges:
                dag_fwd[u].append(v)
                dag_rev[v].append(u)

    # Find roots (no incoming edges in the DAG)
    roots = [nid for nid in node_ids
             if nid not in dag_rev or len(dag_rev[nid]) == 0]
    if not roots:
        roots = [node_ids[0]] if node_ids else []

    # BFS longest path on the DAG (no cycles → guaranteed to terminate)
    layers: dict[str, int] = {}
    queue: deque[str] = deque()
    for r in roots:
        layers[r] = 0
        queue.append(r)

    while queue:
        nid = queue.popleft()
        cur_layer = layers[nid]
        for child in dag_fwd.get(nid, []):
            new_layer = cur_layer + 1
            if child not in layers or layers[child] < new_layer:
                layers[child] = new_layer
                queue.append(child)

    # Assign any disconnected nodes
    for nid in node_ids:
        if nid not in layers:
            layers[nid] = 0

    return layers


def _order_within_layers(spec: DiagramSpec, layers: dict[str, int]) -> dict[str, int]:
    """Determine left-to-right order within each layer using barycenter method.

    Uses a multi-pass barycenter heuristic to minimise edge crossings
    while respecting group constraints (nodes in the same group stay adjacent).

    Algorithm:
      1. Initial ordering: group hierarchy + original JSON order
      2. Forward sweep (layer 0→N): for each layer, compute the barycenter
         (average position of connected nodes in the PREVIOUS layer) and
         reorder.  Nodes without upstream connections keep their position.
      3. Backward sweep (layer N→0): same, but using the NEXT layer.
      4. Repeat for a fixed number of iterations (default 4).

    Group constraint: within each layer, nodes are first clustered by their
    group-sort-key.  Barycenter reordering only happens WITHIN each group
    cluster.  Groups themselves are reordered by their members' average
    barycenter to improve inter-group crossing reduction.
    """
    fwd, rev = _build_adjacency(spec)

    node_group = {n.id: n.group for n in spec.nodes}
    group_order = {g.id: i for i, g in enumerate(spec.groups)}

    def _group_sort_key(nid: str) -> tuple:
        gid = node_group.get(nid)
        if gid is None:
            return (9999,)
        chain = []
        cur = gid
        while cur is not None:
            chain.append(group_order.get(cur, 9999))
            parent_grp = next((g for g in spec.groups if g.id == cur), None)
            cur = parent_grp.parent if parent_grp else None
        chain.reverse()
        return tuple(chain)

    # Group nodes by layer
    layer_nodes: dict[int, list[str]] = defaultdict(list)
    for n in spec.nodes:
        layer_nodes[layers[n.id]].append(n.id)
    num_layers = max(layer_nodes.keys()) + 1 if layer_nodes else 0

    # Original order as fallback tiebreaker
    orig_idx = {n.id: i for i, n in enumerate(spec.nodes)}

    # ── Step 1: Initial ordering (group-clustered, original order) ──
    orders: dict[str, float] = {}
    for layer_idx in sorted(layer_nodes.keys()):
        nids = layer_nodes[layer_idx]
        nids_sorted = sorted(nids, key=lambda nid: (
            _group_sort_key(nid), orig_idx.get(nid, 0)
        ))
        for i, nid in enumerate(nids_sorted):
            orders[nid] = float(i)

    # ── Helper: group-constrained reorder ──
    def _group_constrained_reorder(nids: list[str], bary: dict[str, float]):
        """Reorder *nids* by barycenter while keeping group clusters intact.

        1. Partition nids into group-clusters (preserving cluster order).
        2. Within each cluster, sort nodes by their barycenter value.
        3. Reorder clusters by the average barycenter of their members.
        """
        # Partition into clusters
        clusters: list[tuple[tuple, list[str]]] = []
        for nid in nids:
            gk = _group_sort_key(nid)
            if clusters and clusters[-1][0] == gk:
                clusters[-1][1].append(nid)
            else:
                clusters.append((gk, [nid]))

        # Sort within each cluster by barycenter
        for _, members in clusters:
            members.sort(key=lambda nid: bary.get(nid, orders.get(nid, 0.0)))

        # Sort clusters by average barycenter of members
        def _cluster_bary(cluster):
            _, members = cluster
            vals = [bary.get(nid, orders.get(nid, 0.0)) for nid in members]
            return sum(vals) / len(vals) if vals else 0.0

        clusters.sort(key=_cluster_bary)

        # Flatten and assign sequential orders
        result = []
        for _, members in clusters:
            result.extend(members)
        return result

    # ── Step 2-4: Iterative barycenter sweeps ──
    NUM_ITERATIONS = 4

    for iteration in range(NUM_ITERATIONS):
        # ── Forward sweep: layer 0 → N-1 ──
        for layer_idx in range(1, num_layers):
            nids = layer_nodes.get(layer_idx, [])
            if not nids:
                continue

            bary: dict[str, float] = {}
            for nid in nids:
                # Barycenter = average order of upstream neighbors (previous layers)
                parents = rev.get(nid, [])
                parent_orders = [orders[p] for p in parents if p in orders]
                if parent_orders:
                    bary[nid] = sum(parent_orders) / len(parent_orders)
                # else: no barycenter → keep current position

            reordered = _group_constrained_reorder(nids, bary)
            for i, nid in enumerate(reordered):
                orders[nid] = float(i)

        # ── Backward sweep: layer N-1 → 0 ──
        for layer_idx in range(num_layers - 2, -1, -1):
            nids = layer_nodes.get(layer_idx, [])
            if not nids:
                continue

            bary: dict[str, float] = {}
            for nid in nids:
                # Barycenter = average order of downstream neighbors (next layers)
                children = fwd.get(nid, [])
                child_orders = [orders[c] for c in children if c in orders]
                if child_orders:
                    bary[nid] = sum(child_orders) / len(child_orders)

            reordered = _group_constrained_reorder(nids, bary)
            for i, nid in enumerate(reordered):
                orders[nid] = float(i)

    # Convert to integer orders
    final_orders: dict[str, int] = {}
    for layer_idx in sorted(layer_nodes.keys()):
        nids = layer_nodes[layer_idx]
        nids_sorted = sorted(nids, key=lambda nid: orders.get(nid, 0.0))
        for i, nid in enumerate(nids_sorted):
            final_orders[nid] = i

    return final_orders


def _compute_layout_v2(spec: DiagramSpec, content_top: float = 0.8) -> list[NodePosition]:
    """Group-aware layout algorithm (試案② full implementation).

    2-pass approach where groups "own" rectangular regions:
      Pass 1 (bottom-up): Compute required cross-axis size for each group.
      Pass 2 (top-down): Allocate cross-axis bands and place nodes within.

    Key concepts:
      - Main axis = direction of flow (X for LR/RL, Y for TB/BT)
      - Cross axis = perpendicular (Y for LR/RL, X for TB/BT)
      - Each group gets a fixed cross-axis band; nodes are placed within it.
      - Ungrouped nodes are treated as members of an implicit root group.
      - Groups without parent are treated as top-level sections.

    This solves the fundamental problem of v1 where nodes were placed
    without considering group containment, causing group boxes to overlap.
    """
    layout = spec.layout
    layers = _assign_layers(spec)
    orders = _order_within_layers(spec, layers)

    nw = layout.node_width
    nh = layout.node_height
    hg = layout.h_gap
    vg = layout.v_gap

    is_horizontal = spec.direction in ("LR", "RL")
    is_reversed = spec.direction in ("BT", "RL")

    # Diamond nodes need extra height
    diamond_ids = {n.id for n in spec.nodes if n.shape == "diamond"}
    node_h_map = {}
    for n in spec.nodes:
        node_h_map[n.id] = nh * 1.6 if n.id in diamond_ids else nh

    num_layers = max(layers.values(), default=0) + 1

    margin_x = 0.8
    margin_y_top = content_top
    margin_y_bot = 0.3

    # ── Axis abstraction ──
    # main = layer direction, cross = perpendicular
    if is_horizontal:
        main_node = nw           # node size along main axis
        cross_node = nh          # node size along cross axis
        main_gap = vg            # gap between layers (main)
        cross_gap = hg           # gap between nodes in same layer (cross)
    else:
        main_node = nh
        cross_node = nw
        main_gap = vg
        cross_gap = hg

    # ── Group membership ──
    group_direct_nodes: dict[str, list[str]] = defaultdict(list)
    ungrouped_nodes: list[str] = []
    for n in spec.nodes:
        if n.group:
            group_direct_nodes[n.group].append(n.id)
        else:
            ungrouped_nodes.append(n.id)

    # Group nesting constants
    _DEPTH_PAD = {0: 0.25, 1: 0.18, 2: 0.12}
    _LABEL_H = 0.25

    def _all_node_ids(gid: str) -> list[str]:
        """All node IDs in a group and its descendants."""
        result = list(group_direct_nodes.get(gid, []))
        for child_gid in spec.group_children(gid):
            result.extend(_all_node_ids(child_gid))
        return result

    def _layer_span(nids: list[str]) -> tuple[int, int]:
        """(min_layer, max_layer) for a set of node IDs."""
        if not nids:
            return (0, 0)
        lyrs = [layers[nid] for nid in nids]
        return (min(lyrs), max(lyrs))

    def _max_cross_slots(nids: list[str]) -> int:
        """Max number of nodes at any single layer (determines cross-axis size)."""
        if not nids:
            return 0
        layer_counts: dict[int, int] = defaultdict(int)
        for nid in nids:
            layer_counts[layers[nid]] += 1
        return max(layer_counts.values())

    # ── Pass 1: Bottom-up cross-axis size calculation ──
    # Returns the required cross-axis size (in inches) for a group.

    def _group_cross_size(gid: str) -> float:
        """Total cross-axis size (inches) for a group, including padding."""
        depth = spec.group_depth(gid)
        pad = _DEPTH_PAD.get(depth, 0.10)
        children = spec.group_children(gid)

        if children:
            # Parent group: stack children + direct nodes along cross-axis
            total = 0.0
            for i, child_gid in enumerate(children):
                if i > 0:
                    total += cross_gap * 0.5  # reduced gap between sibling groups
                total += _group_cross_size(child_gid)

            # Direct member nodes (not in any child group)
            direct_nids = group_direct_nodes.get(gid, [])
            if direct_nids:
                max_slots = _max_cross_slots(direct_nids)
                if total > 0:
                    total += cross_gap * 0.5
                total += max_slots * cross_node + (max_slots - 1) * cross_gap
        else:
            # Leaf group: size from its member nodes
            all_nids = _all_node_ids(gid)
            max_slots = _max_cross_slots(all_nids)
            total = max_slots * cross_node + max(0, (max_slots - 1)) * cross_gap

        return total + 2 * pad + _LABEL_H

    # ── Per-layer main-axis gap computation (Algorithm A: group-boundary-aware) ──
    # Three-tier gap sizing:
    #   1. Intra-group: compact (just enough to see layer separation)
    #   2. Inter-group boundary: sized to accommodate group padding + label
    #   3. Unrelated groups: normal visual gap
    #
    # Algorithm A ensures that when a group boundary falls between layers L and L+1,
    # the gap is large enough for: ending group's bottom padding + minimum inter-group
    # visual gap + starting group's top padding + label height.

    COMPACT_GAP_RATIO = 0.25   # intra-group gap = main_gap * 0.25
    MIN_GROUP_VISUAL_GAP = 0.10  # minimum visual gap between group bounding boxes

    node_group_map = {n.id: n.group for n in spec.nodes}
    layer_to_nodes: dict[int, list[str]] = defaultdict(list)
    for nid, lyr in layers.items():
        layer_to_nodes[lyr].append(nid)

    # Precompute: which groups have their max/min layer at each layer
    group_layer_ranges_local: dict[str, tuple[int, int]] = {}
    for g in spec.groups:
        g_nids = [n.id for n in spec.nodes if n.group == g.id]
        if g_nids:
            g_lyrs = [layers[nid] for nid in g_nids if nid in layers]
            if g_lyrs:
                group_layer_ranges_local[g.id] = (min(g_lyrs), max(g_lyrs))

    def _groups_ending_at(lyr: int) -> set[str]:
        """Groups whose last layer is exactly lyr."""
        return {gid for gid, (_, mx) in group_layer_ranges_local.items() if mx == lyr}

    def _groups_starting_at(lyr: int) -> set[str]:
        """Groups whose first layer is exactly lyr."""
        return {gid for gid, (mn, _) in group_layer_ranges_local.items() if mn == lyr}

    # Compute gap between layer L and L+1
    layer_gaps = {}  # layer_idx → gap size (in inches)
    for lyr in range(num_layers - 1):
        groups_this = {node_group_map.get(nid) for nid in layer_to_nodes[lyr]}
        groups_next = {node_group_map.get(nid) for nid in layer_to_nodes[lyr + 1]}
        groups_this.discard(None)
        groups_next.discard(None)

        if groups_this and groups_next and groups_this == groups_next:
            # Same group(s) in both layers → compact gap
            layer_gaps[lyr] = main_gap * COMPACT_GAP_RATIO
        else:
            # Check if a group boundary falls between these layers
            ending = _groups_ending_at(lyr)
            starting = _groups_starting_at(lyr + 1)

            if ending and starting:
                # Group boundary: ensure enough space for padding + label + visual gap
                # ending group bottom padding + visual gap + starting group top padding + label
                max_end_depth = max((spec.group_depth(gid) for gid in ending), default=0)
                max_start_depth = max((spec.group_depth(gid) for gid in starting), default=0)
                end_pad = _DEPTH_PAD.get(max_end_depth, 0.10)
                start_pad = _DEPTH_PAD.get(max_start_depth, 0.10)

                boundary_gap = end_pad + MIN_GROUP_VISUAL_GAP + start_pad + _LABEL_H
                # The gap between layer centers is: main_node (for the ending node) + gap
                # The group bbox extends by pad below the ending node and pad+label above
                # the starting node. So: gap >= boundary_gap
                layer_gaps[lyr] = max(boundary_gap, main_gap * COMPACT_GAP_RATIO)
            elif groups_this & groups_next:
                # Partial overlap (some groups continue, some change) → medium gap
                layer_gaps[lyr] = main_gap * 0.5
            else:
                # Different groups, no boundary alignment → standard gap
                layer_gaps[lyr] = main_gap * 0.5

    # Build cumulative main-axis position for each layer
    layer_main_pos: dict[int, float] = {0: 0.0}
    for lyr in range(1, num_layers):
        gap = layer_gaps.get(lyr - 1, main_gap * 0.5)
        layer_main_pos[lyr] = layer_main_pos[lyr - 1] + main_node + gap

    # Total main-axis size using variable gaps
    total_main_var = layer_main_pos.get(num_layers - 1, 0.0) + main_node

    # ── Pass 2: Top-down coordinate allocation ──
    # Allocate cross-axis bands and place nodes within them.

    node_positions: dict[str, tuple] = {}  # nid → (x, y, w, h, layer, order)

    def _place_nodes_in_band(nids: list[str], cross_origin: float):
        """Place a set of nodes within a cross-axis band.

        Main-axis position is determined by precomputed layer_main_pos.
        Cross-axis position is within the allocated band.
        """
        if not nids:
            return
        layer_groups_local: dict[int, list[str]] = defaultdict(list)
        for nid in nids:
            layer_groups_local[layers[nid]].append(nid)

        for lyr, lyr_nids in layer_groups_local.items():
            lyr_nids.sort(key=lambda nid: orders.get(nid, 0))
            for i, nid in enumerate(lyr_nids):
                m_pos = layer_main_pos.get(lyr, 0.0)
                c_pos = cross_origin + i * (cross_node + cross_gap)
                this_h = node_h_map.get(nid, nh)

                if is_horizontal:
                    node_positions[nid] = (m_pos, c_pos, nw, this_h, lyr, orders.get(nid, 0))
                else:
                    node_positions[nid] = (c_pos, m_pos, nw, this_h, lyr, orders.get(nid, 0))

    def _allocate_group(gid: str, cross_origin: float):
        """Allocate a cross-axis band for a group and place its nodes.

        Main-axis positions are global (determined by layer index).
        Only cross-axis positions are allocated by group hierarchy.
        """
        depth = spec.group_depth(gid)
        pad = _DEPTH_PAD.get(depth, 0.10)
        inner_cross = cross_origin + pad + _LABEL_H

        children = spec.group_children(gid)
        current_cross = inner_cross

        if children:
            # Allocate children along cross-axis
            for i, child_gid in enumerate(children):
                if i > 0:
                    current_cross += cross_gap * 0.5
                _allocate_group(child_gid, current_cross)
                current_cross += _group_cross_size(child_gid)

            # Direct member nodes (after child groups in cross-axis)
            direct_nids = group_direct_nodes.get(gid, [])
            if direct_nids:
                if current_cross > inner_cross:
                    current_cross += cross_gap * 0.5
                _place_nodes_in_band(direct_nids, current_cross)
        else:
            # Leaf group: place all nodes within this group's band
            all_nids = _all_node_ids(gid)
            _place_nodes_in_band(all_nids, inner_cross)

    # ── Compute total required sizes ──
    top_groups = spec.top_level_groups()

    # ── Column packing via interval coloring ──
    # Groups with non-overlapping layer ranges share the same cross-axis column,
    # dramatically reducing wasted space in diagrams with many sequential groups.

    # Compute layer ranges for each top-level group
    group_layer_ranges: dict[str, tuple[int, int]] = {}
    for g in top_groups:
        all_nids = _all_node_ids(g.id)
        if all_nids:
            g_layers = [layers[nid] for nid in all_nids if nid in layers]
            if g_layers:
                group_layer_ranges[g.id] = (min(g_layers), max(g_layers))

    # Greedy interval coloring: assign groups to columns
    # Sort by min layer, then by span length (shorter first)
    sorted_tl = sorted(
        [g for g in top_groups if g.id in group_layer_ranges],
        key=lambda g: (group_layer_ranges[g.id][0],
                       group_layer_ranges[g.id][1] - group_layer_ranges[g.id][0])
    )

    columns: list[list] = []           # column_idx → [group objects]
    column_intervals: list[list] = []  # column_idx → [(min_layer, max_layer)]
    group_column: dict[str, int] = {}

    for g in sorted_tl:
        g_min, g_max = group_layer_ranges[g.id]
        assigned = False
        for col_idx in range(len(columns)):
            # Check overlap with all groups in this column
            fits = True
            for (e_min, e_max) in column_intervals[col_idx]:
                if g_min <= e_max and g_max >= e_min:  # strict overlap
                    fits = False
                    break
            if fits:
                columns[col_idx].append(g)
                column_intervals[col_idx].append((g_min, g_max))
                group_column[g.id] = col_idx
                assigned = True
                break
        if not assigned:
            columns.append([g])
            column_intervals.append([(g_min, g_max)])
            group_column[g.id] = len(columns) - 1

    # Groups not in layer_ranges (empty groups) go in column 0
    for g in top_groups:
        if g.id not in group_column:
            if columns:
                columns[0].append(g)
                group_column[g.id] = 0
            else:
                columns.append([g])
                group_column[g.id] = 0
                column_intervals.append([])

    # ── Reorder columns to minimize cross-column edge jumps ──
    # Heuristic: the column with the most total edges to other columns
    # should be in the center; columns with fewer connections on the sides.
    if len(columns) > 2:
        # Count inter-column edges per column
        col_edge_count = [0] * len(columns)
        node_to_group_map = {n.id: n.group for n in spec.nodes if n.group}
        for e in spec.edges:
            fg = node_to_group_map.get(e.from_id)
            tg = node_to_group_map.get(e.to_id)
            if fg and tg and fg != tg:
                fc = group_column.get(fg)
                tc = group_column.get(tg)
                if fc is not None and tc is not None and fc != tc:
                    col_edge_count[fc] += 1
                    col_edge_count[tc] += 1

        # Sort: highest connectivity columns in the middle
        indexed = list(enumerate(col_edge_count))
        indexed.sort(key=lambda x: x[1], reverse=True)
        n_cols = len(columns)
        new_order = [0] * n_cols
        # Place highest connectivity at center, then alternate left/right
        positions_lr = []
        mid = n_cols // 2
        for i in range(n_cols):
            if i % 2 == 0:
                positions_lr.append(mid + i // 2)
            else:
                positions_lr.append(mid - (i + 1) // 2)
        for rank, (orig_idx, _) in enumerate(indexed):
            new_order[orig_idx] = positions_lr[rank]

        reordered_columns = [None] * n_cols
        for orig_idx, new_idx in enumerate(new_order):
            reordered_columns[new_idx] = columns[orig_idx]
        columns = [c for c in reordered_columns if c is not None]

        # Rebuild group_column mapping
        group_column.clear()
        for col_idx, col_groups in enumerate(columns):
            for g in col_groups:
                group_column[g.id] = col_idx

    # ── Compute cross-axis size per column ──
    # Each column's size = max group cross-axis size among groups in that column
    column_cross_sizes = []
    for col_groups in columns:
        max_size = max((_group_cross_size(g.id) for g in col_groups), default=0)
        column_cross_sizes.append(max_size)

    # Total cross-axis = sum of column sizes + gaps
    total_cross = sum(column_cross_sizes) + max(0, (len(columns) - 1)) * cross_gap

    if ungrouped_nodes:
        max_ug_slots = _max_cross_slots(ungrouped_nodes)
        ungrouped_cross_size = max_ug_slots * cross_node + max(0, (max_ug_slots - 1)) * cross_gap
        total_cross += cross_gap + ungrouped_cross_size

    # Total main-axis (using variable gaps from layer_main_pos)
    total_main = total_main_var
    if total_main <= 0:
        total_main = main_node

    # ── Scale to fit slide ──
    if is_horizontal:
        avail_main = SLIDE_W - 2 * margin_x
        avail_cross = SLIDE_H - margin_y_top - margin_y_bot
    else:
        avail_main = SLIDE_H - margin_y_top - margin_y_bot
        avail_cross = SLIDE_W - 2 * margin_x

    scale_main = avail_main / total_main if total_main > 0 else 1.0
    scale_cross = avail_cross / total_cross if total_cross > 0 else 1.0
    scale = min(scale_main, scale_cross, 1.0)  # never scale up

    # Apply scale to all sizing parameters
    s_main_node = main_node * scale
    s_cross_node = cross_node * scale
    s_main_gap = main_gap * scale
    s_cross_gap = cross_gap * scale
    s_nw = nw * scale
    s_nh = nh * scale

    # Scale the parameters used by nested functions
    # (We re-run the allocation with scaled parameters)
    # For simplicity, compute at original scale then apply uniform scale + offset.

    # ── Algorithm B: Connectivity-based cross-axis offset within columns ──
    # Instead of centering every group in its column, shift each group toward
    # the weighted average cross-axis position of its connected neighbor groups
    # in OTHER columns. This reduces line crossings and x-overlap between
    # vertically adjacent groups in the same column.

    # First compute column cross-axis origins (still needed for absolute positions)
    column_cross_origins = []
    _cur = 0.0
    for col_idx, col_groups in enumerate(columns):
        column_cross_origins.append(_cur)
        _cur += column_cross_sizes[col_idx] + cross_gap

    # Compute column center positions for connectivity-based offset
    column_centers = [
        column_cross_origins[i] + column_cross_sizes[i] / 2
        for i in range(len(columns))
    ]

    # Build group-to-group edge weights
    _ntg = {n.id: n.group for n in spec.nodes if n.group}
    _g2g_weights: dict[str, dict[str, int]] = defaultdict(lambda: defaultdict(int))
    for e in spec.edges:
        fg, tg = _ntg.get(e.from_id), _ntg.get(e.to_id)
        if fg and tg and fg != tg:
            _g2g_weights[fg][tg] += 1
            _g2g_weights[tg][fg] += 1

    for col_idx, col_groups in enumerate(columns):
        col_origin = column_cross_origins[col_idx]
        col_size = column_cross_sizes[col_idx]

        for g in col_groups:
            g_size = _group_cross_size(g.id)
            max_offset = col_size - g_size  # available slack within column

            if max_offset <= 0.01:
                # Group fills the column — no room to shift
                _allocate_group(g.id, col_origin)
                continue

            # Compute weighted average cross-axis position of connected groups
            # in OTHER columns
            weighted_sum = 0.0
            total_weight = 0
            for neighbor_gid, weight in _g2g_weights.get(g.id, {}).items():
                n_col = group_column.get(neighbor_gid)
                if n_col is not None and n_col != col_idx:
                    weighted_sum += column_centers[n_col] * weight
                    total_weight += weight

            if total_weight > 0:
                # Target: shift toward weighted average of neighbors
                target_center = weighted_sum / total_weight
                my_center_if_centered = col_origin + col_size / 2

                # How much to shift: proportional to how far the target is
                # from our centered position. Clamp to [0, max_offset].
                shift_direction = target_center - my_center_if_centered
                # Normalize: full shift if target is in another column,
                # proportional otherwise
                shift_ratio = min(abs(shift_direction) / (col_size + cross_gap + 1e-6), 1.0)
                if shift_direction > 0:
                    # Shift right (increase cross_origin)
                    g_offset = max_offset * shift_ratio
                else:
                    # Shift left (decrease cross_origin, but clamped to 0)
                    g_offset = max_offset * (1.0 - shift_ratio)
            else:
                # No cross-column connections: center in column
                g_offset = max_offset / 2

            _allocate_group(g.id, col_origin + g_offset)

    if ungrouped_nodes:
        _place_nodes_in_band(ungrouped_nodes, _cur)

    # ── Cross-axis stretch: use available space when main axis is bottleneck ──
    # When scale_main < scale_cross, the cross-axis has leftover space.
    # Stretch cross-axis SPACING (not node sizes) to fill the available width.
    # This spreads groups apart while keeping nodes proportional.
    MAX_CROSS_STRETCH = 2.5  # cap to avoid overly sparse layout
    if scale_main < scale_cross and total_cross > 0:
        cross_stretch = min(avail_cross / (total_cross * scale), MAX_CROSS_STRETCH)
    else:
        cross_stretch = 1.0

    # ── Transform: scale + cross-stretch + center + margin ──
    positions = []

    # Compute cross-axis center of the raw layout for stretch pivot
    all_raw_cross = []
    for nid, (raw_x, raw_y, raw_w, raw_h, lyr, ord_) in node_positions.items():
        if is_horizontal:
            all_raw_cross.append(raw_y + raw_h / 2)
        else:
            all_raw_cross.append(raw_x + raw_w / 2)
    raw_cross_center = (min(all_raw_cross) + max(all_raw_cross)) / 2 if all_raw_cross else 0

    for nid, (raw_x, raw_y, raw_w, raw_h, lyr, ord_) in node_positions.items():
        if is_horizontal:
            # Main axis = X, Cross axis = Y
            sx = raw_x * scale
            # Stretch Y positions around center, keeping node height unchanged
            y_center = raw_y + raw_h / 2
            sy = ((y_center - raw_cross_center) * cross_stretch + raw_cross_center) * scale - raw_h * scale / 2
            positions.append(NodePosition(
                node_id=nid, x=sx, y=sy,
                w=raw_w * scale, h=raw_h * scale,
                layer=lyr, order=ord_, scale=scale,
            ))
        else:
            # Main axis = Y, Cross axis = X
            sy = raw_y * scale
            # Stretch X positions around center, keeping node width unchanged
            x_center = raw_x + raw_w / 2
            sx = ((x_center - raw_cross_center) * cross_stretch + raw_cross_center) * scale - raw_w * scale / 2
            positions.append(NodePosition(
                node_id=nid, x=sx, y=sy,
                w=raw_w * scale, h=raw_h * scale,
                layer=lyr, order=ord_, scale=scale,
            ))

    # Center in available space
    if positions:
        if is_horizontal:
            scaled_main = total_main * scale
            scaled_cross = total_cross * scale * cross_stretch
            offset_x = margin_x + (avail_main - scaled_main) / 2
            offset_y = margin_y_top + (avail_cross - scaled_cross) / 2
        else:
            scaled_main = total_main * scale
            scaled_cross = total_cross * scale * cross_stretch
            offset_x = margin_x + (avail_cross - scaled_cross) / 2
            offset_y = margin_y_top + (avail_main - scaled_main) / 2

        for p in positions:
            p.x += offset_x
            p.y += offset_y

    # Handle reversed directions
    if is_reversed and positions:
        for p in positions:
            if is_horizontal:
                p.x = SLIDE_W - p.x - p.w
            else:
                p.y = SLIDE_H - p.y - p.h

    return positions


# ══════════════════════════════════════════════
# Swimlane layout algorithm (v3)
# ══════════════════════════════════════════════

@dataclass
class LaneInfo:
    """Computed geometry for a single swimlane (in inches)."""
    lane_id: str
    label: str
    cross_origin: float   # start position on cross-axis
    cross_size: float     # size on cross-axis (band width/height)
    style: LaneStyle


def _compute_layout_swimlane(spec: DiagramSpec, content_top: float = 0.8) -> tuple[list[NodePosition], list[LaneInfo]]:
    """Swimlane-aware layout algorithm.

    Lanes divide the cross-axis into bands. Each node is placed within
    its assigned lane's band. Main-axis position is determined by BFS layer.

    For TB/BT: lanes are vertical columns (cross-axis = X), flow is Y.
    For LR/RL: lanes are horizontal rows (cross-axis = Y), flow is X.

    Returns:
        (positions, lane_infos): Node positions and computed lane geometries.
    """
    layout = spec.layout
    layers = _assign_layers(spec)
    orders = _order_within_layers(spec, layers)

    is_horizontal = spec.direction in ("LR", "RL")
    is_reversed = spec.direction in ("BT", "RL")

    nw = layout.node_width
    nh = layout.node_height
    hg = layout.h_gap
    vg = layout.v_gap

    # Diamond nodes need extra height
    diamond_ids = {n.id for n in spec.nodes if n.shape == "diamond"}
    node_h_map = {n.id: (nh * 1.6 if n.id in diamond_ids else nh) for n in spec.nodes}

    num_layers = max(layers.values(), default=0) + 1

    # Margins
    margin = 0.3
    lane_header_size = 0.6  # width (TB) or height (LR) of the lane header

    # Available space
    if is_horizontal:
        # LR/RL: lanes are horizontal rows, header at left
        main_start = margin + lane_header_size + 0.1
        main_avail = SLIDE_W - main_start - margin
        cross_start = content_top + 0.1
        cross_avail = SLIDE_H - cross_start - margin
    else:
        # TB/BT: lanes are vertical columns, header at top
        main_start = content_top + lane_header_size + 0.1
        main_avail = SLIDE_H - main_start - margin
        cross_start = margin
        cross_avail = SLIDE_W - 2 * margin

    # ── Assign nodes to lanes ──
    lane_map = {ln.id: ln for ln in spec.lanes}
    lane_order = [ln.id for ln in spec.lanes]
    n_lanes = len(lane_order)

    # Group nodes by lane
    lane_nodes: dict[str, list[str]] = {lid: [] for lid in lane_order}
    unassigned = []
    for n in spec.nodes:
        if n.lane and n.lane in lane_nodes:
            lane_nodes[n.lane].append(n.id)
        else:
            unassigned.append(n.id)

    # Unassigned nodes go to a virtual "default" lane
    if unassigned:
        lane_order.append("__default__")
        lane_nodes["__default__"] = unassigned
        n_lanes += 1

    # ── Compute lane sizes ──
    # Each lane gets cross-axis size proportional to max nodes per layer
    lane_max_per_layer: dict[str, int] = {}
    for lid in lane_order:
        nids = lane_nodes[lid]
        if not nids:
            lane_max_per_layer[lid] = 1
            continue
        # Count nodes per layer in this lane
        layer_counts: dict[int, int] = defaultdict(int)
        for nid in nids:
            layer_counts[layers[nid]] += 1
        lane_max_per_layer[lid] = max(layer_counts.values(), default=1)

    # Minimum band size: enough for 1 node + padding
    cross_node = nh if is_horizontal else nw
    cross_gap = hg
    min_band = cross_node + cross_gap * 2

    # Distribute cross-axis space proportionally
    total_weight = sum(max(1, lane_max_per_layer[lid]) for lid in lane_order)
    lane_gap = 0.05  # gap between lanes
    usable_cross = cross_avail - lane_gap * (n_lanes - 1)

    lane_infos: list[LaneInfo] = []
    cur_cross = cross_start
    for lid in lane_order:
        weight = max(1, lane_max_per_layer[lid])
        band_size = max(min_band, usable_cross * weight / total_weight)
        style = lane_map[lid].effective_style() if lid in lane_map else LaneStyle()
        label = lane_map[lid].label if lid in lane_map else ""
        lane_infos.append(LaneInfo(
            lane_id=lid, label=label,
            cross_origin=cur_cross, cross_size=band_size,
            style=style,
        ))
        cur_cross += band_size + lane_gap

    # Build lane lookup
    lane_info_map = {li.lane_id: li for li in lane_infos}

    # ── Main-axis layer positions ──
    main_node = nw if is_horizontal else nh
    main_gap = vg
    raw_main_total = num_layers * main_node + (num_layers - 1) * main_gap
    main_scale = min(main_avail / raw_main_total, 1.0) if raw_main_total > 0 else 1.0

    def _layer_main_pos(layer_idx: int) -> float:
        pos = main_start + layer_idx * (main_node + main_gap) * main_scale
        if is_reversed:
            pos = (main_start + main_avail) - (pos - main_start) - main_node * main_scale
        return pos

    # ── Place nodes ──
    # Within each lane, for each layer, center nodes in the cross-axis band
    positions = []
    for lid in lane_order:
        li = lane_info_map[lid]
        nids = lane_nodes[lid]
        if not nids:
            continue

        # Group nodes by layer
        lane_layer_nodes: dict[int, list[str]] = defaultdict(list)
        for nid in nids:
            lane_layer_nodes[layers[nid]].append(nid)

        # Sort by order within each layer
        for lyr in lane_layer_nodes:
            lane_layer_nodes[lyr].sort(key=lambda nid: orders[nid])

        for lyr, lyr_nids in lane_layer_nodes.items():
            n_in = len(lyr_nids)
            for i, nid in enumerate(lyr_nids):
                this_h = node_h_map.get(nid, nh)
                main_pos = _layer_main_pos(lyr)

                # Center nodes within the lane's cross-axis band
                if is_horizontal:
                    # LR/RL: cross = Y, main = X
                    node_cross = this_h
                    total_cross = n_in * node_cross + (n_in - 1) * cross_gap
                    cross_offset = (li.cross_size - total_cross) / 2
                    x = main_pos
                    y = li.cross_origin + cross_offset + i * (node_cross + cross_gap)
                    w = nw * main_scale
                    h = this_h
                else:
                    # TB/BT: cross = X, main = Y
                    node_cross = nw
                    total_cross = n_in * node_cross + (n_in - 1) * cross_gap
                    cross_offset = (li.cross_size - total_cross) / 2
                    x = li.cross_origin + cross_offset + i * (node_cross + cross_gap)
                    y = main_pos
                    w = nw
                    h = this_h * main_scale

                positions.append(NodePosition(
                    node_id=nid, x=x, y=y, w=w, h=h,
                    layer=lyr, order=i,
                ))

    return positions, lane_infos


def compute_layout(spec: DiagramSpec, content_top: float = 0.8) -> list[NodePosition]:
    """Compute positions for all nodes based on topology and direction.

    Dispatches to:
      - v3 (swimlane) when lanes exist
      - v2 (group-aware) when groups exist
      - v1 (simple layer-centering) when neither

    Args:
        spec: The diagram specification
        content_top: Top margin (increased when using template header bar)
    """
    # Use swimlane layout when lanes are defined
    if spec.lanes:
        positions, _ = _compute_layout_swimlane(spec, content_top)
        return positions

    # Use v2 for any diagram with groups (nested or flat)
    if spec.groups:
        return _compute_layout_v2(spec, content_top)

    # ── v1: Simple layer-centering (no groups) ──
    layout = spec.layout
    layers = _assign_layers(spec)
    orders = _order_within_layers(spec, layers)

    # Group nodes by layer
    layer_nodes = defaultdict(list)
    for n in spec.nodes:
        layer_nodes[layers[n.id]].append(n.id)

    # Sort within each layer by order
    for layer_idx in layer_nodes:
        layer_nodes[layer_idx].sort(key=lambda nid: orders[nid])

    is_horizontal = spec.direction in ("LR", "RL")
    is_reversed = spec.direction in ("BT", "RL")

    # Compute sizes
    nw = layout.node_width
    nh = layout.node_height
    hg = layout.h_gap
    vg = layout.v_gap

    # Diamond nodes need extra height
    diamond_ids = {n.id for n in spec.nodes if n.shape == "diamond"}
    node_h_map = {}
    for n in spec.nodes:
        node_h_map[n.id] = nh * 1.6 if n.id in diamond_ids else nh

    num_layers = max(layers.values(), default=0) + 1

    # Margin from slide edges
    margin_x = 0.8
    margin_y_top = content_top
    margin_y_bot = 0.3

    # ── Compute raw positions ──
    positions = []

    for layer_idx in sorted(layer_nodes.keys()):
        nids = layer_nodes[layer_idx]
        n_in_layer = len(nids)

        for i, nid in enumerate(nids):
            this_h = node_h_map.get(nid, nh)

            if is_horizontal:
                x = margin_x + layer_idx * (nw + vg)
                total_h = n_in_layer * nh + (n_in_layer - 1) * hg
                y_start = (SLIDE_H - total_h) / 2
                y = y_start + i * (nh + hg)

                if is_reversed:
                    x = SLIDE_W - margin_x - nw - layer_idx * (nw + vg)
            else:
                total_w = n_in_layer * nw + (n_in_layer - 1) * hg
                x_start = (SLIDE_W - total_w) / 2
                x = x_start + i * (nw + hg)

                # Use per-layer max height for spacing
                y = margin_y_top + sum(
                    max((node_h_map.get(nid2, nh) for nid2 in layer_nodes.get(l, [nid])), default=nh) + vg
                    for l in range(layer_idx)
                )

                if is_reversed:
                    y = SLIDE_H - margin_y_bot - this_h - sum(
                        max((node_h_map.get(nid2, nh) for nid2 in layer_nodes.get(l, [nid])), default=nh) + vg
                        for l in range(layer_idx)
                    )

            positions.append(NodePosition(
                node_id=nid, x=x, y=y, w=nw, h=this_h,
                layer=layer_idx, order=i
            ))

    # ── Auto-scale to fit slide ──
    if positions:
        all_x = [p.x for p in positions]
        all_y = [p.y for p in positions]
        all_r = [p.x + p.w for p in positions]
        all_b = [p.y + p.h for p in positions]

        eff_w = max(all_r) - min(all_x)
        eff_h = max(all_b) - min(all_y)
        eff_left = min(all_x)
        eff_top = min(all_y)

        avail_w = SLIDE_W - 2 * margin_x
        avail_h = SLIDE_H - margin_y_top - margin_y_bot

        scale_x = avail_w / eff_w if eff_w > 0 else 1.0
        scale_y = avail_h / eff_h if eff_h > 0 else 1.0
        scale = min(scale_x, scale_y, 1.0)

        for p in positions:
            p.x = (p.x - eff_left) * scale
            p.y = (p.y - eff_top) * scale
            p.w *= scale
            p.h *= scale

        scaled_w = eff_w * scale
        scaled_h = eff_h * scale
        offset_x = margin_x + (avail_w - scaled_w) / 2
        offset_y = margin_y_top + (avail_h - scaled_h) / 2

        for p in positions:
            p.x += offset_x
            p.y += offset_y

    return positions


# ══════════════════════════════════════════════
# Drawing Helpers (PoC-proven)
# ══════════════════════════════════════════════

def _hex_to_rgb(hex_str: str) -> RGBColor:
    """Convert '#RRGGBB' or 'RRGGBB' to RGBColor."""
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _cp_coords(shape, cp_idx: int, port_offset: float = 0.0) -> tuple[int, int]:
    """Return (x_emu, y_emu) for connection point on a shape.

    Diamond shapes use vertex positions (not bounding-box midpoints).
    cp: 0=top, 1=right, 2=bottom, 3=left
    port_offset: fraction [-0.5, +0.5] to offset along the edge.
                 For top/bottom edges: offset along width (- = left, + = right).
                 For left/right edges: offset along height (- = up, + = down).
                 0.0 = center (default, backwards compatible).
    """
    l, t, w, h = shape.left, shape.top, shape.width, shape.height
    cx, cy = l + w // 2, t + h // 2

    is_diamond = False
    try:
        prst = shape._element.spPr.prstGeom.attrib.get('prst', '')
        is_diamond = (prst == 'diamond')
    except Exception:
        pass

    if is_diamond:
        # Diamond vertices — port_offset not meaningful, use center
        return {0: (cx, t), 1: (l + w, cy), 2: (cx, t + h), 3: (l, cy)}[cp_idx]

    # Rectangle: offset along the edge
    off_px = int(port_offset * w)  # horizontal offset in EMU
    off_py = int(port_offset * h)  # vertical offset in EMU

    return {
        0: (cx + off_px, t),         # top: offset along width
        1: (l + w, cy + off_py),     # right: offset along height
        2: (cx + off_px, t + h),     # bottom: offset along width
        3: (l, cy + off_py),         # left: offset along height
    }[cp_idx]


def _compute_port_offsets(
    spec,
    shape_map: dict,
    layer_map: dict,
    direction: str,
    back_edges: set,
    bus_handled: set,
) -> dict[tuple[str, str], tuple[float, float]]:
    """Compute port offsets for each edge to avoid overlapping connection points.

    Returns a dict mapping (from_id, to_id) → (src_port_offset, tgt_port_offset).

    For TB/BT layouts:
      - Bottom ports (outgoing): distribute along node width based on target x-position
      - Top ports (incoming): distribute along node width based on source x-position

    The ordering is based on the cross-axis position of the connected neighbor,
    so lines naturally fan out in the direction they need to go.
    """
    from pptx.util import Inches as _Inches

    is_vertical = direction in ("TB", "BT")

    # Collect edges per node per side (outgoing from bottom, incoming to top for TB)
    # Key: (node_id, side) where side is 'out' or 'in'
    node_out_edges: dict[str, list] = defaultdict(list)  # node_id → list of edges going out
    node_in_edges: dict[str, list] = defaultdict(list)   # node_id → list of edges coming in

    for edge in spec.edges:
        key = (edge.from_id, edge.to_id)
        if key in back_edges or key in bus_handled:
            continue
        if edge.from_id not in shape_map or edge.to_id not in shape_map:
            continue
        node_out_edges[edge.from_id].append(edge)
        node_in_edges[edge.to_id].append(edge)

    port_offsets: dict[tuple[str, str], tuple[float, float]] = {}

    def _assign_ports(node_id: str, edges: list, is_outgoing: bool):
        """Assign port offsets for a set of edges on the same side of a node."""
        if len(edges) <= 1:
            # Single edge — use center (offset = 0.0)
            for e in edges:
                key = (e.from_id, e.to_id)
                if key not in port_offsets:
                    port_offsets[key] = (0.0, 0.0)
                if is_outgoing:
                    port_offsets[key] = (0.0, port_offsets[key][1])
                else:
                    port_offsets[key] = (port_offsets[key][0], 0.0)
            return

        # Sort edges by the cross-axis position of the OTHER endpoint
        # This makes lines fan out naturally toward their destination
        def _neighbor_cross(e):
            other_id = e.to_id if is_outgoing else e.from_id
            s = shape_map.get(other_id)
            if s is None:
                return 0
            if is_vertical:
                return s.left + s.width // 2  # x-position
            else:
                return s.top + s.height // 2  # y-position

        sorted_edges = sorted(edges, key=_neighbor_cross)

        n = len(sorted_edges)
        # Spread across [-0.35, +0.35] of node width (leave 15% margin on each side)
        max_spread = 0.35
        for i, e in enumerate(sorted_edges):
            if n == 1:
                offset = 0.0
            else:
                offset = -max_spread + (2 * max_spread) * i / (n - 1)

            key = (e.from_id, e.to_id)
            if key not in port_offsets:
                port_offsets[key] = (0.0, 0.0)
            if is_outgoing:
                port_offsets[key] = (offset, port_offsets[key][1])
            else:
                port_offsets[key] = (port_offsets[key][0], offset)

    # Assign ports for all nodes
    for node_id, edges in node_out_edges.items():
        _assign_ports(node_id, edges, is_outgoing=True)
    for node_id, edges in node_in_edges.items():
        _assign_ports(node_id, edges, is_outgoing=False)

    return port_offsets


def _nudge_overlapping_segments(planned_paths: list, nudge_emu: int = 73000):
    """Post-process planned edge paths to separate overlapping horizontal segments.

    A horizontal segment is a pair of consecutive path points with the same y
    coordinate (within tolerance).  When multiple edges share horizontal
    segments at the same y-level with overlapping x-ranges, they are visually
    indistinguishable.  This function detects such conflicts and nudges the
    y-coordinate of conflicting segments apart.

    Each element of *planned_paths* is a dict with at least:
        'path_points': list of (x_emu, y_emu) tuples  (MUTATED in-place)

    nudge_emu ≈ 0.08" — the distance between staggered horizontal segments.
    """
    if not planned_paths:
        return

    Y_TOL = 5000   # ≈ 0.005" — treat as "same y" if within this tolerance

    # ── 1. Collect all horizontal segments ──
    # Each entry: (plan_idx, seg_idx, x_min, x_max, y)
    h_segments: list[tuple[int, int, int, int, int]] = []
    for pi, plan in enumerate(planned_paths):
        pts = plan['path_points']
        for si in range(len(pts) - 1):
            (x1, y1), (x2, y2) = pts[si], pts[si + 1]
            if abs(y1 - y2) <= Y_TOL and abs(x1 - x2) > Y_TOL:
                # Horizontal segment
                xlo, xhi = (min(x1, x2), max(x1, x2))
                y_avg = (y1 + y2) // 2
                h_segments.append((pi, si, xlo, xhi, y_avg))

    if not h_segments:
        return

    # ── 2. Cluster segments by similar y ──
    # Sort by y, then greedily cluster within Y_CLUSTER tolerance
    h_segments.sort(key=lambda s: s[4])
    Y_CLUSTER = nudge_emu  # cluster radius

    clusters: list[list[int]] = []  # each cluster = list of indices into h_segments
    cur_cluster: list[int] = [0]
    cur_y = h_segments[0][4]

    for i in range(1, len(h_segments)):
        if abs(h_segments[i][4] - cur_y) <= Y_CLUSTER:
            cur_cluster.append(i)
        else:
            clusters.append(cur_cluster)
            cur_cluster = [i]
            cur_y = h_segments[i][4]
    clusters.append(cur_cluster)

    # ── 3. Within each cluster, find x-overlapping groups and nudge ──
    for cluster in clusters:
        if len(cluster) <= 1:
            continue

        # Build overlap groups via greedy sweep
        # Sort cluster members by x_min
        members = sorted(cluster, key=lambda idx: h_segments[idx][2])

        # Find connected components of x-overlapping segments
        overlap_groups: list[list[int]] = []
        current_group = [members[0]]
        group_x_max = h_segments[members[0]][3]

        for mi in range(1, len(members)):
            idx = members[mi]
            seg = h_segments[idx]
            if seg[2] < group_x_max:  # x_min < current group's x_max → overlap
                current_group.append(idx)
                group_x_max = max(group_x_max, seg[3])
            else:
                overlap_groups.append(current_group)
                current_group = [idx]
                group_x_max = seg[3]
        overlap_groups.append(current_group)

        # For each overlap group with 2+ members, apply y-nudges
        for og in overlap_groups:
            if len(og) <= 1:
                continue

            n = len(og)
            # Center the stagger around the original y
            base_y = sum(h_segments[idx][4] for idx in og) // n

            for rank, idx in enumerate(og):
                # Spread: -half_span ... +half_span
                nudge = int((rank - (n - 1) / 2) * nudge_emu)
                target_y = base_y + nudge

                pi, si, _, _, orig_y = h_segments[idx]
                pts = planned_paths[pi]['path_points']

                # Shift y of both endpoints of this horizontal segment
                x1, y1 = pts[si]
                x2, y2 = pts[si + 1]
                pts[si] = (x1, target_y)
                pts[si + 1] = (x2, target_y)

                # Also update adjacent vertical segments to connect properly
                # Previous point (if exists) should now reach target_y
                if si > 0:
                    px, py = pts[si - 1]
                    # Only adjust if this was a vertical segment connecting to the h-seg
                    if abs(px - x1) <= Y_TOL:
                        pass  # vertical segment endpoint updated via pts[si]
                # Next point (if exists) should now start from target_y
                if si + 2 < len(pts):
                    nx, ny = pts[si + 2]
                    if abs(nx - x2) <= Y_TOL:
                        pass  # vertical segment start updated via pts[si+1]


def _draw_planned_path(slide, plan: dict, theme):
    """Draw a planned multi-segment edge path and optional label."""
    pts = plan['path_points']
    color = plan['color']
    width = plan['width']
    arrow = plan['arrow']
    dash = plan['dash']

    if len(pts) < 2:
        return

    # Draw all segments except the last one without arrows
    for i in range(len(pts) - 2):
        x1, y1 = pts[i]
        x2, y2 = pts[i + 1]
        _draw_line_xy(slide, x1, y1, x2, y2,
                      color=color, width=width, theme=theme, dash=dash)

    # Last segment gets the arrow
    x1, y1 = pts[-2]
    x2, y2 = pts[-1]
    _draw_line_xy(slide, x1, y1, x2, y2,
                  color=color, width=width, arrow=arrow, theme=theme, dash=dash)


def _bbox_overlap(ax, ay, aw, ah, bx, by, bw, bh) -> bool:
    """Return True if two axis-aligned rectangles overlap (EMU coords)."""
    return ax < bx + bw and ax + aw > bx and ay < by + bh and ay + ah > by


def _place_edge_label(
    path_points: list[tuple[int, int]],
    label_text: str,
    all_shapes: list,
    slide_w: int = None,
    slide_h: int = None,
) -> tuple[int, int, int, int]:
    """Compute (x, y, w, h) in EMU for an edge label, avoiding node overlap.

    Algorithm (案A+ — prioritised candidate list with boundary clamping):
      1. Decompose the connector path into segments.
      2. Generate candidate positions in priority order:
         a. L-route elbow exterior (if path has ≥3 points)
         b. Longest segment midpoint × [normal, small-normal] × [side A, side B]
         c. Longest segment 1/4 and 3/4 × same offsets
         d. All segments × midpoint × same offsets (if longest fails)
      3. Each candidate must pass:  no node overlap AND within slide bounds.
      4. Return the first passing candidate.
      5. Fallback: clamp the highest-priority candidate into slide bounds.
    """
    if slide_w is None:
        slide_w = Inches(10.0)
    if slide_h is None:
        slide_h = Inches(7.5)

    # ── Label dimensions ──
    char_count = len(label_text)
    lbl_w = Inches(max(0.4, min(char_count * 0.08, 1.4)))
    lbl_h = Inches(0.22)

    # Clearance levels: normal and tight (for crowded areas)
    _CLR_NORMAL = Inches(0.08)
    _CLR_TIGHT = Inches(0.03)

    # ── Segments ──
    segments = []
    for i in range(len(path_points) - 1):
        x1, y1 = path_points[i]
        x2, y2 = path_points[i + 1]
        length = ((x2 - x1) ** 2 + (y2 - y1) ** 2) ** 0.5
        segments.append((x1, y1, x2, y2, length))

    if not segments:
        return (Inches(0.1), Inches(0.1), lbl_w, lbl_h)

    # ── Helper: perpendicular offset vectors for a segment ──
    def _normal_offsets(x1, y1, x2, y2, clearance):
        seg_dx = x2 - x1
        seg_dy = y2 - y1
        seg_len = (seg_dx ** 2 + seg_dy ** 2) ** 0.5
        if seg_len < 1:
            return [(0, -clearance), (0, clearance)]
        nx = int(-seg_dy / seg_len * clearance)
        ny = int(seg_dx / seg_len * clearance)
        return [(nx, ny), (-nx, -ny)]

    def _candidate_at(px, py, offset_dx, offset_dy):
        return (px + offset_dx - lbl_w // 2, py + offset_dy - lbl_h // 2)

    # ── Node bounding boxes with padding ──
    # Add visual padding so labels don't "touch" nodes (looks like overlap)
    _NODE_PAD = Inches(0.06)
    node_bboxes = [
        (s.left - _NODE_PAD, s.top - _NODE_PAD,
         s.width + _NODE_PAD * 2, s.height + _NODE_PAD * 2)
        for s in all_shapes
    ]

    # Slide margin (small inset to avoid edge clipping)
    _MARGIN = Inches(0.05)

    def _is_valid(cx, cy):
        """Label must be inside slide AND not overlapping any node (with padding)."""
        # Slide boundary check
        if cx < _MARGIN or cy < _MARGIN:
            return False
        if cx + lbl_w > slide_w - _MARGIN or cy + lbl_h > slide_h - _MARGIN:
            return False
        # Node overlap check (padded bboxes)
        for (nx, ny, nw, nh) in node_bboxes:
            if _bbox_overlap(cx, cy, lbl_w, lbl_h, nx, ny, nw, nh):
                return False
        return True

    def _clamp(cx, cy):
        """Force label inside slide bounds."""
        cx = max(_MARGIN, min(cx, slide_w - lbl_w - _MARGIN))
        cy = max(_MARGIN, min(cy, slide_h - lbl_h - _MARGIN))
        return (cx, cy)

    # ── Generate prioritised candidates ──
    candidates = []

    # (a) L-route elbow exterior — highest priority
    if len(path_points) >= 3:
        for elbow_idx in range(1, len(path_points) - 1):
            ex, ey = path_points[elbow_idx]
            px, py = path_points[elbow_idx - 1]
            qx, qy = path_points[elbow_idx + 1]
            # Vector pointing towards the outer side of the bend
            mid_arm_x = (px + qx) // 2 - ex
            mid_arm_y = (py + qy) // 2 - ey
            arm_len = (mid_arm_x ** 2 + mid_arm_y ** 2) ** 0.5
            for dist in [_CLR_NORMAL + lbl_h // 2, _CLR_TIGHT + lbl_h // 2]:
                if arm_len > 1:
                    off_x = int(mid_arm_x / arm_len * dist)
                    off_y = int(mid_arm_y / arm_len * dist)
                else:
                    off_x, off_y = 0, -dist
                candidates.append(_candidate_at(ex, ey, off_x, off_y))
                # Also try opposite side of elbow
                candidates.append(_candidate_at(ex, ey, -off_x, -off_y))

    # (b,c) Segment-based candidates: longest first, then all others
    # Tier 1: normal/tight perpendicular offsets
    seg_order = sorted(range(len(segments)), key=lambda i: -segments[i][4])
    for seg_idx in seg_order:
        sx1, sy1, sx2, sy2, _ = segments[seg_idx]
        for clearance in [_CLR_NORMAL, _CLR_TIGHT]:
            offsets = _normal_offsets(sx1, sy1, sx2, sy2, clearance)
            for ratio in [0.5, 0.25, 0.75]:
                px = int(sx1 + (sx2 - sx1) * ratio)
                py = int(sy1 + (sy2 - sy1) * ratio)
                for odx, ody in offsets:
                    candidates.append(_candidate_at(px, py, odx, ody))

    # (d) Wide-offset candidates — for crowded areas where the label is
    #     wider than the gap between nodes.  Place label completely to
    #     one side of the connector, far enough to clear nearby nodes.
    #     Uses the path midpoint as anchor, with large perpendicular offset.
    path_mid_x = (path_points[0][0] + path_points[-1][0]) // 2
    path_mid_y = (path_points[0][1] + path_points[-1][1]) // 2
    for wide_dist in [lbl_w + _CLR_NORMAL, lbl_w // 2 + _CLR_NORMAL]:
        # Left / up (negative direction)
        candidates.append((path_mid_x - wide_dist - lbl_w // 2,
                           path_mid_y - lbl_h // 2))
        # Right / down (positive direction)
        candidates.append((path_mid_x + wide_dist - lbl_w // 2,
                           path_mid_y - lbl_h // 2))
        # Also try above / below the nodes
        candidates.append((path_mid_x - lbl_w // 2,
                           path_mid_y - wide_dist - lbl_h // 2))
        candidates.append((path_mid_x - lbl_w // 2,
                           path_mid_y + wide_dist - lbl_h // 2))

    # ── Select first valid candidate ──
    for (cx, cy) in candidates:
        if _is_valid(cx, cy):
            return (cx, cy, lbl_w, lbl_h)

    # ── Fallback: clamp first candidate into slide bounds ──
    cx, cy = _clamp(*candidates[0])
    return (cx, cy, lbl_w, lbl_h)


# ══════════════════════════════════════════════
# Group-aware connector routing (Step 13)
# ══════════════════════════════════════════════

def _classify_edge_route(from_id: str, to_id: str,
                          node_to_group: dict, group_bboxes: dict,
                          pos_map: dict, layer_map: dict,
                          direction: str, back_edges: set,
                          merge_targets: set) -> str:
    """Classify routing strategy based on node/group relationships.

    The routing decision is driven by the structural relationship between
    the source and target nodes' groups, not just geometric position:

      1. Back-edge (cycle)                          → 'back_edge'
      2. Cross-group (laterally-separated groups)   → 'cross_group'
      3. Inter-layer with significant axis offset   → 'l_route'
      4. Otherwise                                  → 'direct'
    """
    if (from_id, to_id) in back_edges:
        return 'back_edge'

    fp = pos_map.get(from_id)
    tp = pos_map.get(to_id)
    if not fp or not tp:
        return 'direct'

    fg = node_to_group.get(from_id)
    tg = node_to_group.get(to_id)
    from_layer = layer_map.get(from_id, 0)
    to_layer = layer_map.get(to_id, 0)
    is_inter_layer = from_layer != to_layer
    is_vertical = direction in ("TB", "BT")

    # ── Cross-group detection ──
    # Manhattan routing is reserved for connections between groups that are:
    #   (a) laterally separated (minimal cross-axis bbox overlap), AND
    #   (b) far apart in the layer hierarchy (Δlayer > 3).
    # This prevents normal tier-to-tier flow connections (Δlayer 1-2)
    # from using Manhattan routing — those cascade naturally with L-routes.
    layer_dist = abs(from_layer - to_layer)
    if (fg and tg and fg != tg
            and fg in group_bboxes and tg in group_bboxes
            and layer_dist > 3):
        fb = group_bboxes[fg]   # (min_x, min_y, max_x, max_y)
        tb = group_bboxes[tg]
        if is_vertical:
            h_overlap = min(fb[2], tb[2]) - max(fb[0], tb[0])
            min_w = min(fb[2] - fb[0], tb[2] - tb[0])
            if min_w > 0 and h_overlap < min_w * 0.3:
                return 'cross_group'
        else:
            v_overlap = min(fb[3], tb[3]) - max(fb[1], tb[1])
            min_h = min(fb[3] - fb[1], tb[3] - tb[1])
            if min_h > 0 and v_overlap < min_h * 0.3:
                return 'cross_group'

    # ── L-route detection ──
    if is_inter_layer:
        if is_vertical:
            h_offset = abs((fp.x + fp.w / 2) - (tp.x + tp.w / 2))
            if h_offset > fp.w * 0.3:
                return 'l_route'
        else:
            v_offset = abs((fp.y + fp.h / 2) - (tp.y + tp.h / 2))
            if v_offset > fp.h * 0.3:
                return 'l_route'

    # Flowchart merge targets with large offset
    if (to_id in merge_targets and is_vertical and is_inter_layer):
        h_offset = abs((fp.x + fp.w / 2) - (tp.x + tp.w / 2))
        if h_offset > 0.5:
            return 'l_route'

    return 'direct'


def _plan_manhattan_route(sa, sb, from_pos, to_pos,
                           fg_bbox, tg_bbox, direction: str,
                           src_port_off: float = 0.0,
                           tgt_port_off: float = 0.0) -> list:
    """Compute Manhattan-route path points WITHOUT drawing.

    Same routing logic as _draw_manhattan_route but only returns the
    list of (x_emu, y_emu) waypoints.  Used by the two-phase rendering
    pipeline so that nudging can adjust coordinates before drawing.
    """
    if direction in ("TB",):
        clear_y = Inches(fg_bbox[3] + 0.15)
        src_bx, src_by = _cp_coords(sa, 2, src_port_off)
        tgt_tx, tgt_ty = _cp_coords(sb, 0, tgt_port_off)

        if clear_y >= tgt_ty:
            src_cx = from_pos.x + from_pos.w / 2
            tgt_cx = to_pos.x + to_pos.w / 2
            if src_cx < tgt_cx:
                sx, sy = _cp_coords(sa, 1)
                tx, ty = _cp_coords(sb, 3)
                cx = Inches((fg_bbox[2] + tg_bbox[0]) / 2)
            else:
                sx, sy = _cp_coords(sa, 3)
                tx, ty = _cp_coords(sb, 1)
                cx = Inches((tg_bbox[2] + fg_bbox[0]) / 2)
            return [(sx, sy), (cx, sy), (cx, ty), (tx, ty)]

        return [(src_bx, src_by), (src_bx, clear_y),
                (tgt_tx, clear_y), (tgt_tx, tgt_ty)]

    elif direction in ("BT",):
        clear_y = Inches(fg_bbox[1] - 0.15)
        src_tx, src_ty = _cp_coords(sa, 0, src_port_off)
        tgt_bx, tgt_by = _cp_coords(sb, 2, tgt_port_off)

        if clear_y <= tgt_by:
            src_cx = from_pos.x + from_pos.w / 2
            tgt_cx = to_pos.x + to_pos.w / 2
            if src_cx < tgt_cx:
                sx, sy = _cp_coords(sa, 1)
                tx, ty = _cp_coords(sb, 3)
                cx = Inches((fg_bbox[2] + tg_bbox[0]) / 2)
            else:
                sx, sy = _cp_coords(sa, 3)
                tx, ty = _cp_coords(sb, 1)
                cx = Inches((tg_bbox[2] + fg_bbox[0]) / 2)
            return [(sx, sy), (cx, sy), (cx, ty), (tx, ty)]

        return [(src_tx, src_ty), (src_tx, clear_y),
                (tgt_bx, clear_y), (tgt_bx, tgt_by)]

    else:   # LR / RL
        src_cy = from_pos.y + from_pos.h / 2
        tgt_cy = to_pos.y + to_pos.h / 2

        if src_cy < tgt_cy:
            src_x, src_y = _cp_coords(sa, 2, src_port_off)
            tgt_x, tgt_y = _cp_coords(sb, 0, tgt_port_off)
            corridor_y = Inches((fg_bbox[3] + tg_bbox[1]) / 2)
        else:
            src_x, src_y = _cp_coords(sa, 0, src_port_off)
            tgt_x, tgt_y = _cp_coords(sb, 2, tgt_port_off)
            corridor_y = Inches((tg_bbox[3] + fg_bbox[1]) / 2)

        return [(src_x, src_y), (src_x, corridor_y),
                (tgt_x, corridor_y), (tgt_x, tgt_y)]


def _draw_manhattan_route(slide, sa, sb, from_pos, to_pos,
                           fg_bbox, tg_bbox, direction: str,
                           color, width, arrow: bool,
                           theme, dash: int = 0,
                           src_port_off: float = 0.0,
                           tgt_port_off: float = 0.0) -> list:
    """Draw Manhattan-routed connector for cross-group connections.

    For TB/BT: exits the source along the **flow axis** (bottom for TB,
    top for BT), routes horizontally at a y-level just past the source
    group boundary, then drops vertically into the target.  This avoids
    the previous approach (side-exit + vertical corridor) which would cut
    through intermediate groups in a cascading diagonal layout.

    For LR/RL: exits source from the side facing the target, routes
    through a corridor between groups (vertical midpoint).

    Returns list of path points (EMU) for label placement.
    """
    if direction in ("TB",):
        # Route: source BOTTOM → short vertical to clear_y
        #      → horizontal at clear_y → vertical drop to target TOP
        # clear_y sits just below the source group, above all tiers below it.
        clear_y = Inches(fg_bbox[3] + 0.15)
        src_bx, src_by = _cp_coords(sa, 2, src_port_off)   # source bottom
        tgt_tx, tgt_ty = _cp_coords(sb, 0, tgt_port_off)   # target top

        if clear_y >= tgt_ty:
            # Edge case: source group bottom is at or below target top.
            # Fall back to side-exit corridor approach.
            src_cx = from_pos.x + from_pos.w / 2
            tgt_cx = to_pos.x + to_pos.w / 2
            if src_cx < tgt_cx:
                sx, sy = _cp_coords(sa, 1)
                tx, ty = _cp_coords(sb, 3)
                cx = Inches((fg_bbox[2] + tg_bbox[0]) / 2)
            else:
                sx, sy = _cp_coords(sa, 3)
                tx, ty = _cp_coords(sb, 1)
                cx = Inches((tg_bbox[2] + fg_bbox[0]) / 2)
            _draw_line_xy(slide, sx, sy, cx, sy,
                          color=color, width=width, theme=theme, dash=dash)
            _draw_line_xy(slide, cx, sy, cx, ty,
                          color=color, width=width, theme=theme, dash=dash)
            _draw_line_xy(slide, cx, ty, tx, ty,
                          color=color, width=width, arrow=arrow, theme=theme, dash=dash)
            return [(sx, sy), (cx, sy), (cx, ty), (tx, ty)]

        # Normal 3-segment route (avoids intermediate groups)
        _draw_line_xy(slide, src_bx, src_by, src_bx, clear_y,
                      color=color, width=width, theme=theme, dash=dash)
        _draw_line_xy(slide, src_bx, clear_y, tgt_tx, clear_y,
                      color=color, width=width, theme=theme, dash=dash)
        _draw_line_xy(slide, tgt_tx, clear_y, tgt_tx, tgt_ty,
                      color=color, width=width, arrow=arrow, theme=theme, dash=dash)
        return [(src_bx, src_by), (src_bx, clear_y),
                (tgt_tx, clear_y), (tgt_tx, tgt_ty)]

    elif direction in ("BT",):
        # Reverse: source TOP → clear_y above source group → horizontal → target BOTTOM
        clear_y = Inches(fg_bbox[1] - 0.15)
        src_tx, src_ty = _cp_coords(sa, 0, src_port_off)   # source top
        tgt_bx, tgt_by = _cp_coords(sb, 2, tgt_port_off)   # target bottom

        if clear_y <= tgt_by:
            # Fallback: side-exit corridor
            src_cx = from_pos.x + from_pos.w / 2
            tgt_cx = to_pos.x + to_pos.w / 2
            if src_cx < tgt_cx:
                sx, sy = _cp_coords(sa, 1)
                tx, ty = _cp_coords(sb, 3)
                cx = Inches((fg_bbox[2] + tg_bbox[0]) / 2)
            else:
                sx, sy = _cp_coords(sa, 3)
                tx, ty = _cp_coords(sb, 1)
                cx = Inches((tg_bbox[2] + fg_bbox[0]) / 2)
            _draw_line_xy(slide, sx, sy, cx, sy,
                          color=color, width=width, theme=theme, dash=dash)
            _draw_line_xy(slide, cx, sy, cx, ty,
                          color=color, width=width, theme=theme, dash=dash)
            _draw_line_xy(slide, cx, ty, tx, ty,
                          color=color, width=width, arrow=arrow, theme=theme, dash=dash)
            return [(sx, sy), (cx, sy), (cx, ty), (tx, ty)]

        _draw_line_xy(slide, src_tx, src_ty, src_tx, clear_y,
                      color=color, width=width, theme=theme, dash=dash)
        _draw_line_xy(slide, src_tx, clear_y, tgt_bx, clear_y,
                      color=color, width=width, theme=theme, dash=dash)
        _draw_line_xy(slide, tgt_bx, clear_y, tgt_bx, tgt_by,
                      color=color, width=width, arrow=arrow, theme=theme, dash=dash)
        return [(src_tx, src_ty), (src_tx, clear_y),
                (tgt_bx, clear_y), (tgt_bx, tgt_by)]

    else:   # LR / RL — side-exit corridor (groups stack vertically)
        src_cy = from_pos.y + from_pos.h / 2
        tgt_cy = to_pos.y + to_pos.h / 2

        if src_cy < tgt_cy:
            src_x, src_y = _cp_coords(sa, 2, src_port_off)    # bottom
            tgt_x, tgt_y = _cp_coords(sb, 0, tgt_port_off)    # top
            corridor_y = Inches((fg_bbox[3] + tg_bbox[1]) / 2)
        else:
            src_x, src_y = _cp_coords(sa, 0, src_port_off)    # top
            tgt_x, tgt_y = _cp_coords(sb, 2, tgt_port_off)    # bottom
            corridor_y = Inches((tg_bbox[3] + fg_bbox[1]) / 2)

        _draw_line_xy(slide, src_x, src_y, src_x, corridor_y,
                      color=color, width=width, theme=theme, dash=dash)
        _draw_line_xy(slide, src_x, corridor_y, tgt_x, corridor_y,
                      color=color, width=width, theme=theme, dash=dash)
        _draw_line_xy(slide, tgt_x, corridor_y, tgt_x, tgt_y,
                      color=color, width=width, arrow=arrow, theme=theme, dash=dash)

        return [(src_x, src_y), (src_x, corridor_y),
                (tgt_x, corridor_y), (tgt_x, tgt_y)]


def _draw_fan_in_bus(slide, target_id, source_edges, shape_map, pos_map,
                      direction, is_flowchart, theme, ds, layout_scale, fonts) -> set:
    """Draw fan-in bus: multiple sources → shared horizontal bus → single target.

    Visual (TB):
        [s1]  [s2]  [s3]
          |     |     |      ← individual stubs
          +-----+-----+      ← shared horizontal bus
                |
             [target]

    Returns set of (from_id, to_id) tuples that were drawn as bus lines.
    """
    sb = shape_map.get(target_id)
    tp = pos_map.get(target_id)
    if not sb or not tp:
        return set()

    sources = []
    for e in source_edges:
        sa = shape_map.get(e.from_id)
        fp = pos_map.get(e.from_id)
        if sa and fp:
            sources.append((e, sa, fp))

    if len(sources) < 3:
        return set()

    handled = set()

    if direction in ("TB",):
        sources.sort(key=lambda s: s[2].x)

        # Bus y = midpoint between deepest source bottom and target top
        src_bottoms = [fp.y + fp.h for _, _, fp in sources]
        bus_y = Inches((max(src_bottoms) + tp.y) / 2)

        tgt_tx, tgt_ty = _cp_coords(sb, 0)   # target top center

        # Use first edge style for bus line color/width
        bus_es = sources[0][0].effective_style()
        bus_color = _hex_to_rgb(bus_es.color)
        bus_width = Pt(bus_es.width)

        # Individual stubs from each source down to bus
        stub_xs = []
        for e, sa, fp in sources:
            es = e.effective_style()
            src_bx, src_by = _cp_coords(sa, 2)   # bottom center
            dash = 2 if es.dash else 0
            _draw_line_xy(slide, src_bx, src_by, src_bx, bus_y,
                          color=_hex_to_rgb(es.color), width=Pt(es.width),
                          theme=theme, dash=dash)
            stub_xs.append(src_bx)

            # Edge label on stub
            if e.label:
                edge_fs = _scaled_font_size(ds.edge_label_font_size, layout_scale)
                lbl = slide.shapes.add_textbox(
                    src_bx + Inches(0.02),
                    (src_by + bus_y) // 2 - Inches(0.08),
                    Inches(0.5), Inches(0.18))
                lp = lbl.text_frame.paragraphs[0]
                lp.text = e.label
                _set_font(lp, fonts.body, edge_fs, es.color, bold=True)

            handled.add((e.from_id, e.to_id))

        # Horizontal bus spanning all stubs + target position
        all_xs = stub_xs + [tgt_tx]
        _draw_line_xy(slide, min(all_xs), bus_y, max(all_xs), bus_y,
                      color=bus_color, width=bus_width, theme=theme)

        # Vertical drop from bus to target
        bus_arrow = bus_es.arrow if is_flowchart else False
        _draw_line_xy(slide, tgt_tx, bus_y, tgt_tx, tgt_ty,
                      color=bus_color, width=bus_width,
                      arrow=bus_arrow, theme=theme)

    elif direction in ("BT",):
        sources.sort(key=lambda s: s[2].x)
        src_tops = [fp.y for _, _, fp in sources]
        bus_y = Inches((min(src_tops) + tp.y + tp.h) / 2)
        tgt_bx, tgt_by = _cp_coords(sb, 2)

        bus_es = sources[0][0].effective_style()
        bus_color = _hex_to_rgb(bus_es.color)
        bus_width = Pt(bus_es.width)

        stub_xs = []
        for e, sa, fp in sources:
            es = e.effective_style()
            src_tx, src_ty = _cp_coords(sa, 0)
            dash = 2 if es.dash else 0
            _draw_line_xy(slide, src_tx, src_ty, src_tx, bus_y,
                          color=_hex_to_rgb(es.color), width=Pt(es.width),
                          theme=theme, dash=dash)
            stub_xs.append(src_tx)
            handled.add((e.from_id, e.to_id))

        all_xs = stub_xs + [tgt_bx]
        _draw_line_xy(slide, min(all_xs), bus_y, max(all_xs), bus_y,
                      color=bus_color, width=bus_width, theme=theme)
        bus_arrow = bus_es.arrow if is_flowchart else False
        _draw_line_xy(slide, tgt_bx, bus_y, tgt_bx, tgt_by,
                      color=bus_color, width=bus_width,
                      arrow=bus_arrow, theme=theme)

    return handled


def _draw_fan_out_bus(slide, source_id, target_edges, shape_map, pos_map,
                       direction, is_flowchart, theme, ds, layout_scale, fonts) -> set:
    """Draw fan-out bus: single source → shared horizontal bus → multiple targets.

    Visual (TB):
             [source]
                |
          +-----+-----+      ← shared horizontal bus
          |     |     |      ← individual stubs
        [t1]  [t2]  [t3]

    Returns set of (from_id, to_id) tuples that were drawn as bus lines.
    """
    sa = shape_map.get(source_id)
    fp = pos_map.get(source_id)
    if not sa or not fp:
        return set()

    targets = []
    for e in target_edges:
        sb = shape_map.get(e.to_id)
        tp = pos_map.get(e.to_id)
        if sb and tp:
            targets.append((e, sb, tp))

    if len(targets) < 3:
        return set()

    handled = set()

    if direction in ("TB",):
        targets.sort(key=lambda t: t[2].x)

        # Bus y = midpoint between source bottom and shallowest target top
        tgt_tops = [tp.y for _, _, tp in targets]
        bus_y = Inches((fp.y + fp.h + min(tgt_tops)) / 2)

        src_bx, src_by = _cp_coords(sa, 2)   # source bottom center

        bus_es = targets[0][0].effective_style()
        bus_color = _hex_to_rgb(bus_es.color)
        bus_width = Pt(bus_es.width)

        # Individual stubs from bus down to each target
        stub_xs = []
        for e, sb, tp in targets:
            es = e.effective_style()
            tgt_tx, tgt_ty = _cp_coords(sb, 0)   # top center
            dash = 2 if es.dash else 0
            tgt_arrow = es.arrow if is_flowchart else False
            _draw_line_xy(slide, tgt_tx, bus_y, tgt_tx, tgt_ty,
                          color=_hex_to_rgb(es.color), width=Pt(es.width),
                          arrow=tgt_arrow, theme=theme, dash=dash)
            stub_xs.append(tgt_tx)

            if e.label:
                edge_fs = _scaled_font_size(ds.edge_label_font_size, layout_scale)
                lbl = slide.shapes.add_textbox(
                    tgt_tx + Inches(0.02),
                    (bus_y + tgt_ty) // 2 - Inches(0.08),
                    Inches(0.5), Inches(0.18))
                lp = lbl.text_frame.paragraphs[0]
                lp.text = e.label
                _set_font(lp, fonts.body, edge_fs, es.color, bold=True)

            handled.add((e.from_id, e.to_id))

        # Horizontal bus spanning source + all targets
        all_xs = stub_xs + [src_bx]
        _draw_line_xy(slide, min(all_xs), bus_y, max(all_xs), bus_y,
                      color=bus_color, width=bus_width, theme=theme)

        # Vertical from source down to bus
        _draw_line_xy(slide, src_bx, src_by, src_bx, bus_y,
                      color=bus_color, width=bus_width, theme=theme)

    elif direction in ("BT",):
        targets.sort(key=lambda t: t[2].x)
        tgt_bottoms = [tp.y + tp.h for _, _, tp in targets]
        bus_y = Inches((fp.y + max(tgt_bottoms)) / 2)
        src_tx, src_ty = _cp_coords(sa, 0)

        bus_es = targets[0][0].effective_style()
        bus_color = _hex_to_rgb(bus_es.color)
        bus_width = Pt(bus_es.width)

        stub_xs = []
        for e, sb, tp in targets:
            es = e.effective_style()
            tgt_bx, tgt_by = _cp_coords(sb, 2)
            dash = 2 if es.dash else 0
            tgt_arrow = es.arrow if is_flowchart else False
            _draw_line_xy(slide, tgt_bx, bus_y, tgt_bx, tgt_by,
                          color=_hex_to_rgb(es.color), width=Pt(es.width),
                          arrow=tgt_arrow, theme=theme, dash=dash)
            stub_xs.append(tgt_bx)
            handled.add((e.from_id, e.to_id))

        all_xs = stub_xs + [src_tx]
        _draw_line_xy(slide, min(all_xs), bus_y, max(all_xs), bus_y,
                      color=bus_color, width=bus_width, theme=theme)
        _draw_line_xy(slide, src_tx, src_ty, src_tx, bus_y,
                      color=bus_color, width=bus_width, theme=theme)

    return handled


def _set_arrow(connector, end='tail', arrow_type='arrow'):
    """Add an arrowhead via OOXML."""
    ln = connector._element.find(f'.//{{{_a}}}ln')
    if ln is None:
        return
    tag = f'{{{_a}}}tailEnd' if end == 'tail' else f'{{{_a}}}headEnd'
    existing = ln.find(tag)
    if existing is not None:
        ln.remove(existing)
    el = etree.SubElement(ln, tag)
    el.set('type', arrow_type)
    el.set('w', 'med')
    el.set('len', 'med')


def _detect_cp(shape_a, shape_b, direction: str = "TB",
               layer_a: int = 0, layer_b: int = 0) -> tuple[int, int]:
    """Auto-detect best connection point pair based on direction and layers.

    For inter-layer connections (parent→child), strongly prefers the
    direction-aligned axis:
      TB/BT → vertical (bottom→top or top→bottom)
      LR/RL → horizontal (right→left or left→right)

    For same-layer connections, uses relative position to pick best side.

    cp indices: 0=top, 1=right, 2=bottom, 3=left
    """
    ax = shape_a.left + shape_a.width // 2
    ay = shape_a.top + shape_a.height // 2
    bx = shape_b.left + shape_b.width // 2
    by = shape_b.top + shape_b.height // 2
    dx, dy = bx - ax, by - ay

    is_inter_layer = (layer_a != layer_b)

    if is_inter_layer:
        # Inter-layer: use direction-aligned axis
        if direction in ("TB", "BT"):
            # Vertical: parent bottom → child top (or reversed for BT)
            if dy > 0:
                return (2, 0)  # A below → B above  (A.bottom → B.top)
            else:
                return (0, 2)  # A above → B below
        else:  # LR, RL
            if dx > 0:
                return (1, 3)  # A left of B → A.right → B.left
            else:
                return (3, 1)
    else:
        # Same layer: use relative position (original logic)
        if abs(dx) > abs(dy):
            return (1, 3) if dx > 0 else (3, 1)
        return (2, 0) if dy > 0 else (0, 2)


def _is_diamond(shape) -> bool:
    try:
        return shape._element.spPr.prstGeom.attrib.get('prst', '') == 'diamond'
    except Exception:
        return False


def _draw_connector(slide, shape_a, shape_b,
                    begin_cp=None, end_cp=None,
                    color=None, width=Pt(2), arrow=False, snap=True,
                    theme: ThemeConfig = DEFAULT_THEME):
    """Draw a connector between two shapes. Returns the connector shape."""
    if begin_cp is None or end_cp is None:
        begin_cp, end_cp = _detect_cp(shape_a, shape_b)

    x1, y1 = _cp_coords(shape_a, begin_cp)
    x2, y2 = _cp_coords(shape_b, end_cp)

    conn = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT, x1, y1, x2, y2
    )

    if snap:
        try:
            conn.begin_connect(shape_a, begin_cp)
            conn.end_connect(shape_b, end_cp)
        except Exception:
            pass

    conn.line.color.rgb = color or _hex_to_rgb(theme.diagram_style.edge_color)
    conn.line.width = width

    if arrow:
        _set_arrow(conn)

    return conn


def _draw_line_xy(slide, x1, y1, x2, y2, color=None, width=Pt(1.5), arrow=False,
                  theme: ThemeConfig = DEFAULT_THEME, dash: int = 0):
    """Draw a simple line between two EMU coordinates (no snap).

    Args:
        dash: Dash style (0=solid, 2=dash, 3=dash-dot, etc.).
              Maps to MSO_LINE_DASH_STYLE values.
    """
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT, x1, y1, x2, y2
    )
    conn.line.color.rgb = color or _hex_to_rgb(theme.diagram_style.edge_color)
    conn.line.width = width
    if dash:
        conn.line.dash_style = dash
    if arrow:
        _set_arrow(conn)
    return conn


def _set_font(run_or_para, font_name: str, font_size: float,
              font_color: str, bold: bool = False):
    """Apply consistent font settings to a run or paragraph's font."""
    run_or_para.font.name = font_name
    run_or_para.font.size = Pt(font_size)
    run_or_para.font.color.rgb = _hex_to_rgb(font_color)
    run_or_para.font.bold = bold


def _scaled_font_size(base_size: float, scale: float, min_size: float = 5.0) -> float:
    """Compute font size adjusted for layout scale factor.

    When the layout is shrunk (scale < 1.0), fonts are proportionally
    reduced but never below min_size.
    """
    if scale >= 1.0:
        return base_size
    return max(base_size * scale, min_size)


def _enable_auto_shrink(text_frame):
    """Enable PowerPoint's built-in text auto-shrink for a text frame.

    This sets fontAutofit so that PowerPoint automatically reduces the
    font size if text exceeds the shape bounds.
    """
    try:
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        # Fallback: set via lxml if python-pptx property doesn't work
        try:
            txBody = text_frame._txBody
            bodyPr = txBody.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
            if bodyPr is not None:
                # Remove existing autofit elements
                for child in list(bodyPr):
                    tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                    if tag in ('noAutofit', 'spAutoFit', 'normAutofit'):
                        bodyPr.remove(child)
                # Add normAutofit (font shrink to fit)
                ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                etree.SubElement(bodyPr, f'{{{ns}}}normAutofit')
        except Exception:
            pass


# ══════════════════════════════════════════════
# Header Bar (template-style)
# ══════════════════════════════════════════════

def _draw_header_bar(slide, title: str, subtitle: str = "",
                     theme: ThemeConfig = DEFAULT_THEME):
    """Draw a Midnight Executive-style header bar at the top of the slide.

    Matches the template's _header_bar(): navy background, white title,
    ice_blue subtitle.
    """
    ds = theme.diagram_style
    fonts = theme.fonts

    # Navy background bar (full width, 1.15" tall)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(SLIDE_W), Inches(1.15)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = _hex_to_rgb(ds.header_bar_color)
    bar.line.fill.background()

    # Left accent bar (Midnight Executive motif)
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(0.08), Inches(1.15)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = _hex_to_rgb(f"#{theme.palette.accent}")
    accent_bar.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.18), Inches(10), Inches(0.55)
    )
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    _set_font(p, fonts.heading, 28, ds.header_font_color, bold=True)

    # Subtitle text
    if subtitle:
        sub_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.72), Inches(10), Inches(0.35)
        )
        sp = sub_box.text_frame.paragraphs[0]
        sp.text = subtitle
        _set_font(sp, fonts.body, 12, ds.header_subtitle_color, bold=False)


def _draw_swimlanes(slide, lane_infos: list[LaneInfo], direction: str,
                     content_top: float, theme: ThemeConfig = DEFAULT_THEME):
    """Draw swimlane bands, headers, and separator lines.

    For TB/BT: vertical columns with headers at top.
    For LR/RL: horizontal rows with headers at left.
    """
    is_horizontal = direction in ("LR", "RL")
    fonts = theme.fonts
    lane_header_size = 0.6  # matches layout algorithm

    # Alternating band colors for readability
    band_colors = ["#F8FAFC", "#EFF3F8"]

    for i, li in enumerate(lane_infos):
        if li.lane_id == "__default__":
            continue  # don't draw band for unassigned nodes

        ls = li.style
        band_fill = ls.band_fill or band_colors[i % 2]

        if is_horizontal:
            # ── Horizontal lanes (LR/RL): rows ──
            # Band: full width, within lane's Y range
            band = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(li.cross_origin),
                Inches(SLIDE_W), Inches(li.cross_size),
            )
            band.fill.solid()
            band.fill.fore_color.rgb = _hex_to_rgb(band_fill)
            band.line.fill.background()

            # Header: vertical strip at left
            header = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(li.cross_origin),
                Inches(lane_header_size), Inches(li.cross_size),
            )
            header.fill.solid()
            header.fill.fore_color.rgb = _hex_to_rgb(ls.header_fill)
            header.line.fill.background()

            # Header label (vertical centering)
            hdr_txt = slide.shapes.add_textbox(
                Inches(0.05), Inches(li.cross_origin),
                Inches(lane_header_size - 0.1), Inches(li.cross_size),
            )
            hdr_txt.text_frame.word_wrap = True
            p = hdr_txt.text_frame.paragraphs[0]
            p.text = li.label
            _set_font(p, fonts.heading, 11, ls.header_font_color, bold=True)
            p.alignment = PP_ALIGN.CENTER
            # Vertical center
            hdr_txt.text_frame.paragraphs[0].space_before = Pt(0)

            # Separator line at bottom of lane
            if i < len(lane_infos) - 1:
                sep_y = li.cross_origin + li.cross_size
                _draw_line_xy(slide,
                    Inches(0), Inches(sep_y),
                    Inches(SLIDE_W), Inches(sep_y),
                    color=_hex_to_rgb(ls.border), width=Pt(ls.border_width), theme=theme)

        else:
            # ── Vertical lanes (TB/BT): columns ──
            # Band: full height below header, within lane's X range
            band_top = content_top
            band_height = SLIDE_H - band_top

            band = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(li.cross_origin), Inches(band_top),
                Inches(li.cross_size), Inches(band_height),
            )
            band.fill.solid()
            band.fill.fore_color.rgb = _hex_to_rgb(band_fill)
            band.line.fill.background()

            # Header: horizontal strip at top (below slide header bar)
            header = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(li.cross_origin), Inches(content_top),
                Inches(li.cross_size), Inches(lane_header_size),
            )
            header.fill.solid()
            header.fill.fore_color.rgb = _hex_to_rgb(ls.header_fill)
            header.line.fill.background()

            # Header label
            hdr_txt = slide.shapes.add_textbox(
                Inches(li.cross_origin), Inches(content_top + 0.1),
                Inches(li.cross_size), Inches(lane_header_size - 0.2),
            )
            p = hdr_txt.text_frame.paragraphs[0]
            p.text = li.label
            _set_font(p, fonts.heading, 11, ls.header_font_color, bold=True)
            p.alignment = PP_ALIGN.CENTER

            # Separator line at right of lane
            if i < len(lane_infos) - 1:
                sep_x = li.cross_origin + li.cross_size
                _draw_line_xy(slide,
                    Inches(sep_x), Inches(band_top),
                    Inches(sep_x), Inches(SLIDE_H),
                    color=_hex_to_rgb(ls.border), width=Pt(ls.border_width), theme=theme)


# ══════════════════════════════════════════════
# Post-layout: Group Overlap Separation
# ══════════════════════════════════════════════

def _separate_overlapping_groups(
    spec: DiagramSpec,
    positions: list,  # list of NodePosition
    min_gap: float = 0.15,
    slide_width: float = 13.333,
    slide_height: float = 7.5,
    margin: float = 0.3,
    margin_top: Optional[float] = None,
) -> list:
    """Push apart non-nested groups whose bounding boxes overlap (Algorithm C).

    Handles BOTH axes:
      - X-axis: groups with y-overlap need minimum x-gap
      - Y-axis: groups with x-overlap need minimum y-gap

    Uses scaled padding (proportional to layout_scale) to compute realistic bboxes.

    Returns updated list of NodePosition with adjusted coordinates.
    """
    if not spec.groups or len(positions) < 2:
        return positions

    layout_scale = positions[0].scale if positions else 1.0
    _pad_s = max(layout_scale, 0.4)
    _BASE_PAD = {0: 0.25, 1: 0.18, 2: 0.12}
    _BASE_LABEL = 0.25
    scaled_pads = {k: max(0.06, v * _pad_s) for k, v in _BASE_PAD.items()}
    scaled_label = max(0.10, _BASE_LABEL * _pad_s)

    pos_map = {p.node_id: p for p in positions}
    node_group = {n.id: n.group for n in spec.nodes if n.group}

    def _compute_bboxes() -> dict[str, list[float]]:
        """Compute group bboxes with scaled padding."""
        bboxes: dict[str, list[float]] = {}
        for grp in sorted(spec.groups, key=lambda g: spec.group_depth(g.id), reverse=True):
            depth = spec.group_depth(grp.id)
            pad = scaled_pads.get(depth, 0.06)
            members = [n.id for n in spec.nodes if n.group == grp.id]
            pts_x, pts_y, pts_xe, pts_ye = [], [], [], []
            for nid in members:
                if nid in pos_map:
                    p = pos_map[nid]
                    pts_x.append(p.x); pts_y.append(p.y)
                    pts_xe.append(p.x + p.w); pts_ye.append(p.y + p.h)
            for child_gid in spec.group_children(grp.id):
                if child_gid in bboxes:
                    cb = bboxes[child_gid]
                    pts_x.append(cb[0]); pts_y.append(cb[1])
                    pts_xe.append(cb[2]); pts_ye.append(cb[3])
            if pts_x:
                bboxes[grp.id] = [
                    min(pts_x) - pad,
                    min(pts_y) - pad - scaled_label,
                    max(pts_xe) + pad,
                    max(pts_ye) + pad,
                ]
        return bboxes

    # Identify top-level groups
    top_level = [g for g in spec.groups if g.parent is None]

    # Map nodes to their top-level group
    def _find_top_level(gid: str) -> Optional[str]:
        visited = set()
        while gid and gid not in visited:
            visited.add(gid)
            grp_obj = next((g for g in spec.groups if g.id == gid), None)
            if not grp_obj:
                return None
            if grp_obj.parent is None:
                return gid
            gid = grp_obj.parent
        return None

    node_top_group: dict[str, str] = {}
    for nid, gid in node_group.items():
        tl = _find_top_level(gid)
        if tl:
            node_top_group[nid] = tl

    # ── Iterative overlap resolution (max 5 iterations) ──
    for iteration in range(5):
        bboxes = _compute_bboxes()
        tl_with_bbox = [g for g in top_level if g.id in bboxes]
        if len(tl_with_bbox) < 2:
            break

        # Find worst overlap
        worst_overlap = 0.0
        worst_pair = None
        worst_axis = None  # 'x' or 'y'

        for i, g1 in enumerate(tl_with_bbox):
            for g2 in tl_with_bbox[i+1:]:
                b1, b2 = bboxes[g1.id], bboxes[g2.id]
                x_ov = max(0, min(b1[2], b2[2]) - max(b1[0], b2[0]))
                y_ov = max(0, min(b1[3], b2[3]) - max(b1[1], b2[1]))
                if x_ov > 0 and y_ov > 0:
                    # True overlap — resolve along the axis with less overlap
                    resolve_amount = min(x_ov, y_ov) + min_gap
                    if resolve_amount > worst_overlap:
                        worst_overlap = resolve_amount
                        worst_pair = (g1, g2)
                        worst_axis = 'x' if x_ov < y_ov else 'y'

        if worst_pair is None:
            break  # No overlaps!

        g1, g2 = worst_pair
        b1, b2 = bboxes[g1.id], bboxes[g2.id]

        if worst_axis == 'x':
            # Push apart on x-axis
            x_ov = max(0, min(b1[2], b2[2]) - max(b1[0], b2[0]))
            push = (x_ov + min_gap) / 2
            c1 = (b1[0] + b1[2]) / 2
            c2 = (b2[0] + b2[2]) / 2
            if c1 <= c2:
                shift_1, shift_2 = -push, push
            else:
                shift_1, shift_2 = push, -push
            for p in positions:
                tl = node_top_group.get(p.node_id)
                if tl == g1.id:
                    p.x += shift_1
                elif tl == g2.id:
                    p.x += shift_2
        else:
            # Push apart on y-axis
            y_ov = max(0, min(b1[3], b2[3]) - max(b1[1], b2[1]))
            push = (y_ov + min_gap) / 2
            c1 = (b1[1] + b1[3]) / 2
            c2 = (b2[1] + b2[3]) / 2
            if c1 <= c2:
                shift_1, shift_2 = -push, push
            else:
                shift_1, shift_2 = push, -push
            for p in positions:
                tl = node_top_group.get(p.node_id)
                if tl == g1.id:
                    p.y += shift_1
                elif tl == g2.id:
                    p.y += shift_2

    # ── Rescale to fit slide ──
    all_x = [p.x for p in positions]
    all_xe = [p.x + p.w for p in positions]
    all_y = [p.y for p in positions]
    all_ye = [p.y + p.h for p in positions]

    # X-axis fit
    cur_min_x, cur_max_x = min(all_x), max(all_xe)
    cur_w = cur_max_x - cur_min_x
    avail_w = slide_width - 2 * margin
    if cur_w > avail_w and cur_w > 0:
        sx = avail_w / cur_w
        for p in positions:
            p.x = margin + (p.x - cur_min_x) * sx
            p.w *= sx
            p.scale *= sx
    elif cur_min_x < margin or cur_max_x > slide_width - margin:
        off = margin - cur_min_x + (avail_w - cur_w) / 2
        for p in positions:
            p.x += off

    # Y-axis fit (respect header bar via margin_top)
    y_top = margin_top if margin_top is not None else margin
    y_bot = margin
    all_y = [p.y for p in positions]
    all_ye = [p.y + p.h for p in positions]
    cur_min_y, cur_max_y = min(all_y), max(all_ye)
    cur_h = cur_max_y - cur_min_y
    avail_h = slide_height - y_top - y_bot
    if cur_h > avail_h and cur_h > 0:
        sy = avail_h / cur_h
        for p in positions:
            p.y = y_top + (p.y - cur_min_y) * sy
            p.h *= sy
            p.scale *= sy
    elif cur_min_y < y_top or cur_max_y > slide_height - y_bot:
        off = y_top - cur_min_y + (avail_h - cur_h) / 2
        for p in positions:
            p.y += off

    return positions


# ══════════════════════════════════════════════
# Main Renderer
# ══════════════════════════════════════════════

def render_diagram(spec: DiagramSpec,
                   prs: Optional[Presentation] = None,
                   slide_layout_idx: int = 6,
                   template_path: Optional[str] = None,
                   theme: ThemeConfig = DEFAULT_THEME,
                   use_header_bar: bool = True) -> Presentation:
    """Render a DiagramSpec to a PPTX slide.

    Args:
        spec: Validated DiagramSpec
        prs: Existing Presentation to add slide to (creates new if None)
        slide_layout_idx: Slide layout index (6 = blank)
        template_path: Path to template PPTX for layout inheritance.
                       If provided and prs is None, loads the template.
        theme: ThemeConfig for colors and fonts
        use_header_bar: If True, draw a themed header bar at top

    Returns:
        The Presentation object
    """
    ds = theme.diagram_style
    fonts = theme.fonts

    if prs is None:
        if template_path:
            prs = Presentation(template_path)
        else:
            prs = Presentation()
            prs.slide_width = Emu(12192000)   # 13.333"
            prs.slide_height = Emu(6858000)   # 7.5"

    slide = prs.slides.add_slide(prs.slide_layouts[slide_layout_idx])

    # ── Background (non-template only) ──
    if not template_path and ds.slide_bg:
        bg = slide.background
        bg_fill = bg.fill
        bg_fill.solid()
        bg_fill.fore_color.rgb = _hex_to_rgb(ds.slide_bg)

    # ── Header bar ──
    content_top = 0.8  # default when no header bar
    if use_header_bar and spec.title:
        _draw_header_bar(slide, spec.title, theme=theme)
        content_top = 1.35  # below header bar
    elif spec.title:
        # Simple title textbox (no header bar)
        txBox = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.15), Inches(10), Inches(0.45)
        )
        p = txBox.text_frame.paragraphs[0]
        p.text = spec.title
        _set_font(p, fonts.heading, ds.title_font_size,
                  ds.title_font_color, bold=ds.title_font_bold)

    # ── Compute layout ──
    lane_infos: list[LaneInfo] = []
    if spec.lanes:
        positions, lane_infos = _compute_layout_swimlane(spec, content_top)
    else:
        positions = compute_layout(spec, content_top=content_top)

    # ── Post-layout: separate overlapping groups ──
    if spec.groups and not spec.lanes:
        positions = _separate_overlapping_groups(
            spec, positions, margin_top=content_top)

    pos_map = {p.node_id: p for p in positions}

    # ── Determine global scale factor for font adjustment ──
    layout_scale = positions[0].scale if positions else 1.0

    # ── Draw swimlane bands (before nodes, so they appear behind) ──
    if lane_infos:
        _draw_swimlanes(slide, lane_infos, spec.direction, content_top, theme)

    # ── Draw shapes ──
    shape_map = {}  # node_id → pptx shape object
    node_map = {n.id: n for n in spec.nodes}

    for pos in positions:
        node = node_map[pos.node_id]
        style = spec.resolve_node_style(node)

        # ── Check for icon ──
        icon_png = None
        if node.icon:
            icon_png = get_icon_png_path(node.icon)

        if icon_png:
            # ── Icon node: icon image on top + label text box below ──
            # Layout: icon takes upper 60% of node area, label takes lower 40%
            icon_ratio = 0.60
            icon_h = pos.h * icon_ratio
            label_h = pos.h * (1.0 - icon_ratio)

            # Icon is square, centered horizontally within node width
            icon_size = min(icon_h, pos.w * 0.8)  # slightly smaller than full width
            icon_x = pos.x + (pos.w - icon_size) / 2
            icon_y = pos.y

            # Add icon image
            pic = slide.shapes.add_picture(
                str(icon_png),
                Inches(icon_x), Inches(icon_y),
                Inches(icon_size), Inches(icon_size)
            )

            # Add label text box below icon
            label_y = pos.y + icon_h
            label_box = slide.shapes.add_textbox(
                Inches(pos.x), Inches(label_y),
                Inches(pos.w), Inches(label_h)
            )
            tf = label_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = node.label
            label_font_size = _scaled_font_size(min(style.font_size, 9), layout_scale)
            _set_font(p, fonts.body, label_font_size,
                      f"#{theme.palette.navy}", bold=style.font_bold)
            p.alignment = PP_ALIGN.CENTER
            _enable_auto_shrink(tf)

            # Create an invisible anchor shape for connector attachment
            # This shape spans the full node area for correct connection routing
            anchor = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(pos.x), Inches(pos.y), Inches(pos.w), Inches(pos.h)
            )
            anchor.fill.background()
            anchor.line.fill.background()
            # Clear any default text
            anchor.text_frame.paragraphs[0].text = ""

            shape_map[node.id] = anchor
        else:
            # ── Standard shape node (no icon) ──
            mso_shape = SHAPE_MAP.get(node.shape, MSO_SHAPE.RECTANGLE)
            shape = slide.shapes.add_shape(
                mso_shape,
                Inches(pos.x), Inches(pos.y), Inches(pos.w), Inches(pos.h)
            )

            # Fill
            fill_color = style.fill or f"#{theme.palette.navy}"
            shape.fill.solid()
            shape.fill.fore_color.rgb = _hex_to_rgb(fill_color)

            # Border
            if style.border:
                shape.line.color.rgb = _hex_to_rgb(style.border)
                shape.line.width = Pt(style.border_width)
                if style.border_dash:
                    shape.line.dash_style = 2
            else:
                shape.line.fill.background()

            # Text — apply themed fonts consistently
            tf = shape.text_frame
            tf.word_wrap = True

            if node.sublabel:
                # Two-line: sublabel (small, body font) + label (large, heading font)
                p = tf.paragraphs[0]
                p.text = node.sublabel
                sub_fs = _scaled_font_size(max(style.font_size - 3, 7), layout_scale)
                _set_font(p, fonts.body, sub_fs,
                          style.font_color, bold=False)
                p.alignment = PP_ALIGN.CENTER

                p2 = tf.add_paragraph()
                p2.text = node.label
                main_fs = _scaled_font_size(style.font_size, layout_scale)
                _set_font(p2, fonts.heading, main_fs,
                          style.font_color, bold=style.font_bold)
                p2.alignment = PP_ALIGN.CENTER
            else:
                p = tf.paragraphs[0]
                p.text = node.label
                font_name = fonts.body
                if node.shape in ("rounded_rect", "circle", "oval"):
                    font_name = fonts.heading if style.font_bold else fonts.body
                node_fs = _scaled_font_size(style.font_size, layout_scale)
                _set_font(p, font_name, node_fs,
                          style.font_color, bold=style.font_bold)
                p.alignment = PP_ALIGN.CENTER

            _enable_auto_shrink(tf)
            shape_map[node.id] = shape

    # ── Draw group zones (nest-aware) ──
    # Compute bounding box for each group, bottom-up (leaves first).
    # A parent group's bbox encloses all its child group bboxes + direct member nodes.

    # Padding per nesting depth — scales with layout to maintain proportionality
    # and avoid bbox overlaps when nodes are small. Floor ensures readability.
    _BASE_DEPTH_PAD = {0: 0.25, 1: 0.18, 2: 0.12}
    _BASE_LABEL_H = 0.25
    _pad_scale = max(layout_scale, 0.4)  # floor at 0.4 to keep labels readable
    DEPTH_PAD = {k: max(0.06, v * _pad_scale) for k, v in _BASE_DEPTH_PAD.items()}
    LABEL_HEIGHT = max(0.10, _BASE_LABEL_H * _pad_scale)

    # Sort groups by depth (deepest first) for bottom-up bbox calculation
    group_map_local = {g.id: g for g in spec.groups}
    group_bboxes: dict[str, tuple[float, float, float, float]] = {}  # id → (min_x, min_y, max_x, max_y)

    def _groups_sorted_by_depth_desc() -> list[Group]:
        """Return groups sorted deepest-first."""
        return sorted(spec.groups, key=lambda g: spec.group_depth(g.id), reverse=True)

    for grp in _groups_sorted_by_depth_desc():
        depth = spec.group_depth(grp.id)
        pad = DEPTH_PAD.get(depth, 0.10)

        # Collect bounding points from:
        # 1. Direct member nodes
        all_node_ids = [n.id for n in spec.nodes if n.group == grp.id]
        points_x_min, points_y_min = [], []
        points_x_max, points_y_max = [], []

        for nid in all_node_ids:
            if nid in pos_map:
                p = pos_map[nid]
                points_x_min.append(p.x)
                points_y_min.append(p.y)
                points_x_max.append(p.x + p.w)
                points_y_max.append(p.y + p.h)

        # 2. Child group bounding boxes
        for child_gid in spec.group_children(grp.id):
            if child_gid in group_bboxes:
                cb = group_bboxes[child_gid]
                points_x_min.append(cb[0])
                points_y_min.append(cb[1])
                points_x_max.append(cb[2])
                points_y_max.append(cb[3])

        if not points_x_min:
            continue  # empty group

        min_x = min(points_x_min) - pad
        min_y = min(points_y_min) - pad - LABEL_HEIGHT
        max_x = max(points_x_max) + pad
        max_y = max(points_y_max) + pad

        group_bboxes[grp.id] = (min_x, min_y, max_x, max_y)

    # Draw groups from outermost (depth 0) to innermost (depth N)
    # so that outer borders appear behind inner ones
    groups_draw_order = sorted(spec.groups,
                               key=lambda g: spec.group_depth(g.id))

    for grp in groups_draw_order:
        if grp.id not in group_bboxes:
            continue

        min_x, min_y, max_x, max_y = group_bboxes[grp.id]
        depth = spec.group_depth(grp.id)
        gs = grp.effective_style()

        zone = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(min_x), Inches(min_y),
            Inches(max_x - min_x), Inches(max_y - min_y)
        )

        # Light fill for outer groups, transparent for inner
        if gs.fill:
            zone.fill.solid()
            zone.fill.fore_color.rgb = _hex_to_rgb(gs.fill)
        else:
            zone.fill.background()

        zone.line.color.rgb = _hex_to_rgb(gs.border)
        zone.line.width = Pt(ds.group_border_width)
        if gs.border_dash:
            zone.line.dash_style = 2

        # Group label at top-left — width matches group bbox, font scales
        label_font_size = _scaled_font_size(
            max(ds.group_label_font_size - depth, 7), layout_scale, min_size=5.0)
        group_bbox_w = max_x - min_x
        label_w = min(group_bbox_w - 0.1, 3.0)  # cap at 3", leave margin
        if label_w < 0.3:
            label_w = 0.3
        lbl = slide.shapes.add_textbox(
            Inches(min_x + 0.08), Inches(min_y + 0.04),
            Inches(label_w), Inches(LABEL_HEIGHT)
        )
        lp = lbl.text_frame.paragraphs[0]
        lp.text = grp.label
        _set_font(lp, fonts.body, label_font_size,
                  gs.border, bold=True)
        _enable_auto_shrink(lbl.text_frame)

    # ── Draw connectors ──
    # Routing strategy is determined by the structural relationship between
    # source and target nodes' groups (Step 13 improvements):
    #   - Same group / flow-adjacent groups → direct or L-route
    #   - Laterally-separated groups → Manhattan routing (B)
    #   - Fan-in/out (3+ edges, same group) → shared bus lines (C)
    #   - Back-edges (cycles) → U-turn routing
    is_flowchart = (spec.type == "flowchart")
    direction = spec.direction

    # Build layer map for direction-aware cp detection
    layer_map = _assign_layers(spec)

    # Detect back-edges (cycles) — need special U-turn routing
    fwd_adj, _rev_adj = _build_adjacency(spec)
    back_edges = _find_back_edges(fwd_adj, [n.id for n in spec.nodes])

    # Build node-to-group membership lookup
    node_to_group = {n.id: n.group for n in spec.nodes}

    # Compute diagram right margin for back-edge routing
    all_right_edges = [p.x + p.w for p in positions]
    diagram_right_x = max(all_right_edges) if all_right_edges else 8.0

    # Detect merge targets (nodes with >1 incoming edge) for L-routing
    in_degree = defaultdict(int)
    for e in spec.edges:
        in_degree[e.to_id] += 1
    merge_targets = {nid for nid, deg in in_degree.items() if deg > 1}

    is_vertical = direction in ("TB", "BT")

    # ── Bus line drawing: explicit bus_group + auto-detection ──
    #
    # Two mechanisms for bus consolidation:
    #  1. Explicit bus_group: user-specified grouping key per edge
    #  2. Auto-detection: merge fan-in/fan-out patterns when ALL conditions met:
    #     (a) Fan-in (N sources → 1 target) or fan-out (1 source → N targets)
    #     (b) N ≥ 2 edges
    #     (c) No edge in the group has a label (labels = explicit distinct semantics)
    #     (d) All edges share the same visual style (color, dash, width)
    #     (e) Not back-edges
    #     When any edge has a label or different style, the entire group is kept
    #     individual — the label signals "this connection is special".

    bus_handled: set[tuple[str, str]] = set()

    # ── 1. Explicit bus_group ──
    bus_groups: dict[str, list] = defaultdict(list)
    for edge in spec.edges:
        if edge.bus_group and (edge.from_id, edge.to_id) not in back_edges:
            bus_groups[edge.bus_group].append(edge)

    for bg_key, bg_edges in bus_groups.items():
        if len(bg_edges) < 2:
            continue
        targets = {e.to_id for e in bg_edges}
        sources = {e.from_id for e in bg_edges}
        if len(targets) == 1:
            target_id = next(iter(targets))
            handled = _draw_fan_in_bus(
                slide, target_id, bg_edges, shape_map, pos_map,
                direction, is_flowchart, theme, ds, layout_scale, fonts)
            bus_handled.update(handled)
        elif len(sources) == 1:
            source_id = next(iter(sources))
            handled = _draw_fan_out_bus(
                slide, source_id, bg_edges, shape_map, pos_map,
                direction, is_flowchart, theme, ds, layout_scale, fonts)
            bus_handled.update(handled)

    # ── 2. Auto-detection of mergeable fan-in / fan-out patterns ──
    def _edge_style_key(e):
        """Canonical style tuple for merge-compatibility comparison."""
        es = e.effective_style()
        return (es.color, es.dash, es.width)

    # Build candidate groups: fan-out (by source) and fan-in (by target)
    fanout_candidates: dict[str, list] = defaultdict(list)  # source_id → edges
    fanin_candidates: dict[str, list] = defaultdict(list)   # target_id → edges

    for edge in spec.edges:
        key = (edge.from_id, edge.to_id)
        if key in bus_handled or key in back_edges:
            continue
        if edge.bus_group:  # already handled above
            continue
        fanout_candidates[edge.from_id].append(edge)
        fanin_candidates[edge.to_id].append(edge)

    def _is_auto_mergeable(edges: list) -> bool:
        """Check if a set of edges meets all auto-merge conditions."""
        if len(edges) < 2:
            return False
        # Condition: no labels on ANY edge in the group
        if any(e.label for e in edges):
            return False
        # Condition: all edges same style
        styles = {_edge_style_key(e) for e in edges}
        if len(styles) > 1:
            return False
        return True

    # Fan-out: 1 source → N targets (N ≥ 2)
    for src_id, edges in fanout_candidates.items():
        if len(edges) < 2:
            continue
        if not _is_auto_mergeable(edges):
            continue
        # Check none already handled
        if any((e.from_id, e.to_id) in bus_handled for e in edges):
            continue
        handled = _draw_fan_out_bus(
            slide, src_id, edges, shape_map, pos_map,
            direction, is_flowchart, theme, ds, layout_scale, fonts)
        bus_handled.update(handled)

    # Fan-in: N sources → 1 target (N ≥ 2)
    for tgt_id, edges in fanin_candidates.items():
        if len(edges) < 2:
            continue
        # Filter out edges already drawn as fan-out bus
        remaining = [e for e in edges
                     if (e.from_id, e.to_id) not in bus_handled]
        if not _is_auto_mergeable(remaining):
            continue
        handled = _draw_fan_in_bus(
            slide, tgt_id, remaining, shape_map, pos_map,
            direction, is_flowchart, theme, ds, layout_scale, fonts)
        bus_handled.update(handled)

    # ── Compute port offsets to reduce line overlaps ──
    port_offsets = _compute_port_offsets(
        spec, shape_map, layer_map, direction, back_edges, bus_handled)

    # ── Per-edge routing — Two-Phase Approach ──
    # Phase 1: Compute paths for L-route and Manhattan without drawing.
    #          Back-edges and direct connectors are drawn immediately.
    # Phase 2: Nudge overlapping horizontal segments across all planned paths.
    # Phase 3: Draw all planned paths and labels.

    planned_paths: list[dict] = []   # deferred L-route / Manhattan paths

    for edge in spec.edges:
        # Skip edges already drawn as bus lines
        if (edge.from_id, edge.to_id) in bus_handled:
            continue

        sa = shape_map.get(edge.from_id)
        sb = shape_map.get(edge.to_id)
        if sa is None or sb is None:
            continue

        es = edge.effective_style()
        color = _hex_to_rgb(es.color)
        width = Pt(es.width)
        arrow = es.arrow if is_flowchart else False
        dash = 2 if es.dash else 0   # Apply dash from edge style

        from_pos = pos_map.get(edge.from_id)
        to_pos = pos_map.get(edge.to_id)
        from_layer = layer_map.get(edge.from_id, 0)
        to_layer = layer_map.get(edge.to_id, 0)
        is_inter_layer = (from_layer != to_layer)

        # Classify routing strategy based on group relationships
        route_type = _classify_edge_route(
            edge.from_id, edge.to_id, node_to_group, group_bboxes,
            pos_map, layer_map, direction, back_edges, merge_targets)

        # Direction-aware connection point detection
        begin_cp, end_cp = _detect_cp(
            sa, sb, direction=direction,
            layer_a=from_layer, layer_b=to_layer
        )

        # Never use snap for inter-layer — prevents python-pptx overrides
        use_snap = False if is_inter_layer else (not _is_diamond(sa))

        # Retrieve port offsets for this edge
        src_port_off, tgt_port_off = port_offsets.get(
            (edge.from_id, edge.to_id), (0.0, 0.0))

        # ── Back-edge (cycle) — draw immediately (U-turn, rarely overlaps) ──
        if route_type == 'back_edge' and from_pos and to_pos:
            margin_x = Inches(diagram_right_x + 0.5)
            dash_style = 3   # dash-dot for back-edges

            if direction in ("TB",):
                src_rx = sa.left + sa.width
                src_cy = sa.top + sa.height // 2
                tgt_rx = sb.left + sb.width
                tgt_cy = sb.top + sb.height // 2
                _draw_line_xy(slide, src_rx, src_cy, margin_x, src_cy,
                              color=color, width=width, theme=theme, dash=dash_style)
                _draw_line_xy(slide, margin_x, src_cy, margin_x, tgt_cy,
                              color=color, width=width, theme=theme, dash=dash_style)
                _draw_line_xy(slide, margin_x, tgt_cy, tgt_rx, tgt_cy,
                              color=color, width=width, arrow=True, theme=theme,
                              dash=dash_style)
            elif direction in ("BT",):
                src_rx = sa.left + sa.width
                src_cy = sa.top + sa.height // 2
                tgt_rx = sb.left + sb.width
                tgt_cy = sb.top + sb.height // 2
                _draw_line_xy(slide, src_rx, src_cy, margin_x, src_cy,
                              color=color, width=width, theme=theme, dash=dash_style)
                _draw_line_xy(slide, margin_x, src_cy, margin_x, tgt_cy,
                              color=color, width=width, theme=theme, dash=dash_style)
                _draw_line_xy(slide, margin_x, tgt_cy, tgt_rx, tgt_cy,
                              color=color, width=width, arrow=True, theme=theme,
                              dash=dash_style)
            else:
                all_bottoms = [p.y + p.h for p in positions]
                margin_y = Inches(max(all_bottoms) + 0.5)
                src_cx = sa.left + sa.width // 2
                src_by = sa.top + sa.height
                tgt_cx = sb.left + sb.width // 2
                tgt_by = sb.top + sb.height
                _draw_line_xy(slide, src_cx, src_by, src_cx, margin_y,
                              color=color, width=width, theme=theme, dash=dash_style)
                _draw_line_xy(slide, src_cx, margin_y, tgt_cx, margin_y,
                              color=color, width=width, theme=theme, dash=dash_style)
                _draw_line_xy(slide, tgt_cx, margin_y, tgt_cx, tgt_by,
                              color=color, width=width, arrow=True, theme=theme,
                              dash=dash_style)

            if edge.label:
                lbl_x = margin_x + Inches(0.05)
                lbl_y = (sa.top + sa.height // 2 + sb.top + sb.height // 2) // 2
                lbl = slide.shapes.add_textbox(lbl_x, lbl_y - Inches(0.15),
                                               Inches(0.6), Inches(0.3))
                lp = lbl.text_frame.paragraphs[0]
                lp.text = edge.label
                edge_fs = _scaled_font_size(ds.edge_label_font_size, layout_scale)
                _set_font(lp, fonts.body, edge_fs, es.color, bold=True)
            continue

        # ── Cross-group Manhattan routing — plan path (don't draw yet) ──
        elif route_type == 'cross_group' and from_pos and to_pos:
            fg = node_to_group.get(edge.from_id)
            tg = node_to_group.get(edge.to_id)
            if fg in group_bboxes and tg in group_bboxes:
                path_pts = _plan_manhattan_route(
                    sa, sb, from_pos, to_pos,
                    group_bboxes[fg], group_bboxes[tg],
                    direction, src_port_off, tgt_port_off)
                planned_paths.append({
                    'edge': edge, 'path_points': path_pts,
                    'color': color, 'width': width,
                    'arrow': arrow, 'dash': dash,
                })
            else:
                # Fallback to direct if group bbox missing — draw immediately
                c = _draw_connector(slide, sa, sb,
                                    begin_cp=begin_cp, end_cp=end_cp,
                                    color=color, width=width,
                                    arrow=arrow, snap=use_snap, theme=theme)
                if dash:
                    c.line.dash_style = dash
                edge_path_points = [
                    _cp_coords(sa, begin_cp, src_port_off),
                    _cp_coords(sb, end_cp, tgt_port_off)]
                if edge.label and edge_path_points:
                    all_node_shapes = [shape_map[nid] for nid in shape_map]
                    lbl_x, lbl_y, lbl_w, lbl_h = _place_edge_label(
                        edge_path_points, edge.label, all_node_shapes,
                        slide_w=prs.slide_width, slide_h=prs.slide_height)
                    lbl = slide.shapes.add_textbox(lbl_x, lbl_y, lbl_w, lbl_h)
                    lp = lbl.text_frame.paragraphs[0]
                    lp.text = edge.label
                    edge_fs = _scaled_font_size(ds.edge_label_font_size, layout_scale)
                    _set_font(lp, fonts.body, edge_fs, es.color, bold=True)
            continue

        # ── L-shaped routing — plan path (don't draw yet) ──
        elif route_type == 'l_route' and from_pos and to_pos:
            _BEND_RATIO = 0.80
            path_pts = []

            if direction in ("TB",):
                src_bx, src_by = _cp_coords(sa, 2, src_port_off)
                tgt_tx, tgt_ty = _cp_coords(sb, 0, tgt_port_off)
                gap = tgt_ty - src_by
                mid_y = src_by + int(gap * _BEND_RATIO)
                path_pts = [
                    (src_bx, src_by), (src_bx, mid_y),
                    (tgt_tx, mid_y), (tgt_tx, tgt_ty)]
            elif direction in ("BT",):
                src_tx, src_ty = _cp_coords(sa, 0, src_port_off)
                tgt_bx, tgt_by = _cp_coords(sb, 2, tgt_port_off)
                gap = src_ty - tgt_by
                mid_y = src_ty - int(gap * _BEND_RATIO)
                path_pts = [
                    (src_tx, src_ty), (src_tx, mid_y),
                    (tgt_bx, mid_y), (tgt_bx, tgt_by)]
            else:
                src_rx, src_ry = _cp_coords(sa, 1, src_port_off)
                tgt_lx, tgt_ly = _cp_coords(sb, 3, tgt_port_off)
                gap = tgt_lx - src_rx
                mid_x = src_rx + int(gap * _BEND_RATIO)
                path_pts = [
                    (src_rx, src_ry), (mid_x, src_ry),
                    (mid_x, tgt_ly), (tgt_lx, tgt_ly)]

            planned_paths.append({
                'edge': edge, 'path_points': path_pts,
                'color': color, 'width': width,
                'arrow': arrow, 'dash': dash,
            })
            continue

        else:
            # Standard direct connector — draw immediately
            c = _draw_connector(slide, sa, sb,
                                begin_cp=begin_cp, end_cp=end_cp,
                                color=color, width=width,
                                arrow=arrow, snap=use_snap, theme=theme)
            if dash:
                c.line.dash_style = dash
            edge_path_points = [
                _cp_coords(sa, begin_cp, src_port_off),
                _cp_coords(sb, end_cp, tgt_port_off)]

            if edge.label and edge_path_points:
                all_node_shapes = [shape_map[nid] for nid in shape_map]
                lbl_x, lbl_y, lbl_w, lbl_h = _place_edge_label(
                    edge_path_points, edge.label, all_node_shapes,
                    slide_w=prs.slide_width, slide_h=prs.slide_height)
                lbl = slide.shapes.add_textbox(lbl_x, lbl_y, lbl_w, lbl_h)
                lp = lbl.text_frame.paragraphs[0]
                lp.text = edge.label
                edge_fs = _scaled_font_size(ds.edge_label_font_size, layout_scale)
                _set_font(lp, fonts.body, edge_fs, es.color, bold=True)

    # ── Phase 2: Nudge overlapping horizontal segments ──
    _nudge_overlapping_segments(planned_paths)

    # ── Phase 3: Draw all planned paths + labels ──
    all_node_shapes = [shape_map[nid] for nid in shape_map]
    for plan in planned_paths:
        _draw_planned_path(slide, plan, theme)

        edge = plan['edge']
        if edge.label and plan['path_points']:
            es = edge.effective_style()
            lbl_x, lbl_y, lbl_w, lbl_h = _place_edge_label(
                plan['path_points'], edge.label, all_node_shapes,
                slide_w=prs.slide_width, slide_h=prs.slide_height)
            lbl = slide.shapes.add_textbox(lbl_x, lbl_y, lbl_w, lbl_h)
            lp = lbl.text_frame.paragraphs[0]
            lp.text = edge.label
            edge_fs = _scaled_font_size(ds.edge_label_font_size, layout_scale)
            _set_font(lp, fonts.body, edge_fs, es.color, bold=True)

    return prs


# ══════════════════════════════════════════════
# CLI entry point
# ══════════════════════════════════════════════

def render_from_json(json_str: str, output_path: str,
                     prs: Optional[Presentation] = None,
                     template_path: Optional[str] = None,
                     theme: ThemeConfig = DEFAULT_THEME) -> str:
    """Parse JSON, render to PPTX, save. Returns output path."""
    spec = parse_diagram_json(json_str)
    prs = render_diagram(spec, prs, template_path=template_path, theme=theme)
    prs.save(output_path)
    return output_path


if __name__ == "__main__":
    import json

    theme = DEFAULT_THEME
    prs = None  # accumulate all test slides in one file

    # ── Test 1: Flowchart ──
    flow_json = json.dumps({
        "type": "flowchart",
        "direction": "TB",
        "title": "API認証フロー",
        "classDefs": {
            "process": {"fill": "#1E2761", "border": "#3B82F6", "font_color": "#FFFFFF"},
            "decision": {"fill": "#F59E0B", "font_color": "#1E293B", "font_size": 10},
            "terminal": {"fill": "#3B82F6"},
            "error": {"fill": "#2D3A6E", "font_color": "#FFFFFF"},
        },
        "nodes": [
            {"id": "start", "label": "開始", "shape": "rounded_rect", "class": "terminal"},
            {"id": "proc1", "label": "リクエスト受付", "shape": "rect", "class": "process"},
            {"id": "auth", "label": "認証OK？", "shape": "diamond", "class": "decision"},
            {"id": "ok", "label": "データ処理", "shape": "rect", "class": "process"},
            {"id": "ng", "label": "エラー返却", "shape": "rect", "class": "error"},
            {"id": "resp", "label": "レスポンス生成", "shape": "rect", "class": "process"},
            {"id": "end", "label": "終了", "shape": "rounded_rect", "class": "terminal"},
        ],
        "edges": [
            {"from": "start", "to": "proc1"},
            {"from": "proc1", "to": "auth"},
            {"from": "auth", "to": "ok", "label": "Yes"},
            {"from": "auth", "to": "ng", "label": "No"},
            {"from": "ok", "to": "resp"},
            {"from": "ng", "to": "end"},
            {"from": "resp", "to": "end"},
        ],
        "groups": [],
        "layout": {"node_width": 2.4, "node_height": 0.7, "h_gap": 1.0, "v_gap": 0.6},
    })
    spec = parse_diagram_json(flow_json)
    prs = render_diagram(spec, prs, theme=theme)
    print("  ✓ Flowchart rendered")

    # ── Test 2: Network Diagram ──
    net_json = json.dumps({
        "type": "network",
        "direction": "TB",
        "title": "ネットワーク構成図",
        "classDefs": {
            "external": {"fill": "#94A3B8", "border": "#1E293B", "font_color": "#FFFFFF"},
            "firewall": {"fill": "#F59E0B", "font_color": "#1E293B", "font_size": 9},
            "core": {"fill": "#1E2761", "border": "#3B82F6"},
            "switch": {"fill": "#2D3A6E", "border": "#3B82F6", "font_size": 10},
            "web": {"fill": "#3B82F6", "font_size": 9},
            "app": {"fill": "#06B6D4", "font_size": 9},
            "db": {"fill": "#1E2761", "font_size": 9},
        },
        "nodes": [
            {"id": "inet", "label": "Internet", "shape": "rounded_rect", "class": "external"},
            {"id": "fw", "label": "Firewall", "shape": "diamond", "class": "firewall"},
            {"id": "router", "label": "Core Router", "shape": "rounded_rect", "class": "core"},
            {"id": "sw1", "label": "Switch-A", "shape": "rounded_rect", "class": "switch", "group": "access"},
            {"id": "sw2", "label": "Switch-B", "shape": "rounded_rect", "class": "switch", "group": "access"},
            {"id": "sw3", "label": "Switch-C", "shape": "rounded_rect", "class": "switch", "group": "access"},
            {"id": "web1", "label": "Web-01", "shape": "rect", "class": "web", "group": "servers"},
            {"id": "web2", "label": "Web-02", "shape": "rect", "class": "web", "group": "servers"},
            {"id": "app1", "label": "App-01", "shape": "rect", "class": "app", "group": "servers"},
            {"id": "app2", "label": "App-02", "shape": "rect", "class": "app", "group": "servers"},
            {"id": "db1", "label": "DB-01", "shape": "rect", "class": "db", "group": "servers"},
            {"id": "db2", "label": "DB-02", "shape": "rect", "class": "db", "group": "servers"},
        ],
        "edges": [
            {"from": "inet", "to": "fw"},
            {"from": "fw", "to": "router"},
            {"from": "router", "to": "sw1"},
            {"from": "router", "to": "sw2"},
            {"from": "router", "to": "sw3"},
            {"from": "sw1", "to": "web1"},
            {"from": "sw1", "to": "web2"},
            {"from": "sw2", "to": "app1"},
            {"from": "sw2", "to": "app2"},
            {"from": "sw3", "to": "db1"},
            {"from": "sw3", "to": "db2"},
        ],
        "groups": [
            {"id": "access", "label": "Access Layer", "style": {"border": "#06B6D4", "border_dash": True}},
            {"id": "servers", "label": "Server Farm", "style": {"border": "#1E2761", "border_dash": True}},
        ],
        "layout": {"node_width": 1.6, "node_height": 0.6, "h_gap": 0.3, "v_gap": 0.5},
    })
    spec = parse_diagram_json(net_json)
    prs = render_diagram(spec, prs, theme=theme)
    print("  ✓ Network diagram rendered")

    # ── Test 3: Organization Chart ──
    org_json = json.dumps({
        "type": "orgchart",
        "direction": "TB",
        "title": "組織図",
        "classDefs": {
            "ceo": {"fill": "#141B41", "border": "#3B82F6", "font_size": 12},
            "vp_eng": {"fill": "#1E2761", "border": "#3B82F6"},
            "vp_sales": {"fill": "#1E2761", "border": "#06B6D4"},
            "vp_ops": {"fill": "#1E2761", "border": "#F59E0B"},
            "team_eng": {"fill": "#2D3A6E", "border": "#3B82F6", "font_size": 10},
            "team_sales": {"fill": "#2D3A6E", "border": "#06B6D4", "font_size": 10},
            "team_ops": {"fill": "#2D3A6E", "border": "#F59E0B", "font_size": 10},
        },
        "nodes": [
            {"id": "ceo", "label": "田中 太郎", "sublabel": "CEO", "shape": "rounded_rect", "class": "ceo"},
            {"id": "vp1", "label": "鈴木 花子", "sublabel": "VP of Engineering", "shape": "rounded_rect", "class": "vp_eng"},
            {"id": "vp2", "label": "佐藤 次郎", "sublabel": "VP of Sales", "shape": "rounded_rect", "class": "vp_sales"},
            {"id": "vp3", "label": "山田 美咲", "sublabel": "VP of Operations", "shape": "rounded_rect", "class": "vp_ops"},
            {"id": "t1", "label": "高橋 一郎", "sublabel": "Backend", "shape": "rounded_rect", "class": "team_eng"},
            {"id": "t2", "label": "伊藤 真理", "sublabel": "Frontend", "shape": "rounded_rect", "class": "team_eng"},
            {"id": "t3", "label": "渡辺 健太", "sublabel": "Japan Sales", "shape": "rounded_rect", "class": "team_sales"},
            {"id": "t4", "label": "小林 あい", "sublabel": "Global Sales", "shape": "rounded_rect", "class": "team_sales"},
            {"id": "t5", "label": "加藤 翔", "sublabel": "DevOps", "shape": "rounded_rect", "class": "team_ops"},
            {"id": "t6", "label": "松本 由美", "sublabel": "QA", "shape": "rounded_rect", "class": "team_ops"},
        ],
        "edges": [
            {"from": "ceo", "to": "vp1", "style": {"color": "#3B82F6"}},
            {"from": "ceo", "to": "vp2", "style": {"color": "#06B6D4"}},
            {"from": "ceo", "to": "vp3", "style": {"color": "#F59E0B"}},
            {"from": "vp1", "to": "t1", "style": {"color": "#3B82F6"}},
            {"from": "vp1", "to": "t2", "style": {"color": "#3B82F6"}},
            {"from": "vp2", "to": "t3", "style": {"color": "#06B6D4"}},
            {"from": "vp2", "to": "t4", "style": {"color": "#06B6D4"}},
            {"from": "vp3", "to": "t5", "style": {"color": "#F59E0B"}},
            {"from": "vp3", "to": "t6", "style": {"color": "#F59E0B"}},
        ],
        "groups": [],
        "layout": {"node_width": 1.8, "node_height": 0.8, "h_gap": 0.3, "v_gap": 0.8},
    })
    spec = parse_diagram_json(org_json)
    prs = render_diagram(spec, prs, theme=theme)
    print("  ✓ Org chart rendered")

    out = '/sessions/stoic-zealous-fermi/mnt/presentations/DiagramRenderer_Test.pptx'
    prs.save(out)
    print(f"\n✅ All 3 tests saved: {out}")
