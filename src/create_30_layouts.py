#!/usr/bin/env python3
"""
Midnight Executive Template — 30 Layout Patterns
=================================================
Data/Analysis Report focused template with rich decorative elements.

Design principles:
  - ALL decorative elements baked into layout XML
  - Placeholders are TEXT ONLY — users just type
  - 30 varied patterns prevent "template" look
  - Geometric shapes, color blocks, panels for visual richness
  - No accent lines under titles (per design guidelines)
  - Visual motif: thick left accent bars + geometric panels

Technical:
  - First 11 layouts overwrite default slots
  - Layouts 12-30 added via OPC SlideLayoutPart
"""

from pptx import Presentation
from pptx.parts.slide import SlideLayoutPart
from pptx.opc.package import PackURI
from pptx.oxml import parse_xml
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from lxml import etree

# ════════════════════════════════════════════
# CONSTANTS
# ════════════════════════════════════════════
SLIDE_W = 12192000   # 13.333" in EMU
SLIDE_H = 6858000    # 7.5" in EMU

def emu(inches):
    return int(inches * 914400)

# Color palette — Midnight Executive
C = {
    'navy':        '1E2761',
    'dark_navy':   '141B41',
    'ice_blue':    'CADCFC',
    'white':       'FFFFFF',
    'light_gray':  'F5F7FA',
    'panel_gray':  'EDF0F7',
    'mid_gray':    '94A3B8',
    'dark_text':   '1E293B',
    'accent':      '3B82F6',
    'accent_dark': '2563EB',
    'teal':        '06B6D4',
    'amber':       'F59E0B',
    'soft_navy':   '2D3A6E',
    'card_bg':     'F0F4FF',
}

F_HEAD = 'Georgia'
F_BODY = 'Calibri'

# XML namespaces
NS = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
}

def _tag(ns_local):
    if ':' in ns_local:
        prefix, local = ns_local.split(':', 1)
        return f'{{{NS[prefix]}}}{local}'
    return ns_local


# ════════════════════════════════════════════
# XML ELEMENT BUILDERS
# ════════════════════════════════════════════

def El(tag, attrib=None, children=None, text=None):
    elem = etree.Element(_tag(tag), {k: v for k, v in (attrib or {}).items()})
    if text is not None:
        elem.text = text
    for child in (children or []):
        if child is not None:
            elem.append(child)
    return elem

def Sub(parent, tag, attrib=None, text=None):
    elem = etree.SubElement(parent, _tag(tag), {k: v for k, v in (attrib or {}).items()})
    if text is not None:
        elem.text = text
    return elem


# ─── Decorative shapes ───

def xml_rect(sid, name, x, y, w, h, color):
    sp = El('p:sp')
    nv = Sub(sp, 'p:nvSpPr')
    Sub(nv, 'p:cNvPr', {'id': str(sid), 'name': name})
    cnv = Sub(nv, 'p:cNvSpPr')
    Sub(cnv, 'a:spLocks', {'noGrp': '1', 'noSelect': '0', 'noRot': '1', 'noMove': '1'})
    Sub(nv, 'p:nvPr')
    spPr = Sub(sp, 'p:spPr')
    xfrm = Sub(spPr, 'a:xfrm')
    Sub(xfrm, 'a:off', {'x': str(x), 'y': str(y)})
    Sub(xfrm, 'a:ext', {'cx': str(w), 'cy': str(h)})
    Sub(spPr, 'a:prstGeom', {'prst': 'rect'})
    sf = Sub(spPr, 'a:solidFill')
    Sub(sf, 'a:srgbClr', {'val': color})
    ln = Sub(spPr, 'a:ln')
    Sub(ln, 'a:noFill')
    return sp

def xml_roundrect(sid, name, x, y, w, h, color, radius='16667'):
    sp = El('p:sp')
    nv = Sub(sp, 'p:nvSpPr')
    Sub(nv, 'p:cNvPr', {'id': str(sid), 'name': name})
    cnv = Sub(nv, 'p:cNvSpPr')
    Sub(cnv, 'a:spLocks', {'noGrp': '1', 'noSelect': '0', 'noRot': '1', 'noMove': '1'})
    Sub(nv, 'p:nvPr')
    spPr = Sub(sp, 'p:spPr')
    xfrm = Sub(spPr, 'a:xfrm')
    Sub(xfrm, 'a:off', {'x': str(x), 'y': str(y)})
    Sub(xfrm, 'a:ext', {'cx': str(w), 'cy': str(h)})
    geom = Sub(spPr, 'a:prstGeom', {'prst': 'roundRect'})
    avLst = Sub(geom, 'a:avLst')
    Sub(avLst, 'a:gd', {'name': 'adj', 'fmla': f'val {radius}'})
    sf = Sub(spPr, 'a:solidFill')
    Sub(sf, 'a:srgbClr', {'val': color})
    ln = Sub(spPr, 'a:ln')
    Sub(ln, 'a:noFill')
    return sp

def xml_circle(sid, name, cx, cy, r, color):
    sp = El('p:sp')
    nv = Sub(sp, 'p:nvSpPr')
    Sub(nv, 'p:cNvPr', {'id': str(sid), 'name': name})
    cnv = Sub(nv, 'p:cNvSpPr')
    Sub(cnv, 'a:spLocks', {'noGrp': '1'})
    Sub(nv, 'p:nvPr')
    spPr = Sub(sp, 'p:spPr')
    xfrm = Sub(spPr, 'a:xfrm')
    Sub(xfrm, 'a:off', {'x': str(cx - r), 'y': str(cy - r)})
    Sub(xfrm, 'a:ext', {'cx': str(2 * r), 'cy': str(2 * r)})
    Sub(spPr, 'a:prstGeom', {'prst': 'ellipse'})
    sf = Sub(spPr, 'a:solidFill')
    Sub(sf, 'a:srgbClr', {'val': color})
    ln = Sub(spPr, 'a:ln')
    Sub(ln, 'a:noFill')
    return sp


# ─── Placeholder builder ───

def xml_ph(sid, name, ph_type, ph_idx, x, y, w, h,
           font_size=14, font_name=F_BODY, font_color='1E293B',
           bold=False, align='l', prompt='テキストを入力'):
    """Editable placeholder with bullet suppression."""
    sp = El('p:sp')
    nv = Sub(sp, 'p:nvSpPr')
    Sub(nv, 'p:cNvPr', {'id': str(sid), 'name': name})
    cnv = Sub(nv, 'p:cNvSpPr')
    Sub(cnv, 'a:spLocks', {'noGrp': '1'})
    nvPr = Sub(nv, 'p:nvPr')
    ph_attr = {'idx': str(ph_idx)}
    if ph_type:
        ph_attr['type'] = ph_type
    Sub(nvPr, 'p:ph', ph_attr)

    spPr = Sub(sp, 'p:spPr')
    xfrm = Sub(spPr, 'a:xfrm')
    Sub(xfrm, 'a:off', {'x': str(x), 'y': str(y)})
    Sub(xfrm, 'a:ext', {'cx': str(w), 'cy': str(h)})

    txBody = Sub(sp, 'p:txBody')
    Sub(txBody, 'a:bodyPr', {'wrap': 'square', 'anchor': 't', 'anchorCtr': '0'})

    # lstStyle — defines DEFAULT text appearance for this placeholder.
    # When ph.text = "..." replaces content, these defRPr styles are inherited.
    sz = str(font_size * 100)
    lstStyle = Sub(txBody, 'a:lstStyle')

    def _add_defRPr(parent_pPr):
        """Add defRPr with font color/size/name/bold to a pPr element."""
        Sub(parent_pPr, 'a:buNone')
        drp_attr = {'sz': sz, 'dirty': '0'}
        if bold:
            drp_attr['b'] = '1'
        drp = Sub(parent_pPr, 'a:defRPr', drp_attr)
        sf = Sub(drp, 'a:solidFill')
        Sub(sf, 'a:srgbClr', {'val': font_color})
        Sub(drp, 'a:latin', {'typeface': font_name})
        Sub(drp, 'a:ea', {'typeface': font_name})

    defPPr = Sub(lstStyle, 'a:defPPr', {'algn': align})
    _add_defRPr(defPPr)
    for lvl in range(1, 10):
        lvlPPr = Sub(lstStyle, f'a:lvl{lvl}pPr', {'algn': align})
        _add_defRPr(lvlPPr)

    # Prompt text (ghost text shown in edit mode)
    p = Sub(txBody, 'a:p')
    pPr = Sub(p, 'a:pPr', {'algn': align})
    Sub(pPr, 'a:buNone')

    rPr_attr = {'lang': 'ja-JP', 'sz': sz, 'dirty': '0'}
    if bold:
        rPr_attr['b'] = '1'
    r = Sub(p, 'a:r')
    rPr = Sub(r, 'a:rPr', rPr_attr)
    sf = Sub(rPr, 'a:solidFill')
    Sub(sf, 'a:srgbClr', {'val': font_color})
    Sub(rPr, 'a:latin', {'typeface': font_name})
    Sub(rPr, 'a:ea', {'typeface': font_name})
    Sub(r, 'a:t').text = prompt
    return sp


# ═══════════════════════════════════════════
# LAYOUT ASSEMBLY
# ═══════════════════════════════════════════

def _build_layout(name, type_str, elements):
    sld = El('p:sldLayout', {'type': type_str, 'preserve': '1'})
    cSld = Sub(sld, 'p:cSld', {'name': name})
    spTree = Sub(cSld, 'p:spTree')
    nv = Sub(spTree, 'p:nvGrpSpPr')
    Sub(nv, 'p:cNvPr', {'id': '1', 'name': ''})
    Sub(nv, 'p:cNvGrpSpPr')
    Sub(nv, 'p:nvPr')
    grp = Sub(spTree, 'p:grpSpPr')
    xfrm = Sub(grp, 'a:xfrm')
    Sub(xfrm, 'a:off', {'x': '0', 'y': '0'})
    Sub(xfrm, 'a:ext', {'cx': '0', 'cy': '0'})
    Sub(xfrm, 'a:chOff', {'x': '0', 'y': '0'})
    Sub(xfrm, 'a:chExt', {'cx': '0', 'cy': '0'})
    for el in elements:
        spTree.append(el)
    Sub(sld, 'p:clrMapOvr')
    return sld


# ─── Shorthand: common header bar + title + subtitle ───
def _header_bar(sid_base=2):
    """Returns [header_rect, title_ph, subtitle_ph] for standard content slides."""
    return [
        xml_rect(sid_base, 'HeaderBar', 0, 0, SLIDE_W, emu(1.15), C['navy']),
        xml_ph(sid_base+8, 'SlideTitle.Header', 'body', 15,
               emu(0.8), emu(0.1), emu(10), emu(0.52),
               28, F_HEAD, C['white'], True, 'l', 'スライドタイトル'),
        xml_ph(sid_base+9, 'SlideSubtitle.Header', 'body', 16,
               emu(0.8), emu(0.72), emu(10), emu(0.3),
               12, F_BODY, C['ice_blue'], False, 'l', 'サブタイトル'),
    ]

def _slide_num(sid):
    return xml_ph(sid, 'SlideNum.Footer', 'sldNum', 50,
                  emu(12.0), emu(7.05), emu(1.2), emu(0.3),
                  10, F_BODY, C['mid_gray'], False, 'r', '<#>')


# ════════════════════════════════════════════
# 30 LAYOUT DEFINITIONS
# ════════════════════════════════════════════

# ──────────────── GROUP 1: TITLE SLIDES (3) ────────────────

def L00_title_standard():
    """Full navy bg, accent bar, large title."""
    return _build_layout('Title.1Title.Single', 'title', [
        xml_rect(2, 'BG', 0, 0, SLIDE_W, SLIDE_H, C['navy']),
        xml_rect(3, 'AccentBar', 0, 0, emu(0.18), SLIDE_H, C['accent']),
        xml_rect(4, 'BottomPanel', 0, emu(5.2), SLIDE_W, emu(2.3), C['dark_navy']),
        xml_ph(10, 'CategoryLabel.Top', 'body', 10,
               emu(1.2), emu(1.5), emu(9), emu(0.45),
               15, F_BODY, C['accent'], True, 'l', 'CATEGORY LABEL'),
        xml_ph(11, 'Title.Center', 'ctrTitle', 0,
               emu(1.2), emu(2.1), emu(10.5), emu(1.5),
               48, F_HEAD, C['white'], True, 'l', 'プレゼンテーションタイトル'),
        xml_ph(12, 'Subtitle.Center', 'subTitle', 1,
               emu(1.2), emu(3.7), emu(9), emu(0.7),
               20, F_BODY, C['ice_blue'], False, 'l', 'サブタイトル・説明文'),
        xml_ph(13, 'Date.Bottom', 'body', 11,
               emu(1.2), emu(5.6), emu(8), emu(0.4),
               14, F_BODY, C['ice_blue'], False, 'l', '日付  |  組織名  |  部門名'),
        xml_ph(14, 'Footer.Bottom', 'body', 12,
               emu(1.2), emu(6.1), emu(4), emu(0.3),
               11, F_BODY, C['ice_blue'], False, 'l', 'Confidential'),
    ])

def L01_title_report():
    """Report cover: title + metadata panel on right."""
    return _build_layout('Title.1Title.Single+1Meta', 'title', [
        xml_rect(2, 'BG', 0, 0, SLIDE_W, SLIDE_H, C['navy']),
        xml_rect(3, 'AccentBar', 0, 0, emu(0.18), SLIDE_H, C['accent']),
        xml_roundrect(4, 'MetaPanel', emu(8.8), emu(1.2), emu(3.8), emu(4.5),
                      C['dark_navy'], '8000'),
        xml_ph(10, 'CategoryLabel.Top', 'body', 10,
               emu(1.2), emu(1.2), emu(6.5), emu(0.4),
               14, F_BODY, C['accent'], True, 'l', 'DATA ANALYSIS REPORT'),
        xml_ph(11, 'Title.Left', 'ctrTitle', 0,
               emu(1.2), emu(1.8), emu(7.0), emu(1.8),
               44, F_HEAD, C['white'], True, 'l', 'レポートタイトル'),
        xml_ph(12, 'Subtitle.Left', 'subTitle', 1,
               emu(1.2), emu(3.8), emu(7.0), emu(0.6),
               18, F_BODY, C['ice_blue'], False, 'l', 'サブタイトル'),
        xml_ph(13, 'Meta.Right', 'body', 11,
               emu(9.1), emu(1.6), emu(3.2), emu(3.8),
               13, F_BODY, C['ice_blue'], False, 'l',
               '作成日：\n作成者：\nバージョン：\nステータス：\n配布先：'),
        xml_ph(14, 'Footer.Bottom', 'body', 12,
               emu(1.2), emu(6.2), emu(6), emu(0.3),
               11, F_BODY, C['ice_blue'], False, 'l', 'Confidential'),
    ])

def L02_title_split():
    """Split: navy left panel with title, white right."""
    return _build_layout('Title.1Title.Single+1Summary', 'title', [
        xml_rect(2, 'LeftPanel', 0, 0, emu(7.5), SLIDE_H, C['navy']),
        xml_rect(3, 'RightPanel', emu(7.5), 0, emu(5.833), SLIDE_H, C['light_gray']),
        xml_rect(4, 'AccentBar', 0, 0, emu(0.18), SLIDE_H, C['accent']),
        xml_rect(5, 'CornerBlock', emu(7.5), 0, emu(0.4), emu(0.4), C['accent']),
        xml_ph(10, 'CategoryLabel.Left', 'body', 10,
               emu(1.2), emu(1.5), emu(5.5), emu(0.4),
               14, F_BODY, C['accent'], True, 'l', 'REPORT'),
        xml_ph(11, 'Title.Left', 'ctrTitle', 0,
               emu(1.2), emu(2.1), emu(5.8), emu(2.0),
               42, F_HEAD, C['white'], True, 'l', 'タイトル'),
        xml_ph(12, 'Subtitle.Left', 'subTitle', 1,
               emu(1.2), emu(4.3), emu(5.5), emu(0.6),
               16, F_BODY, C['ice_blue'], False, 'l', 'サブタイトル'),
        xml_ph(13, 'Summary.Right', 'body', 11,
               emu(8.0), emu(1.2), emu(4.8), emu(5.5),
               14, F_BODY, C['dark_text'], False, 'l',
               'サマリー・概要テキスト\nまたは画像を配置'),
        xml_ph(14, 'Date.Bottom', 'body', 12,
               emu(1.2), emu(6.2), emu(5), emu(0.3),
               11, F_BODY, C['ice_blue'], False, 'l', '日付  |  組織名'),
    ])


# ──────────────── GROUP 2: SECTION DIVIDERS (3) ────────────────

def L03_section_standard():
    """Dark bg, large section title + description."""
    return _build_layout('Section.1Title.Single', 'secHead', [
        xml_rect(2, 'BG', 0, 0, SLIDE_W, SLIDE_H, C['dark_navy']),
        xml_rect(3, 'AccentBar', 0, 0, emu(0.18), SLIDE_H, C['accent']),
        xml_rect(4, 'TopBlock', emu(1.2), emu(0.8), emu(0.6), emu(0.08), C['accent']),
        xml_ph(10, 'SectionLabel.Top', 'body', 10,
               emu(1.2), emu(1.1), emu(5), emu(0.4),
               13, F_BODY, C['accent'], True, 'l', 'SECTION'),
        xml_ph(11, 'Title.Center', 'body', 15,
               emu(1.2), emu(1.8), emu(10.5), emu(2.2),
               40, F_HEAD, C['white'], True, 'l', 'セクションタイトル'),
        xml_ph(12, 'Description.Bottom', 'body', 11,
               emu(1.2), emu(4.5), emu(10), emu(1.5),
               16, F_BODY, C['ice_blue'], False, 'l', '説明テキストを入力'),
    ])

def L04_section_numbered():
    """Large number + title side by side."""
    return _build_layout('SectionNav.1Title.Single', 'secHead', [
        xml_rect(2, 'BG', 0, 0, SLIDE_W, SLIDE_H, C['dark_navy']),
        xml_roundrect(3, 'NumPanel', emu(0.8), emu(1.0), emu(2.8), emu(5.0),
                      C['soft_navy'], '10000'),
        xml_ph(10, 'SectionNumber.Left', 'body', 10,
               emu(1.0), emu(1.5), emu(2.4), emu(2.5),
               80, F_HEAD, C['accent'], True, 'ctr', '01'),
        xml_ph(11, 'SectionLabel.Left', 'body', 11,
               emu(1.0), emu(4.0), emu(2.4), emu(0.4),
               12, F_BODY, C['ice_blue'], False, 'ctr', 'SECTION'),
        xml_ph(12, 'Title.Right', 'body', 15,
               emu(4.2), emu(1.5), emu(8.5), emu(2.0),
               38, F_HEAD, C['white'], True, 'l', 'セクションタイトル'),
        xml_ph(13, 'Description.Right', 'body', 12,
               emu(4.2), emu(4.0), emu(8.5), emu(2.0),
               16, F_BODY, C['ice_blue'], False, 'l', '説明テキストを入力'),
    ])

def L05_section_banner():
    """Accent banner across the middle."""
    return _build_layout('SectionBreak.1Title.Single', 'secHead', [
        xml_rect(2, 'TopBG', 0, 0, SLIDE_W, emu(2.8), C['dark_navy']),
        xml_rect(3, 'Banner', 0, emu(2.8), SLIDE_W, emu(2.2), C['navy']),
        xml_rect(4, 'BottomBG', 0, emu(5.0), SLIDE_W, emu(2.5), C['dark_navy']),
        xml_rect(5, 'AccentTop', 0, emu(2.8), SLIDE_W, emu(0.06), C['accent']),
        xml_rect(6, 'AccentBottom', 0, emu(4.94), SLIDE_W, emu(0.06), C['accent']),
        xml_ph(10, 'SectionLabel.Top', 'body', 10,
               emu(1.2), emu(1.2), emu(10), emu(0.4),
               13, F_BODY, C['accent'], True, 'l', 'SECTION 01'),
        xml_ph(11, 'Title.Center', 'body', 15,
               emu(1.2), emu(3.0), emu(10.5), emu(1.5),
               44, F_HEAD, C['white'], True, 'l', 'セクションタイトル'),
        xml_ph(12, 'Description.Bottom', 'body', 11,
               emu(1.2), emu(5.3), emu(10), emu(1.2),
               15, F_BODY, C['ice_blue'], False, 'l', '説明テキスト'),
    ])


# ──────────────── GROUP 3: STANDARD CONTENT (4) ────────────────

def L06_content_standard():
    """Header bar + full body area."""
    return _build_layout('Content.1Body.Single', 'obj', [
        *_header_bar(),
        xml_ph(20, 'Body.Center', 'body', 1,
               emu(0.8), emu(1.45), emu(11.7), emu(5.4),
               14, F_BODY, C['dark_text'], False, 'l', 'コンテンツを入力'),
        _slide_num(30),
    ])

def L07_content_sidebar():
    """Body area + right sidebar panel."""
    return _build_layout('Content.1Body.Single+1Notes', 'obj', [
        *_header_bar(),
        xml_roundrect(5, 'SidePanel', emu(9.4), emu(1.35), emu(3.5), emu(5.6),
                      C['panel_gray'], '8000'),
        xml_ph(20, 'Body.Left', 'body', 1,
               emu(0.8), emu(1.45), emu(8.2), emu(5.4),
               14, F_BODY, C['dark_text'], False, 'l', 'メインコンテンツ'),
        xml_ph(21, 'Notes.Right', 'body', 2,
               emu(9.7), emu(1.65), emu(3.0), emu(5.1),
               12, F_BODY, C['dark_text'], False, 'l',
               'サイドバー\n\n補足情報・メモ・\n関連リンク等'),
        _slide_num(30),
    ])

def L08_content_source():
    """Body area + source citation footer bar."""
    return _build_layout('Content.1Body.Single+1Source', 'obj', [
        *_header_bar(),
        xml_rect(5, 'FooterBar', 0, emu(6.5), SLIDE_W, emu(1.0), C['panel_gray']),
        xml_ph(20, 'Body.Center', 'body', 1,
               emu(0.8), emu(1.45), emu(11.7), emu(4.8),
               14, F_BODY, C['dark_text'], False, 'l', 'コンテンツを入力'),
        xml_ph(21, 'Source.Bottom', 'body', 2,
               emu(0.8), emu(6.65), emu(11.0), emu(0.6),
               10, F_BODY, C['mid_gray'], False, 'l',
               '出典：データソースを記載'),
        _slide_num(30),
    ])

def L09_content_callout():
    """Body area + top accent callout box."""
    return _build_layout('Content.1Body.Single+1Callout', 'obj', [
        *_header_bar(),
        xml_roundrect(5, 'Callout', emu(0.7), emu(1.35), emu(11.9), emu(1.2),
                      C['card_bg'], '6000'),
        xml_rect(6, 'CalloutAccent', emu(0.7), emu(1.35), emu(0.08), emu(1.2), C['accent']),
        xml_ph(20, 'Callout.Top', 'body', 2,
               emu(1.1), emu(1.5), emu(11.2), emu(0.85),
               13, F_BODY, C['navy'], False, 'l',
               'キーメッセージ・要約を記入'),
        xml_ph(21, 'Body.Center', 'body', 1,
               emu(0.8), emu(2.8), emu(11.7), emu(4.1),
               14, F_BODY, C['dark_text'], False, 'l', 'コンテンツを入力'),
        _slide_num(30),
    ])


# ──────────────── GROUP 4: MULTI-COLUMN (3) ────────────────

def L10_two_col():
    """Two equal columns."""
    return _build_layout('Column.2Body.Equal', 'twoObj', [
        *_header_bar(),
        xml_ph(20, 'Body1.Left', 'body', 1,
               emu(0.8), emu(1.45), emu(5.5), emu(5.4),
               14, F_BODY, C['dark_text'], False, 'l', '左カラム'),
        xml_ph(21, 'Body2.Right', 'body', 2,
               emu(6.8), emu(1.45), emu(5.7), emu(5.4),
               14, F_BODY, C['dark_text'], False, 'l', '右カラム'),
        _slide_num(30),
    ])

def L11_two_col_wide():
    """Wide left (65%) + narrow right (35%)."""
    return _build_layout('Column.2Body.MainSub', 'twoObj', [
        *_header_bar(),
        xml_ph(20, 'Body1.Left', 'body', 1,
               emu(0.8), emu(1.45), emu(7.4), emu(5.4),
               14, F_BODY, C['dark_text'], False, 'l', 'メインカラム'),
        xml_roundrect(5, 'RightPanel', emu(8.5), emu(1.35), emu(4.3), emu(5.6),
                      C['panel_gray'], '6000'),
        xml_ph(21, 'Body2.Right', 'body', 2,
               emu(8.8), emu(1.55), emu(3.8), emu(5.1),
               14, F_BODY, C['dark_text'], False, 'l', 'サブカラム'),
        _slide_num(30),
    ])

def L12_three_col():
    """Three equal columns."""
    return _build_layout('Column.3Body.Equal', 'threeObj', [
        *_header_bar(),
        xml_ph(20, 'Body1.Left', 'body', 1,
               emu(0.8), emu(1.45), emu(3.5), emu(5.4),
               14, F_BODY, C['dark_text'], False, 'l', 'カラム 1'),
        xml_ph(21, 'Body2.Center', 'body', 2,
               emu(4.7), emu(1.45), emu(3.5), emu(5.4),
               14, F_BODY, C['dark_text'], False, 'l', 'カラム 2'),
        xml_ph(22, 'Body3.Right', 'body', 3,
               emu(8.6), emu(1.45), emu(3.9), emu(5.4),
               14, F_BODY, C['dark_text'], False, 'l', 'カラム 3'),
        _slide_num(30),
    ])


# ──────────────── GROUP 5: KPI / METRICS (4) ────────────────

def L13_kpi_single():
    """Single hero stat center."""
    return _build_layout('KPI.1Value.Single', 'blank', [
        *_header_bar(),
        xml_roundrect(5, 'StatPanel', emu(3.0), emu(1.8), emu(7.3), emu(4.5),
                      C['panel_gray'], '10000'),
        xml_circle(6, 'AccentDot', emu(3.6), emu(2.4), emu(0.15), C['accent']),
        xml_ph(20, 'Value.Center', 'body', 1,
               emu(3.5), emu(2.2), emu(6.5), emu(2.0),
               72, F_HEAD, C['navy'], True, 'ctr', '00.0%'),
        xml_ph(21, 'ValueLabel.Center', 'body', 2,
               emu(3.5), emu(4.2), emu(6.5), emu(0.5),
               18, F_BODY, C['dark_text'], True, 'ctr', '指標名'),
        xml_ph(22, 'ValueDesc.Center', 'body', 3,
               emu(3.5), emu(4.8), emu(6.5), emu(1.0),
               13, F_BODY, C['mid_gray'], False, 'ctr', '補足説明・比較期間等'),
        _slide_num(30),
    ])

def L14_kpi_two():
    """Two stat cards side by side."""
    return _build_layout('KPI.2Value.Equal', 'blank', [
        *_header_bar(),
        xml_roundrect(5, 'Card1', emu(0.7), emu(1.5), emu(5.6), emu(5.2),
                      C['panel_gray'], '10000'),
        xml_roundrect(6, 'Card2', emu(6.7), emu(1.5), emu(5.8), emu(5.2),
                      C['panel_gray'], '10000'),
        xml_rect(7, 'Accent1', emu(0.7), emu(1.5), emu(5.6), emu(0.08), C['accent']),
        xml_rect(8, 'Accent2', emu(6.7), emu(1.5), emu(5.8), emu(0.08), C['teal']),
        xml_ph(20, 'Value1.Left', 'body', 1,
               emu(1.2), emu(2.0), emu(4.6), emu(1.8),
               60, F_HEAD, C['navy'], True, 'ctr', '00%'),
        xml_ph(21, 'ValueLabel1.Left', 'body', 2,
               emu(1.2), emu(3.9), emu(4.6), emu(2.5),
               14, F_BODY, C['dark_text'], False, 'ctr',
               '指標名\n説明テキスト'),
        xml_ph(22, 'Value2.Right', 'body', 3,
               emu(7.2), emu(2.0), emu(4.8), emu(1.8),
               60, F_HEAD, C['navy'], True, 'ctr', '00%'),
        xml_ph(23, 'ValueLabel2.Right', 'body', 4,
               emu(7.2), emu(3.9), emu(4.8), emu(2.5),
               14, F_BODY, C['dark_text'], False, 'ctr',
               '指標名\n説明テキスト'),
        _slide_num(30),
    ])

def L15_kpi_three():
    """Three stat cards in a row."""
    return _build_layout('KPI.3Value.Equal', 'blank', [
        *_header_bar(),
        xml_roundrect(5, 'Card1', emu(0.5), emu(1.5), emu(3.8), emu(5.2),
                      C['panel_gray'], '10000'),
        xml_roundrect(6, 'Card2', emu(4.6), emu(1.5), emu(4.0), emu(5.2),
                      C['panel_gray'], '10000'),
        xml_roundrect(7, 'Card3', emu(8.9), emu(1.5), emu(3.9), emu(5.2),
                      C['panel_gray'], '10000'),
        xml_rect(8, 'A1', emu(0.5), emu(1.5), emu(3.8), emu(0.06), C['accent']),
        xml_rect(9, 'A2', emu(4.6), emu(1.5), emu(4.0), emu(0.06), C['teal']),
        xml_rect(10, 'A3', emu(8.9), emu(1.5), emu(3.9), emu(0.06), C['amber']),
        xml_ph(20, 'Value1.Left', 'body', 1,
               emu(0.8), emu(2.0), emu(3.2), emu(1.5),
               52, F_HEAD, C['navy'], True, 'ctr', '00%'),
        xml_ph(21, 'ValueLabel1.Left', 'body', 2,
               emu(0.8), emu(3.6), emu(3.2), emu(2.8),
               13, F_BODY, C['dark_text'], False, 'ctr', '指標名\n説明'),
        xml_ph(22, 'Value2.Center', 'body', 3,
               emu(4.9), emu(2.0), emu(3.4), emu(1.5),
               52, F_HEAD, C['navy'], True, 'ctr', '00%'),
        xml_ph(23, 'ValueLabel2.Center', 'body', 4,
               emu(4.9), emu(3.6), emu(3.4), emu(2.8),
               13, F_BODY, C['dark_text'], False, 'ctr', '指標名\n説明'),
        xml_ph(24, 'Value3.Right', 'body', 5,
               emu(9.2), emu(2.0), emu(3.3), emu(1.5),
               52, F_HEAD, C['navy'], True, 'ctr', '00%'),
        xml_ph(25, 'ValueLabel3.Right', 'body', 6,
               emu(9.2), emu(3.6), emu(3.3), emu(2.8),
               13, F_BODY, C['dark_text'], False, 'ctr', '指標名\n説明'),
        _slide_num(31),
    ])

def L16_kpi_grid():
    """2x2 grid of four stat cards."""
    return _build_layout('KPI.4Value.Grid', 'blank', [
        *_header_bar(),
        # Top row
        xml_roundrect(5, 'TL', emu(0.5), emu(1.4), emu(5.8), emu(2.5),
                      C['panel_gray'], '8000'),
        xml_roundrect(6, 'TR', emu(6.7), emu(1.4), emu(5.9), emu(2.5),
                      C['panel_gray'], '8000'),
        # Bottom row
        xml_roundrect(7, 'BL', emu(0.5), emu(4.2), emu(5.8), emu(2.5),
                      C['panel_gray'], '8000'),
        xml_roundrect(8, 'BR', emu(6.7), emu(4.2), emu(5.9), emu(2.5),
                      C['panel_gray'], '8000'),
        # Accent strips
        xml_rect(9, 'A1', emu(0.5), emu(1.4), emu(0.06), emu(2.5), C['accent']),
        xml_rect(10, 'A2', emu(6.7), emu(1.4), emu(0.06), emu(2.5), C['teal']),
        xml_rect(11, 'A3', emu(0.5), emu(4.2), emu(0.06), emu(2.5), C['amber']),
        xml_rect(12, 'A4', emu(6.7), emu(4.2), emu(0.06), emu(2.5), C['accent_dark']),
        # Placeholders
        xml_ph(20, 'Value1.TopLeft', 'body', 1,
               emu(1.0), emu(1.6), emu(4.8), emu(2.0),
               14, F_BODY, C['dark_text'], False, 'l',
               '指標 1\n値・説明'),
        xml_ph(21, 'Value2.TopRight', 'body', 2,
               emu(7.2), emu(1.6), emu(5.0), emu(2.0),
               14, F_BODY, C['dark_text'], False, 'l',
               '指標 2\n値・説明'),
        xml_ph(22, 'Value3.BottomLeft', 'body', 3,
               emu(1.0), emu(4.4), emu(4.8), emu(2.0),
               14, F_BODY, C['dark_text'], False, 'l',
               '指標 3\n値・説明'),
        xml_ph(23, 'Value4.BottomRight', 'body', 4,
               emu(7.2), emu(4.4), emu(5.0), emu(2.0),
               14, F_BODY, C['dark_text'], False, 'l',
               '指標 4\n値・説明'),
        _slide_num(31),
    ])


# ──────────────── GROUP 6: CHART / VISUAL (3) ────────────────

def L17_chart_full():
    """Full width chart area."""
    return _build_layout('Chart.1Chart.Single', 'blank', [
        *_header_bar(),
        xml_roundrect(5, 'ChartArea', emu(0.5), emu(1.4), emu(12.3), emu(5.4),
                      C['panel_gray'], '6000'),
        xml_ph(20, 'Body.Center', 'body', 1,
               emu(0.8), emu(1.6), emu(11.7), emu(5.0),
               14, F_BODY, C['mid_gray'], False, 'ctr',
               'グラフ・チャートを配置'),
        _slide_num(30),
    ])

def L18_chart_text():
    """Chart left + analysis text right."""
    return _build_layout('Chart.1Chart.Single+1Analysis', 'blank', [
        *_header_bar(),
        xml_roundrect(5, 'ChartArea', emu(0.5), emu(1.4), emu(7.5), emu(5.4),
                      C['panel_gray'], '6000'),
        xml_ph(20, 'Body.Left', 'body', 1,
               emu(0.8), emu(1.6), emu(6.9), emu(5.0),
               14, F_BODY, C['mid_gray'], False, 'ctr',
               'グラフ・チャートを配置'),
        xml_ph(21, 'Analysis.Right', 'body', 2,
               emu(8.4), emu(1.5), emu(4.4), emu(5.3),
               14, F_BODY, C['dark_text'], False, 'l',
               '分析テキスト\n\n主要な発見事項や\n考察を記入'),
        _slide_num(30),
    ])

def L19_chart_dual():
    """Two chart areas side by side."""
    return _build_layout('Chart.2Chart.Equal', 'blank', [
        *_header_bar(),
        xml_roundrect(5, 'ChartL', emu(0.5), emu(1.4), emu(5.9), emu(5.4),
                      C['panel_gray'], '6000'),
        xml_roundrect(6, 'ChartR', emu(6.7), emu(1.4), emu(5.9), emu(5.4),
                      C['panel_gray'], '6000'),
        xml_ph(20, 'Body1.Left', 'body', 1,
               emu(0.8), emu(1.6), emu(5.3), emu(5.0),
               14, F_BODY, C['mid_gray'], False, 'ctr',
               'チャート 1'),
        xml_ph(21, 'Body2.Right', 'body', 2,
               emu(7.0), emu(1.6), emu(5.3), emu(5.0),
               14, F_BODY, C['mid_gray'], False, 'ctr',
               'チャート 2'),
        _slide_num(30),
    ])


# ──────────────── GROUP 7: TABLE (2) ────────────────

def L20_table_full():
    """Full width table area."""
    return _build_layout('Table.1Table.Single+1Source', 'blank', [
        *_header_bar(),
        xml_ph(20, 'Body.Center', 'body', 1,
               emu(0.5), emu(1.45), emu(12.3), emu(5.0),
               14, F_BODY, C['dark_text'], False, 'l',
               'テーブル・表を配置'),
        xml_ph(21, 'Source.Bottom', 'body', 2,
               emu(0.8), emu(6.6), emu(11.0), emu(0.4),
               10, F_BODY, C['mid_gray'], False, 'l', '出典：'),
        _slide_num(30),
    ])

def L21_table_notes():
    """Table + notes panel."""
    return _build_layout('Table.1Table.Single+1Notes', 'blank', [
        *_header_bar(),
        xml_ph(20, 'Body.Left', 'body', 1,
               emu(0.5), emu(1.45), emu(8.5), emu(5.4),
               14, F_BODY, C['dark_text'], False, 'l',
               'テーブル・表を配置'),
        xml_roundrect(5, 'NotesPanel', emu(9.3), emu(1.35), emu(3.5), emu(5.6),
                      C['card_bg'], '8000'),
        xml_rect(6, 'NotesAccent', emu(9.3), emu(1.35), emu(0.06), emu(5.6), C['accent']),
        xml_ph(21, 'Notes.Right', 'body', 2,
               emu(9.7), emu(1.55), emu(2.9), emu(5.1),
               12, F_BODY, C['dark_text'], False, 'l',
               'コメント・注記\n\nテーブルの補足情報'),
        _slide_num(30),
    ])


# ──────────────── GROUP 8: COMPARISON (2) ────────────────

def L22_comparison_side():
    """Two labeled panels with visual divider."""
    return _build_layout('Compare.2Option.Versus', 'blank', [
        *_header_bar(),
        xml_roundrect(5, 'PanelL', emu(0.5), emu(1.4), emu(5.8), emu(5.4),
                      C['panel_gray'], '8000'),
        xml_roundrect(6, 'PanelR', emu(6.7), emu(1.4), emu(5.9), emu(5.4),
                      C['panel_gray'], '8000'),
        xml_rect(7, 'AccentL', emu(0.5), emu(1.4), emu(5.8), emu(0.06), C['accent']),
        xml_rect(8, 'AccentR', emu(6.7), emu(1.4), emu(5.9), emu(0.06), C['teal']),
        xml_circle(9, 'VS', emu(6.35 * 914400), emu(4.1 * 914400), emu(0.3), C['navy']),
        xml_ph(20, 'OptionLabel1.Left', 'body', 1,
               emu(0.8), emu(1.7), emu(5.2), emu(0.5),
               16, F_BODY, C['accent'], True, 'l', 'オプション A'),
        xml_ph(21, 'OptionContent1.Left', 'body', 2,
               emu(0.8), emu(2.3), emu(5.2), emu(4.2),
               14, F_BODY, C['dark_text'], False, 'l', '内容を入力'),
        xml_ph(22, 'OptionLabel2.Right', 'body', 3,
               emu(7.0), emu(1.7), emu(5.3), emu(0.5),
               16, F_BODY, C['teal'], True, 'l', 'オプション B'),
        xml_ph(23, 'OptionContent2.Right', 'body', 4,
               emu(7.0), emu(2.3), emu(5.3), emu(4.2),
               14, F_BODY, C['dark_text'], False, 'l', '内容を入力'),
        _slide_num(30),
    ])

def L23_comparison_matrix():
    """Comparison matrix — simple layout with content area for table."""
    return _build_layout('Compare.1Matrix.Single', 'blank', [
        *_header_bar(),
        # Content area background panel
        xml_rect(5, 'ContentBG', emu(0.5), emu(1.4), emu(12.3), emu(5.5), C['white']),
        _slide_num(30),
    ])


# ──────────────── GROUP 9: PROCESS / TIMELINE (2) ────────────────

def L24_timeline_horizontal():
    """4 step horizontal process."""
    return _build_layout('Process.4Step.Sequential', 'blank', [
        *_header_bar(),
        # Step connector line
        xml_rect(5, 'ConnLine', emu(1.8), emu(2.9), emu(9.7), emu(0.04), C['ice_blue']),
        # Step circles
        xml_circle(6, 'C1', emu(2.0 * 914400), emu(2.95 * 914400), emu(0.28), C['accent']),
        xml_circle(7, 'C2', emu(5.2 * 914400), emu(2.95 * 914400), emu(0.28), C['accent']),
        xml_circle(8, 'C3', emu(8.4 * 914400), emu(2.95 * 914400), emu(0.28), C['accent']),
        xml_circle(9, 'C4', emu(11.5 * 914400), emu(2.95 * 914400), emu(0.28), C['accent']),
        # Step labels
        xml_ph(20, 'Step1.Left', 'body', 1,
               emu(0.5), emu(3.6), emu(2.8), emu(3.0),
               13, F_BODY, C['dark_text'], False, 'ctr',
               'STEP 1\nタイトル\n説明テキスト'),
        xml_ph(21, 'Step2.CenterLeft', 'body', 2,
               emu(3.7), emu(3.6), emu(2.8), emu(3.0),
               13, F_BODY, C['dark_text'], False, 'ctr',
               'STEP 2\nタイトル\n説明テキスト'),
        xml_ph(22, 'Step3.CenterRight', 'body', 3,
               emu(6.9), emu(3.6), emu(2.8), emu(3.0),
               13, F_BODY, C['dark_text'], False, 'ctr',
               'STEP 3\nタイトル\n説明テキスト'),
        xml_ph(23, 'Step4.Right', 'body', 4,
               emu(10.0), emu(3.6), emu(2.8), emu(3.0),
               13, F_BODY, C['dark_text'], False, 'ctr',
               'STEP 4\nタイトル\n説明テキスト'),
        _slide_num(30),
    ])

def L25_process_vertical():
    """Vertical process steps."""
    return _build_layout('Process.3Step.Sequential', 'blank', [
        *_header_bar(),
        # Step connectors
        xml_rect(5, 'VLine', emu(2.0), emu(1.8), emu(0.04), emu(4.8), C['ice_blue']),
        # Step indicators
        xml_circle(6, 'S1', emu(2.02 * 914400), emu(2.2 * 914400), emu(0.22), C['accent']),
        xml_circle(7, 'S2', emu(2.02 * 914400), emu(3.8 * 914400), emu(0.22), C['accent']),
        xml_circle(8, 'S3', emu(2.02 * 914400), emu(5.4 * 914400), emu(0.22), C['accent']),
        # Step text
        xml_ph(20, 'Step1.Top', 'body', 1,
               emu(2.8), emu(1.7), emu(9.5), emu(1.2),
               14, F_BODY, C['dark_text'], False, 'l',
               'ステップ 1：タイトル\n説明テキスト'),
        xml_ph(21, 'Step2.Center', 'body', 2,
               emu(2.8), emu(3.3), emu(9.5), emu(1.2),
               14, F_BODY, C['dark_text'], False, 'l',
               'ステップ 2：タイトル\n説明テキスト'),
        xml_ph(22, 'Step3.Bottom', 'body', 3,
               emu(2.8), emu(4.9), emu(9.5), emu(1.2),
               14, F_BODY, C['dark_text'], False, 'l',
               'ステップ 3：タイトル\n説明テキスト'),
        _slide_num(30),
    ])


# ──────────────── GROUP 10: SUMMARY (2) ────────────────

def L26_agenda():
    """Agenda / table of contents."""
    return _build_layout('Summary.1Agenda.Single', 'blank', [
        xml_rect(2, 'BG', 0, 0, SLIDE_W, SLIDE_H, C['white']),
        xml_rect(3, 'LeftPanel', 0, 0, emu(4.5), SLIDE_H, C['navy']),
        xml_rect(4, 'AccentBar', 0, 0, emu(0.18), SLIDE_H, C['accent']),
        xml_ph(10, 'Title.Left', 'body', 15,
               emu(0.6), emu(1.0), emu(3.5), emu(1.5),
               36, F_HEAD, C['white'], True, 'l', 'Agenda'),
        xml_ph(11, 'MeetingInfo.Left', 'body', 16,
               emu(0.6), emu(2.8), emu(3.5), emu(3.5),
               14, F_BODY, C['ice_blue'], False, 'l',
               '日時：\n場所：\n所要時間：'),
        xml_ph(20, 'AgendaItems.Right', 'body', 1,
               emu(5.0), emu(0.8), emu(7.5), emu(6.0),
               15, F_BODY, C['dark_text'], False, 'l',
               '01  議題タイトル\n\n02  議題タイトル\n\n03  議題タイトル\n\n04  議題タイトル\n\n05  議題タイトル'),
    ])

def L27_executive_summary():
    """Executive summary with key findings."""
    return _build_layout('Summary.2Block.Equal', 'blank', [
        *_header_bar(),
        xml_roundrect(5, 'FindingsPanel', emu(0.5), emu(1.4), emu(7.8), emu(5.5),
                      C['panel_gray'], '8000'),
        xml_roundrect(6, 'RecPanel', emu(8.6), emu(1.4), emu(4.1), emu(5.5),
                      C['card_bg'], '8000'),
        xml_rect(7, 'RecAccent', emu(8.6), emu(1.4), emu(0.06), emu(5.5), C['accent']),
        xml_ph(20, 'Body1.Left', 'body', 1,
               emu(0.9), emu(1.7), emu(7.0), emu(5.0),
               14, F_BODY, C['dark_text'], False, 'l',
               '主要な発見事項\n\n1. \n2. \n3. '),
        xml_ph(21, 'Body2.Right', 'body', 2,
               emu(9.0), emu(1.7), emu(3.4), emu(5.0),
               13, F_BODY, C['dark_text'], False, 'l',
               '推奨アクション\n\n優先度：高\n\n対応策を記入'),
        _slide_num(30),
    ])


# ──────────────── GROUP 11: CLOSING (2) ────────────────

def L28_closing_thankyou():
    """Thank you + contact."""
    return _build_layout('Closing.1Message.Single', 'blank', [
        xml_rect(2, 'BG', 0, 0, SLIDE_W, SLIDE_H, C['navy']),
        xml_rect(3, 'AccentBar', 0, 0, emu(0.18), SLIDE_H, C['accent']),
        xml_rect(4, 'BottomPanel', 0, emu(4.8), SLIDE_W, emu(2.7), C['dark_navy']),
        xml_ph(10, 'CategoryLabel.Top', 'body', 10,
               emu(1.2), emu(1.5), emu(10), emu(0.5),
               18, F_BODY, C['accent'], True, 'l', 'THANK YOU'),
        xml_ph(11, 'Title.Center', 'ctrTitle', 0,
               emu(1.2), emu(2.2), emu(10.5), emu(1.6),
               42, F_HEAD, C['white'], True, 'l', 'ご質問をお待ちしています'),
        xml_ph(12, 'PresenterName.Bottom', 'body', 11,
               emu(1.2), emu(5.2), emu(8), emu(0.5),
               16, F_BODY, C['ice_blue'], False, 'l', '担当者名  |  役職'),
        xml_ph(13, 'Contact.Bottom', 'body', 12,
               emu(1.2), emu(5.8), emu(8), emu(0.8),
               13, F_BODY, C['ice_blue'], False, 'l',
               'email@example.com  |  内線: 0000\nSlack: #channel'),
    ])

def L29_closing_nextsteps():
    """Next steps + CTA."""
    return _build_layout('Closing.1Steps.Single+1Notes', 'blank', [
        xml_rect(2, 'BG', 0, 0, SLIDE_W, SLIDE_H, C['navy']),
        xml_rect(3, 'AccentBar', 0, 0, emu(0.18), SLIDE_H, C['accent']),
        xml_roundrect(4, 'StepsPanel', emu(0.8), emu(1.0), emu(7.5), emu(5.5),
                      C['dark_navy'], '10000'),
        xml_roundrect(5, 'CTAPanel', emu(8.8), emu(1.0), emu(4.0), emu(5.5),
                      C['soft_navy'], '10000'),
        xml_ph(10, 'Title.Left', 'body', 15,
               emu(1.2), emu(1.3), emu(6.8), emu(0.7),
               28, F_HEAD, C['white'], True, 'l', 'Next Steps'),
        xml_ph(20, 'ActionSteps.Left', 'body', 1,
               emu(1.2), emu(2.2), emu(6.8), emu(4.0),
               15, F_BODY, C['ice_blue'], False, 'l',
               '1. アクション項目\n\n2. アクション項目\n\n3. アクション項目'),
        xml_ph(21, 'Notes.Right', 'body', 2,
               emu(9.1), emu(1.3), emu(3.4), emu(5.0),
               14, F_BODY, C['white'], False, 'l',
               '期限\n\n担当者\n\n連絡先\n\n次回MTG日程'),
    ])


# ════════════════════════════════════════════
# LAYOUT REGISTRY
# ════════════════════════════════════════════

ALL_LAYOUTS = [
    L00_title_standard,      # 0
    L01_title_report,        # 1
    L02_title_split,         # 2
    L03_section_standard,    # 3
    L04_section_numbered,    # 4
    L05_section_banner,      # 5
    L06_content_standard,    # 6
    L07_content_sidebar,     # 7
    L08_content_source,      # 8
    L09_content_callout,     # 9
    L10_two_col,             # 10
    L11_two_col_wide,        # 11  (overwrite slot 11, last default)
    L12_three_col,           # 12  (OPC addition from here)
    L13_kpi_single,          # 13
    L14_kpi_two,             # 14
    L15_kpi_three,           # 15
    L16_kpi_grid,            # 16
    L17_chart_full,          # 17
    L18_chart_text,          # 18
    L19_chart_dual,          # 19
    L20_table_full,          # 20
    L21_table_notes,         # 21
    L22_comparison_side,     # 22
    L23_comparison_matrix,   # 23
    L24_timeline_horizontal, # 24
    L25_process_vertical,    # 25
    L26_agenda,              # 26
    L27_executive_summary,   # 27
    L28_closing_thankyou,    # 28
    L29_closing_nextsteps,   # 29
]


# ════════════════════════════════════════════
# INJECT LAYOUTS
# ════════════════════════════════════════════

def inject_layouts(prs):
    """Overwrite first 11, add remaining via OPC."""
    existing = list(prs.slide_layouts)
    master_part = prs.slide_masters[0].part
    master_elem = master_part._element
    layout_id_lst = master_elem.find(_tag('p:sldLayoutIdLst'))
    max_id = max(int(e.get('id')) for e in layout_id_lst)

    RT_LAYOUT = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout'
    RT_MASTER = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster'
    ct = 'application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml'

    for i, builder in enumerate(ALL_LAYOUTS):
        layout_xml = builder()

        if i < len(existing):
            # Overwrite existing slot
            target = existing[i]._element
            for child in list(target):
                target.remove(child)
            for child in list(layout_xml):
                target.append(child)
            for k, v in layout_xml.attrib.items():
                target.attrib[k] = v
            print(f"  ✓ [{i:2d}] overwrite → {builder.__name__}")
        else:
            # Add new via OPC
            xml_bytes = etree.tostring(layout_xml, xml_declaration=True,
                                       encoding='UTF-8', standalone=True)
            elem = parse_xml(xml_bytes)
            partname = PackURI(f'/ppt/slideLayouts/slideLayout{i+1}.xml')
            new_part = SlideLayoutPart(partname, ct, prs.part.package, elem)
            rId = master_part.relate_to(new_part, RT_LAYOUT)
            new_part.relate_to(master_part, RT_MASTER)
            new_entry = etree.SubElement(layout_id_lst, _tag('p:sldLayoutId'))
            new_entry.set('id', str(max_id + i + 1))
            new_entry.set(_tag('r:id'), rId)
            print(f"  + [{i:2d}] OPC add   → {builder.__name__}")

    return prs


# ════════════════════════════════════════════
# SET PLACEHOLDER TEXT
# ════════════════════════════════════════════

def set_ph(slide, idx, text):
    """Set TEXT ONLY on a placeholder. No style overrides — master styles apply."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            ph.text = text
            return ph
    return None


def add_matrix_table(slide, rows_data, x, y, width, height):
    """Add a styled comparison matrix table to a slide.

    rows_data: list of lists, first row is header, first col of each row is row label.
    Example: [['', 'A社', 'B社', 'C社'],
              ['機能性', '◎', '○', '△'], ...]
    """
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor

    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, Emu(x), Emu(y), Emu(width), Emu(height))
    table = table_shape.table

    # Column widths: first col wider for labels
    label_w = int(width * 0.28)
    data_w = int((width - label_w) / (n_cols - 1))
    table.columns[0].width = Emu(label_w)
    for ci in range(1, n_cols):
        table.columns[ci].width = Emu(data_w)

    for ri, row_data in enumerate(rows_data):
        for ci, cell_text in enumerate(row_data):
            cell = table.cell(ri, ci)
            cell.text = cell_text

            # Style paragraph
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(11)
                para.font.name = F_BODY
                if ri == 0:
                    # Header row
                    para.font.color.rgb = RGBColor.from_string(C['white'])
                    para.font.bold = True
                    para.alignment = 1  # center
                elif ci == 0:
                    # Row label
                    para.font.color.rgb = RGBColor.from_string(C['dark_text'])
                    para.font.bold = True
                else:
                    # Data cell
                    para.font.color.rgb = RGBColor.from_string(C['dark_text'])
                    para.alignment = 1  # center

            # Cell fill
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor.from_string(C['navy'])
            elif ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor.from_string(C['panel_gray'])
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor.from_string(C['white'])

            # Remove default borders by making them transparent
            _set_cell_borders(cell, C['panel_gray'])

    return table_shape


def _set_cell_borders(cell, color):
    """Set thin, subtle borders on a table cell."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    borders_xml = (
        f'<a:lnL w="6350" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill></a:lnL>'
        f'<a:lnR w="6350" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill></a:lnR>'
        f'<a:lnT w="6350" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill></a:lnT>'
        f'<a:lnB w="6350" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill></a:lnB>'
    )
    for border_xml in [borders_xml]:
        for ln_tag in ['lnL', 'lnR', 'lnT', 'lnB']:
            existing = tcPr.findall(f'{{{NS["a"]}}}{ln_tag}')
            for e in existing:
                tcPr.remove(e)
    for ln_tag in ['lnL', 'lnR', 'lnT', 'lnB']:
        ln_xml = (
            f'<a:{ln_tag} w="6350" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
            f'</a:{ln_tag}>'
        )
        tcPr.append(etree.fromstring(ln_xml))


# ════════════════════════════════════════════
# BUILD TEMPLATE
# ════════════════════════════════════════════

if __name__ == '__main__':
    print("Creating 30-layout template...")
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)

    print("\nInjecting layouts...")
    inject_layouts(prs)

    layouts = list(prs.slide_layouts)
    print(f"\nTotal layouts available: {len(layouts)}")

    # ─── Demo slides: ALL 30 layouts, TEXT ONLY (styles from master) ───
    print("\nAdding demo slides (text only, inheriting master styles)...")

    # 0 — タイトル：標準
    s = prs.slides.add_slide(layouts[0])
    set_ph(s, 10, 'DATA ANALYSIS REPORT')
    set_ph(s, 0,  'NextGen CRM Platform\n分析レポート')
    set_ph(s, 1,  '次世代顧客管理プラットフォーム開発プロジェクト')
    set_ph(s, 11, '2026年3月24日  |  株式会社サンプルテック  |  DX推進本部')
    set_ph(s, 12, 'Confidential')

    # 1 — タイトル：レポート表紙
    s = prs.slides.add_slide(layouts[1])
    set_ph(s, 10, 'DATA ANALYSIS REPORT')
    set_ph(s, 0,  'CRM移行\n影響度分析')
    set_ph(s, 1,  'Phase 1 ユーザー行動分析')
    set_ph(s, 11, '作成日：2026-03-24\n作成者：分析チーム\nバージョン：v2.1\nステータス：Final\n配布先：経営会議')
    set_ph(s, 12, 'Confidential')

    # 2 — タイトル：スプリット
    s = prs.slides.add_slide(layouts[2])
    set_ph(s, 10, 'REPORT')
    set_ph(s, 0,  '四半期レビュー\n2026 Q1')
    set_ph(s, 1,  '営業本部 パフォーマンス分析')
    set_ph(s, 11, 'エグゼクティブサマリー\n\n売上は前年比 +12% で推移。\n新規獲得は目標の 95% を達成。\nチャーン率は改善傾向。')
    set_ph(s, 12, '2026年3月  |  営業企画部')

    # 3 — セクション：標準
    s = prs.slides.add_slide(layouts[3])
    set_ph(s, 10, 'SECTION 01')
    set_ph(s, 15, '現状分析')
    set_ph(s, 11, '既存CRMの利用状況とユーザー行動パターンを\nデータに基づいて分析します。')

    # 4 — セクション：ナンバー付
    s = prs.slides.add_slide(layouts[4])
    set_ph(s, 10, '02')
    set_ph(s, 11, 'SECTION')
    set_ph(s, 15, 'データ分析結果')
    set_ph(s, 12, '収集した1年分のユーザーログデータから\n主要なインサイトを抽出しました。')

    # 5 — セクション：バナー
    s = prs.slides.add_slide(layouts[5])
    set_ph(s, 10, 'SECTION 03')
    set_ph(s, 15, 'システム比較と提案')
    set_ph(s, 11, '現行CRMと候補3システムの詳細比較を行い、\n最適な移行先を提案します。')

    # 6 — コンテンツ：標準
    s = prs.slides.add_slide(layouts[6])
    set_ph(s, 15, '利用状況サマリー')
    set_ph(s, 16, 'Usage Summary')
    set_ph(s, 1,  '月間アクティブユーザー数\n全ライセンスユーザー5,200名のうち、月間でログインするユーザーは平均3,800名（73%）。\n部門別では営業部門が最も高い利用率（92%）を示し、管理部門が最も低い（45%）。\n\n主要機能の利用頻度\n顧客検索：日平均12,000回\nレポート作成：日平均800回\nダッシュボード閲覧：日平均2,400回')

    # 7 — コンテンツ：サイドバー付
    s = prs.slides.add_slide(layouts[7])
    set_ph(s, 15, 'ユーザー満足度調査')
    set_ph(s, 16, 'User Satisfaction Survey')
    set_ph(s, 1,  '調査結果の概要\n\n回答数：1,200名（回答率 68%）\n\n総合満足度は5段階中3.2で前年から0.3pt低下。\n特に「レスポンス速度」と「モバイル対応」の\n2項目で不満が集中している。\n\n部門別では営業部門の満足度が最も低く、\nフィールド作業でのモバイル利用ニーズが\n満たされていないことが主因。')
    set_ph(s, 2,  '関連データ\n\nNPS: +12\n前年: +18\n\n回答率推移:\n2024: 72%\n2025: 68%\n\n次回調査:\n2026年9月')

    # 8 — コンテンツ：出典バー付
    s = prs.slides.add_slide(layouts[8])
    set_ph(s, 15, '市場動向')
    set_ph(s, 16, 'Market Trends')
    set_ph(s, 1,  'CRM市場の成長\n\nグローバルCRM市場は2025年に約800億ドル規模に到達。\nAI搭載CRMのシェアは前年比 +35% で急成長中。\n\n国内市場ではクラウド型CRMの導入率が\n前年の62%から78%に上昇。\nオンプレミス型からの移行が加速している。')
    set_ph(s, 2,  '出典：Gartner Market Report 2025, IDC Japan CRM調査 2025年版')

    # 9 — コンテンツ：コールアウト付
    s = prs.slides.add_slide(layouts[9])
    set_ph(s, 15, 'パフォーマンス分析')
    set_ph(s, 16, 'Performance Analysis')
    set_ph(s, 2,  '重要：現行システムの応答速度は平均3.2秒で、ユーザー満足度低下の主因')
    set_ph(s, 1,  '応答時間の推移\n2023年Q1から四半期ごとに平均レスポンスタイムが悪化。\n特にレポート生成機能は5秒以上かかるケースが40%を占める。\n\n影響度\nアンケート調査（n=1,200）で「遅い」と回答した割合は68%。\n業務効率への影響として、1人あたり週平均45分のロスが発生。')

    # 10 — 2カラム：均等
    s = prs.slides.add_slide(layouts[10])
    set_ph(s, 15, 'スコープ定義')
    set_ph(s, 16, 'In Scope / Out of Scope')
    set_ph(s, 1,  '対象範囲（In Scope）\n\n▸ 顧客データ統合基盤\n▸ AI分析エンジン\n▸ 営業支援モジュール\n▸ モバイルアプリ\n▸ 管理者ダッシュボード')
    set_ph(s, 2,  '対象外（Out of Scope）\n\n▸ 基幹系システム（ERP）刷新\n▸ コールセンターシステム\n▸ 海外拠点対応\n▸ 5年以前のデータ移行')

    # 11 — 2カラム：ワイド左
    s = prs.slides.add_slide(layouts[11])
    set_ph(s, 15, 'リスク分析')
    set_ph(s, 16, 'Risk Assessment')
    set_ph(s, 1,  'プロジェクトリスク一覧\n\n1. データ移行の品質リスク\n既存データの整合性チェックに想定以上の工数がかかる可能性。\n\n2. ユーザー定着リスク\n新UIへの習熟に時間がかかり、一時的な生産性低下が予想される。\n\n3. ベンダーロックイン\nクラウドサービスへの依存度が高まる。')
    set_ph(s, 2,  '対策方針\n\n移行ツールの\n事前検証を\n2段階で実施\n\nトレーニング\nプログラムを\n並行準備')

    # 12 — 3カラム
    s = prs.slides.add_slide(layouts[12])
    set_ph(s, 15, 'プロジェクト目標')
    set_ph(s, 16, 'Objectives')
    set_ph(s, 1,  'ビジネス目標\n\n▸ 顧客LTV +20%\n▸ リードタイム -30%\n▸ チャーン率 -15%')
    set_ph(s, 2,  '技術目標\n\n▸ 応答速度 <1秒\n▸ 可用性 99.9%\n▸ データ統合率 100%')
    set_ph(s, 3,  '組織目標\n\n▸ モバイル活用率 90%+\n▸ 部門間データ共有\n▸ AI意思決定速度 2倍')

    # 13 — KPI：シングル
    s = prs.slides.add_slide(layouts[13])
    set_ph(s, 15, '顧客満足度スコア')
    set_ph(s, 16, 'Customer Satisfaction Score')
    set_ph(s, 1,  '73.2%')
    set_ph(s, 2,  'NPS（Net Promoter Score）')
    set_ph(s, 3,  '前年比 +5.8pt  |  目標値: 80%  |  業界平均: 68%')

    # 14 — KPI：2カード
    s = prs.slides.add_slide(layouts[14])
    set_ph(s, 15, '重要指標')
    set_ph(s, 16, 'Key Metrics')
    set_ph(s, 1,  '+20%')
    set_ph(s, 2,  '顧客LTV\n前年比成長率')
    set_ph(s, 3,  '-30%')
    set_ph(s, 4,  '営業リードタイム\n短縮率')

    # 15 — KPI：3カード
    s = prs.slides.add_slide(layouts[15])
    set_ph(s, 15, '主要パフォーマンス指標')
    set_ph(s, 16, 'Key Performance Indicators')
    set_ph(s, 1,  '+20%')
    set_ph(s, 2,  '顧客LTV\n2027年度末目標')
    set_ph(s, 3,  '-30%')
    set_ph(s, 4,  '営業リードタイム\n初回接触〜成約')
    set_ph(s, 5,  '<1秒')
    set_ph(s, 6,  'システム応答速度\n99パーセンタイル')

    # 16 — KPI：4グリッド
    s = prs.slides.add_slide(layouts[16])
    set_ph(s, 15, '四半期ダッシュボード')
    set_ph(s, 16, 'Q1 2026 Dashboard')
    set_ph(s, 1,  '売上達成率\n98.5%\n目標比 -1.5pt')
    set_ph(s, 2,  '新規顧客数\n142社\n前年同期比 +23%')
    set_ph(s, 3,  'チャーン率\n2.1%\n目標: 3.0%以下')
    set_ph(s, 4,  'ARPU\n¥45,200\n前年同期比 +8%')

    # 17 — チャート：フル幅
    s = prs.slides.add_slide(layouts[17])
    set_ph(s, 15, '売上推移')
    set_ph(s, 16, 'Revenue Trend')
    set_ph(s, 1,  '（ここに売上推移のグラフを配置）')

    # 18 — チャート＋テキスト
    s = prs.slides.add_slide(layouts[18])
    set_ph(s, 15, 'ユーザー行動トレンド')
    set_ph(s, 16, 'User Behavior Trends')
    set_ph(s, 1,  '（月別アクティブユーザー数の\n推移グラフを配置）')
    set_ph(s, 2,  '分析結果\n\n月間アクティブユーザー数は\n6月に最大値（4,200名）を記録。\n\n8月は夏季休暇の影響で\n一時的に減少するが、\n9月以降は回復傾向。')

    # 19 — チャート：2面
    s = prs.slides.add_slide(layouts[19])
    set_ph(s, 15, 'チャネル別分析')
    set_ph(s, 16, 'Channel Analysis')
    set_ph(s, 1,  '（流入チャネル別\nコンバージョン率）')
    set_ph(s, 2,  '（チャネル別\n顧客単価推移）')

    # 20 — テーブル：フル幅
    s = prs.slides.add_slide(layouts[20])
    set_ph(s, 15, '機能比較一覧')
    set_ph(s, 16, 'Feature Comparison')
    set_ph(s, 1,  '（比較表を配置）')
    set_ph(s, 2,  '出典：各社公式ドキュメント 2026年3月時点')

    # 21 — テーブル＋ノート
    s = prs.slides.add_slide(layouts[21])
    set_ph(s, 15, 'コスト試算')
    set_ph(s, 16, 'Cost Estimation')
    set_ph(s, 1,  '（年間コスト試算表を配置）')
    set_ph(s, 2,  '注記\n\n初年度は移行コストを\n含むため割高。\n\n2年目以降は\n年間 ¥12M の\nコスト削減を見込む。')

    # 22 — 比較：サイドバイサイド
    s = prs.slides.add_slide(layouts[22])
    set_ph(s, 15, 'システム比較')
    set_ph(s, 16, 'System Comparison')
    set_ph(s, 1,  '現行CRM')
    set_ph(s, 2,  'レスポンス: 3.2秒\nモバイル: 非対応\nAI機能: なし\n月額コスト: ¥850/user\nAPI: 制限あり\nカスタマイズ: 低')
    set_ph(s, 3,  '新CRM（提案）')
    set_ph(s, 4,  'レスポンス: 0.8秒\nモバイル: 完全対応\nAI機能: 予測分析搭載\n月額コスト: ¥1,200/user\nAPI: フルオープン\nカスタマイズ: 高')

    # 23 — 比較：マトリクス (with real Table)
    s = prs.slides.add_slide(layouts[23])
    set_ph(s, 15, 'ベンダー評価マトリクス')
    set_ph(s, 16, 'Vendor Evaluation Matrix')
    add_matrix_table(s, [
        ['評価項目',     'A社', 'B社', 'C社'],
        ['機能性',       '◎',  '○',  '△'],
        ['拡張性',       '○',  '◎',  '○'],
        ['コスト',       '△',  '◎',  '○'],
        ['サポート',     '◎',  '○',  '◎'],
        ['セキュリティ', '○',  '△',  '◎'],
        ['総合評価',     'A',   'A+',  'B+'],
    ], emu(0.6), emu(1.5), emu(12.1), emu(5.2))

    # 24 — タイムライン：横
    s = prs.slides.add_slide(layouts[24])
    set_ph(s, 15, '導入ロードマップ')
    set_ph(s, 16, 'Implementation Roadmap')
    set_ph(s, 1,  'Phase 1\n要件定義\n2026 Q2')
    set_ph(s, 2,  'Phase 2\n開発・テスト\n2026 Q3-Q4')
    set_ph(s, 3,  'Phase 3\nパイロット\n2027 Q1')
    set_ph(s, 4,  'Phase 4\n全社展開\n2027 Q2')

    # 25 — プロセス：縦
    s = prs.slides.add_slide(layouts[25])
    set_ph(s, 15, 'データ移行プロセス')
    set_ph(s, 16, 'Data Migration Process')
    set_ph(s, 1,  'ステップ 1：データ棚卸し\n既存データの品質・整合性を評価し、移行対象を確定')
    set_ph(s, 2,  'ステップ 2：クレンジング＆変換\nデータ形式の統一、重複排除、欠損値の補完を実施')
    set_ph(s, 3,  'ステップ 3：検証＆本番投入\nテスト環境での検証後、段階的に本番環境へ移行')

    # 26 — アジェンダ
    s = prs.slides.add_slide(layouts[26])
    set_ph(s, 15, 'Agenda')
    set_ph(s, 16, '2026年3月24日\n会議室B-3F\n所要時間：90分')
    set_ph(s, 1,  '01  プロジェクト概要と目的\n\n02  現状分析データの共有\n\n03  システム比較と推奨案\n\n04  導入ロードマップ\n\n05  Q&A・ネクストステップ')

    # 27 — エグゼクティブサマリー
    s = prs.slides.add_slide(layouts[27])
    set_ph(s, 15, 'エグゼクティブサマリー')
    set_ph(s, 16, 'Executive Summary')
    set_ph(s, 1,  '主要な発見事項\n\n1. 現行CRMの応答速度は業界平均の3倍遅い\n2. ユーザー満足度は2年連続で低下傾向\n3. モバイル非対応による機会損失が年間¥50M\n4. AI活用による予測精度向上の余地が大きい')
    set_ph(s, 2,  '推奨アクション\n\n優先度：高\n\nB社CRMへの移行を\n2026年Q2に開始し、\n2027年Q2までに\n全社展開を完了する')

    # 28 — クロージング：感謝
    s = prs.slides.add_slide(layouts[28])
    set_ph(s, 10, 'THANK YOU')
    set_ph(s, 0,  'ご質問・ご意見をお待ちしています')
    set_ph(s, 11, 'プロジェクトマネージャー：山田 太郎')
    set_ph(s, 12, 'taro.yamada@sampletech.co.jp  |  内線: 1234\nSlack: #proj-nextgen-crm')

    # 29 — クロージング：ネクストステップ
    s = prs.slides.add_slide(layouts[29])
    set_ph(s, 15, 'Next Steps')
    set_ph(s, 1,  '1. 要件定義書の最終レビュー（4/1まで）\n\n2. ベンダー選定委員会の開催（4/8）\n\n3. 予算承認プロセスの開始（4/15）\n\n4. キックオフMTGの日程調整')
    set_ph(s, 2,  '期限：2026年4月末\n\n担当：山田太郎\n\n連絡先：\ntaro.yamada@\nsampletech.co.jp\n\n次回MTG：\n2026年4月8日 14:00')

    # ─── Save ───
    out = '/sessions/stoic-zealous-fermi/mnt/presentations/Midnight_Executive_30_Template.pptx'
    prs.save(out)
    print(f"\n✅ Saved: {out}")
    print(f"   Slides: {len(prs.slides)}")
    print(f"   Layouts: {len(list(prs.slide_layouts))}")
    print(f"\n   Layout list:")
    for i, layout in enumerate(prs.slide_layouts):
        print(f"     [{i:2d}] {layout.name}")
